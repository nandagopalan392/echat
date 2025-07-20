import os
import tempfile
import shutil
import subprocess
import logging
from pathlib import Path
from typing import List, Dict, Optional
import base64
from io import BytesIO
from PIL import Image

logger = logging.getLogger(__name__)

class PPTXImageConverter:
    """Convert PPTX slides to images using LibreOffice or alternative methods"""
    
    def __init__(self, temp_dir: Optional[str] = None):
        self.temp_dir = temp_dir or tempfile.gettempdir()
        self.libreoffice_available = self._check_libreoffice()
        
    def _check_libreoffice(self) -> bool:
        """Check if LibreOffice is available"""
        try:
            result = subprocess.run(['libreoffice', '--version'], 
                                 capture_output=True, text=True, timeout=10)
            return result.returncode == 0
        except (subprocess.TimeoutExpired, FileNotFoundError, subprocess.SubprocessError):
            return False
    
    def convert_pptx_to_images(self, pptx_path: str, output_format: str = 'png') -> List[Dict]:
        """
        Convert PPTX slides to images
        
        Args:
            pptx_path: Path to the PPTX file
            output_format: Image format (png, jpg)
            
        Returns:
            List of slide information with image data
        """
        if self.libreoffice_available:
            return self._convert_with_libreoffice(pptx_path, output_format)
        else:
            return self._convert_with_fallback(pptx_path, output_format)
    
    def _convert_with_libreoffice(self, pptx_path: str, output_format: str) -> List[Dict]:
        """Convert using LibreOffice"""
        slides = []
        
        try:
            # Create temporary directory for conversion
            with tempfile.TemporaryDirectory() as temp_output_dir:
                # Convert PPTX to images using LibreOffice
                cmd = [
                    'libreoffice',
                    '--headless',
                    '--convert-to', f'{output_format}',
                    '--outdir', temp_output_dir,
                    pptx_path
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
                
                if result.returncode != 0:
                    logger.error(f"LibreOffice conversion failed: {result.stderr}")
                    return self._convert_with_fallback(pptx_path, output_format)
                
                # Find generated image files
                output_files = list(Path(temp_output_dir).glob(f'*.{output_format}'))
                output_files.sort()  # Ensure proper order
                
                for i, image_path in enumerate(output_files, 1):
                    try:
                        # Read and encode image
                        with open(image_path, 'rb') as img_file:
                            img_data = img_file.read()
                            img_b64 = base64.b64encode(img_data).decode('utf-8')
                        
                        # Get image dimensions
                        with Image.open(image_path) as img:
                            width, height = img.size
                        
                        slides.append({
                            'slide_number': i,
                            'image_data': img_b64,
                            'image_format': output_format,
                            'width': width,
                            'height': height,
                            'file_size': len(img_data)
                        })
                        
                    except Exception as e:
                        logger.error(f"Error processing slide {i}: {e}")
                        continue
                
        except subprocess.TimeoutExpired:
            logger.error("LibreOffice conversion timed out")
            return self._convert_with_fallback(pptx_path, output_format)
        except Exception as e:
            logger.error(f"LibreOffice conversion error: {e}")
            return self._convert_with_fallback(pptx_path, output_format)
        
        return slides
    
    def _convert_with_fallback(self, pptx_path: str, output_format: str) -> List[Dict]:
        """Fallback method: Generate placeholder images with slide content"""
        slides = []
        
        try:
            from pptx import Presentation
            
            prs = Presentation(pptx_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # Extract text content
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                content = '\n\n'.join(slide_text) if slide_text else f"Slide {slide_num}"
                
                # Generate a simple text-based image
                placeholder_image = self._create_text_placeholder(content, slide_num)
                
                slides.append({
                    'slide_number': slide_num,
                    'image_data': placeholder_image,
                    'image_format': 'png',
                    'width': 800,
                    'height': 600,
                    'is_placeholder': True,
                    'content': content
                })
                
        except Exception as e:
            logger.error(f"Fallback conversion error: {e}")
        
        return slides
    
    def _create_text_placeholder(self, text: str, slide_number: int) -> str:
        """Create a simple text-based placeholder image"""
        try:
            # Create a simple image with text
            img_width, img_height = 800, 600
            img = Image.new('RGB', (img_width, img_height), color='white')
            
            # For now, create a simple colored rectangle with slide number
            # In a real implementation, you'd use PIL.ImageDraw to add text
            from PIL import ImageDraw, ImageFont
            
            draw = ImageDraw.Draw(img)
            
            # Draw background
            draw.rectangle([50, 50, img_width-50, img_height-50], fill='#f0f8ff', outline='#4169e1', width=2)
            
            # Draw slide number
            try:
                font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48)
            except:
                font = ImageFont.load_default()
            
            title_text = f"Slide {slide_number}"
            draw.text((100, 100), title_text, fill='#2c3e50', font=font)
            
            # Draw content (truncated)
            try:
                content_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
            except:
                content_font = ImageFont.load_default()
            
            # Wrap and draw text
            wrapped_text = self._wrap_text(text[:500], 80)  # Limit to 500 chars, 80 chars per line
            y_pos = 180
            for line in wrapped_text.split('\n')[:15]:  # Max 15 lines
                draw.text((100, y_pos), line, fill='#34495e', font=content_font)
                y_pos += 25
            
            # Convert to base64
            buffer = BytesIO()
            img.save(buffer, format='PNG')
            img_data = buffer.getvalue()
            return base64.b64encode(img_data).decode('utf-8')
            
        except Exception as e:
            logger.error(f"Error creating placeholder image: {e}")
            # Return a minimal fallback
            return ""
    
    def _wrap_text(self, text: str, width: int) -> str:
        """Simple text wrapping"""
        words = text.split()
        lines = []
        current_line = []
        
        for word in words:
            if len(' '.join(current_line + [word])) <= width:
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                    current_line = [word]
                else:
                    lines.append(word)
        
        if current_line:
            lines.append(' '.join(current_line))
        
        return '\n'.join(lines)

def convert_pptx_to_slide_images(pptx_file_path: str) -> Dict:
    """
    Main function to convert PPTX to slide images
    
    Returns:
        Dictionary with slide images and metadata
    """
    converter = PPTXImageConverter()
    slides = converter.convert_pptx_to_images(pptx_file_path)
    
    return {
        'type': 'presentation_images',
        'slides': slides,
        'total_slides': len(slides),
        'conversion_method': 'libreoffice' if converter.libreoffice_available else 'placeholder'
    }
