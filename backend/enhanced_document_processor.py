"""
Enhanced Document Processor with Multiple Chunking Methods
Handles different document types and chunking strategies with table-aware processing
"""

import os
import logging
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import json
import re
from dataclasses import dataclass

from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter, TokenTextSplitter
from langchain_community.document_loaders import PyPDFLoader

from chunking_config import ChunkingMethod, ChunkingConfig, get_chunking_config_manager
from table_extraction import (
    get_table_extractor, detect_table_in_text, 
    create_row_based_chunks, create_semantic_table_chunks, create_table_chunk_metadata
)

import os
import logging
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import json
import re
from dataclasses import dataclass

from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter, TokenTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from chunking_config import ChunkingMethod, ChunkingConfig, get_chunking_config_manager

logger = logging.getLogger(__name__)

# Try to import additional libraries for enhanced document processing
try:
    from langchain_community.document_loaders import Docx2txtLoader, CSVLoader, TextLoader
    from langchain_community.document_loaders import UnstructuredPowerPointLoader
    ADVANCED_LOADERS_AVAILABLE = True
except ImportError:
    logger.warning("Advanced document loaders not available. Some features may be limited.")
    ADVANCED_LOADERS_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    logger.warning("Pandas not available. Table processing may be limited.")
    PANDAS_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    logger.warning("OCR libraries not available. Image text extraction disabled.")
    OCR_AVAILABLE = False

@dataclass
class ChunkingResult:
    """Result of document chunking operation"""
    chunks: List[Document]
    metadata: Dict[str, Any]
    method_used: ChunkingMethod
    config_used: ChunkingConfig
    warnings: List[str] = None
    
    def __post_init__(self):
        if self.warnings is None:
            self.warnings = []

class EnhancedDocumentProcessor:
    """Enhanced document processor with multiple chunking methods"""
    
    def __init__(self):
        self.config_manager = get_chunking_config_manager()
        self.supported_extensions = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.doc': self._process_docx,
            '.txt': self._process_text,
            '.md': self._process_text,
            '.csv': self._process_csv,
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,
            '.ppt': self._process_presentation,
            '.pptx': self._process_presentation,
            '.html': self._process_html,
            '.json': self._process_json,
            '.eml': self._process_email,
            '.jpg': self._process_image,
            '.jpeg': self._process_image,
            '.png': self._process_image,
            '.gif': self._process_image,
            '.tif': self._process_image,
            '.tiff': self._process_image
        }
    
    def process_document(self, file_path: str, method: ChunkingMethod = None, 
                        config: ChunkingConfig = None, user_id: str = None, original_filename: str = None) -> ChunkingResult:
        """
        Process document with specified chunking method and configuration
        """
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext not in self.supported_extensions:
            raise ValueError(f"Unsupported file format: {file_ext}")
        
        # Use original filename if provided, otherwise use the file path name
        # Always extract just the filename (no path) for consistent metadata storage
        if original_filename:
            source_filename = Path(original_filename).name
        else:
            source_filename = Path(file_path).name
        if method is None:
            from chunking_config import FileFormatSupport
            method = FileFormatSupport.get_optimal_method(file_ext[1:])  # Remove dot
        
        # Get configuration if not specified
        if config is None:
            config = self.config_manager.get_config(method, user_id)
        
        # Validate configuration
        warnings = self.config_manager.validate_config(config)
        
        logger.info(f"Processing {file_path} with method {method.value}")
        
        # Load document content
        loader_func = self.supported_extensions[file_ext]
        raw_documents = loader_func(file_path)
        
        if not raw_documents:
            raise ValueError(f"No content extracted from {file_path}")
        
        # Apply chunking method
        chunks = self._apply_chunking_method(raw_documents, method, config)
        
        # Add metadata to chunks - clear existing metadata first to avoid complex objects
        for i, chunk in enumerate(chunks):
            chunk.metadata = {  # Replace instead of update to avoid complex objects
                'source_file': source_filename,
                'chunk_index': i,
                'chunking_method': method.value,
                'total_chunks': len(chunks)
            }
        
        # Debug: Log metadata after setting
        if chunks:
            logger.info(f"Enhanced processor - chunk metadata set: {chunks[0].metadata}")
        
        # Create processing metadata
        metadata = {
            'source_file': source_filename,
            'file_size': os.path.getsize(file_path),
            'chunk_count': len(chunks),
            'method_used': method.value,
            'config_used': config.to_dict(),
            'total_tokens': sum(len(chunk.page_content.split()) for chunk in chunks)
        }
        
        return ChunkingResult(
            chunks=chunks,
            metadata=metadata,
            method_used=method,
            config_used=config,
            warnings=warnings
        )
    
    def _apply_chunking_method(self, documents: List[Document], method: ChunkingMethod, 
                              config: ChunkingConfig) -> List[Document]:
        """Apply specific chunking method to documents with automatic table detection"""
        
        # First, check if any document contains tables and extract them
        all_chunks = []
        table_extractor = get_table_extractor()
        
        for doc in documents:
            doc_chunks = []
            source_file = doc.metadata.get('source', '')
            
            # Try to extract tables if we have a file path
            extracted_tables = []
            if source_file and Path(source_file).exists():
                try:
                    extracted_tables = table_extractor.extract_tables(str(source_file))
                    if extracted_tables:
                        logger.info(f"Extracted {len(extracted_tables)} tables from {source_file}")
                except Exception as e:
                    logger.warning(f"Table extraction failed for {source_file}: {e}")
            
            # Process extracted tables as separate chunks using proper row-based chunking
            for table_info in extracted_tables:
                table_data = table_info['data']
                
                if hasattr(table_data, 'shape') and not table_data.empty:
                    # Use proper row-based chunking for tables
                    if len(table_data) <= 20:  # Small tables: one chunk per row
                        table_chunks = create_row_based_chunks(table_data, table_info)
                    else:  # Large tables: semantic grouping
                        table_chunks = create_semantic_table_chunks(table_data, table_info, max_rows_per_chunk=10)
                    
                    # Convert to Document objects
                    for chunk_data in table_chunks:
                        chunk_metadata = chunk_data['metadata']
                        chunk_metadata['chunking_method'] = f"{method.value}_table_aware"
                        chunk_metadata.update(doc.metadata)  # Include original document metadata
                        
                        chunk = Document(
                            page_content=chunk_data['content'],
                            metadata=chunk_metadata
                        )
                        doc_chunks.append(chunk)
            
            # Now process the document content based on the selected method
            # but make it table-aware for text content too
            if method == ChunkingMethod.GENERAL:
                method_chunks = self._chunk_general_table_aware([doc], config)
            elif method == ChunkingMethod.QA:
                method_chunks = self._chunk_qa_table_aware([doc], config)
            elif method == ChunkingMethod.RESUME:
                method_chunks = self._chunk_resume_table_aware([doc], config)
            elif method == ChunkingMethod.TABLE:
                method_chunks = self._chunk_table([doc], config)
            elif method == ChunkingMethod.PRESENTATION:
                method_chunks = self._chunk_presentation([doc], config)
            elif method == ChunkingMethod.PICTURE:
                method_chunks = self._chunk_picture([doc], config)
            elif method == ChunkingMethod.EMAIL:
                method_chunks = self._chunk_email([doc], config)
            else:
                # Default to general chunking
                method_chunks = self._chunk_general_table_aware([doc], config)
            
            doc_chunks.extend(method_chunks)
            all_chunks.extend(doc_chunks)
        
        return all_chunks
    
    def _chunk_general(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Standard recursive character text splitting"""
        # Parse delimiter string into list of separators
        separators = config.delimiter.split('|') if '|' in config.delimiter else [config.delimiter]
        separators = [sep.replace('\\n', '\n').replace('\\t', '\t') for sep in separators]
        
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=config.chunk_token_num,
            chunk_overlap=config.chunk_overlap,
            separators=separators,
            length_function=len
        )
        
        return splitter.split_documents(documents)
    
    def _chunk_qa(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Q&A format chunking - split on question/answer patterns"""
        chunks = []
        qa_patterns = [r'Q:', r'A:', r'Question:', r'Answer:', r'\d+\.', r'Q\d+', r'A\d+']
        
        for doc in documents:
            content = doc.page_content
            
            # Find Q&A boundaries
            boundaries = []
            for pattern in qa_patterns:
                for match in re.finditer(pattern, content, re.IGNORECASE):
                    boundaries.append(match.start())
            
            boundaries = sorted(set(boundaries))
            boundaries.append(len(content))
            
            # Create chunks based on Q&A boundaries
            for i in range(len(boundaries) - 1):
                start, end = boundaries[i], boundaries[i + 1]
                chunk_text = content[start:end].strip()
                
                if len(chunk_text) > 50:  # Minimum chunk size
                    chunk = Document(
                        page_content=chunk_text,
                        metadata={'qa_chunk': True}  # Only simple metadata, no complex inheritance
                    )
                    chunks.append(chunk)
        
        return chunks
    
    def _chunk_resume(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Resume-specific chunking based on sections"""
        chunks = []
        section_patterns = [
            r'EXPERIENCE|WORK EXPERIENCE|EMPLOYMENT',
            r'EDUCATION|ACADEMIC',
            r'SKILLS|TECHNICAL SKILLS|COMPETENCIES',
            r'PROJECTS|PROJECT EXPERIENCE',
            r'CERTIFICATIONS|CERTIFICATES',
            r'SUMMARY|OBJECTIVE|PROFILE'
        ]
        
        for doc in documents:
            content = doc.page_content
            
            # Find section boundaries
            boundaries = [0]  # Start with beginning
            for pattern in section_patterns:
                for match in re.finditer(pattern, content, re.IGNORECASE):
                    boundaries.append(match.start())
            
            boundaries = sorted(set(boundaries))
            boundaries.append(len(content))
            
            # Create chunks for each section
            for i in range(len(boundaries) - 1):
                start, end = boundaries[i], boundaries[i + 1]
                chunk_text = content[start:end].strip()
                
                if len(chunk_text) > 100:  # Minimum meaningful section size
                    chunk = Document(
                        page_content=chunk_text,
                        metadata={'resume_section': True}  # Only simple metadata
                    )
                    chunks.append(chunk)
        
        return chunks
    

    
    def _chunk_table(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Table-aware chunking for structured data"""
        chunks = []
        table_extractor = get_table_extractor()
        
        for doc in documents:
            content = doc.page_content
            source_file = doc.metadata.get('source', '')
            
            # First, try to extract tables using the table extractor if we have a file path
            if source_file and Path(source_file).exists():
                try:
                    extracted_tables = table_extractor.extract_tables(str(source_file))
                    
                    if extracted_tables:
                        # Process extracted tables using proper row-based chunking
                        for table_info in extracted_tables:
                            table_data = table_info['data']
                            
                            if hasattr(table_data, 'shape') and not table_data.empty:  # pandas DataFrame
                                # Use proper row-based chunking for tables
                                if len(table_data) <= 20:  # Small tables: one chunk per row
                                    table_chunks = create_row_based_chunks(table_data, table_info)
                                else:  # Large tables: semantic grouping
                                    table_chunks = create_semantic_table_chunks(table_data, table_info, max_rows_per_chunk=10)
                                
                                # Convert to Document objects
                                for chunk_data in table_chunks:
                                    chunk_metadata = chunk_data['metadata']
                                    chunk_metadata.update(doc.metadata)  # Include original document metadata
                                    
                                    chunk = Document(
                                        page_content=chunk_data['content'],
                                        metadata=chunk_metadata
                                    )
                                    chunks.append(chunk)
                        
                        # If we successfully extracted tables, return them
                        if chunks:
                            return chunks
                            
                except Exception as e:
                    logger.warning(f"Table extraction failed for {source_file}: {e}")
            
            # Fallback: detect table structure in text content
            if detect_table_in_text(content):
                lines = content.split('\n')
                
                # Check if content looks like CSV/TSV
                if any('\t' in line or ',' in line for line in lines[:5]):
                    # Process as tabular data
                    delimiter = '\t' if '\t' in content else ','
                    
                    # Group rows into chunks
                    header = lines[0] if lines else ""
                    data_lines = lines[1:] if len(lines) > 1 else []
                    
                    chunk_size = max(10, config.chunk_token_num // 50)  # Rough estimate of rows per chunk
                    
                    for i in range(0, len(data_lines), chunk_size):
                        chunk_lines = [header] + data_lines[i:i + chunk_size]
                        chunk_content = '\n'.join(chunk_lines)
                        
                        chunk = Document(
                            page_content=chunk_content,
                            metadata={
                                'chunk_type': 'table',
                                'table_chunk': True, 
                                'chunk_rows': len(chunk_lines) - 1,
                                'extraction_method': 'text_detection'
                            }
                        )
                        chunks.append(chunk)
                else:
                    # Table-like structure but not CSV/TSV - preserve structure
                    table_chunks = self._chunk_table_like_text(content, config)
                    chunks.extend(table_chunks)
            else:
                # Not clearly tabular, use regular chunking
                chunks.extend(self._chunk_general([doc], config))
        
        return chunks

    def _chunk_general_table_aware(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """General chunking with table detection and special handling"""
        chunks = []
        
        for doc in documents:
            content = doc.page_content;
            
            # Check if content contains table-like structures
            if detect_table_in_text(content):
                logger.info(f"Detected table-like content in document, applying table-aware chunking")
                
                # Split content into table and non-table sections
                table_chunks = self._extract_table_sections_from_text(content, config)
                for chunk in table_chunks:
                    chunk.metadata.update({
                        'chunking_method': 'general_table_aware',
                        'contains_tables': True
                    })
                chunks.extend(table_chunks)
            else:
                # Regular general chunking for non-table content
                regular_chunks = self._chunk_general([doc], config)
                for chunk in regular_chunks:
                    chunk.metadata.update({
                        'chunking_method': 'general',
                        'contains_tables': False
                    })
                chunks.extend(regular_chunks);
        
        return chunks

    def _chunk_qa_table_aware(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Q&A chunking with table detection and special handling"""
        chunks = []
        
        for doc in documents:
            content = doc.page_content;
            
            # Check if content contains table-like structures
            if detect_table_in_text(content):
                logger.info(f"Detected table-like content in Q&A document, applying mixed chunking")
                
                # For Q&A documents with tables, we need to be more careful
                # Split by Q&A patterns first, then check each section for tables
                qa_sections = self._split_qa_sections(content);
                
                for i, section in enumerate(qa_sections):
                    if detect_table_in_text(section):
                        # This Q&A section contains a table
                        table_chunks = self._extract_table_sections_from_text(section, config);
                        for chunk in table_chunks:
                            chunk.metadata.update({
                                'chunking_method': 'qa_table_aware',
                                'qa_section': i,
                                'contains_tables': True
                            });
                        chunks.extend(table_chunks);
                    else:
                        # Regular Q&A processing
                        section_doc = Document(page_content=section, metadata=doc.metadata.copy());
                        qa_chunks = self._chunk_qa([section_doc], config);
                        for chunk in qa_chunks:
                            chunk.metadata.update({
                                'chunking_method': 'qa',
                                'qa_section': i,
                                'contains_tables': False
                            });
                        chunks.extend(qa_chunks);
            else:
                # Regular Q&A chunking
                qa_chunks = self._chunk_qa([doc], config);
                for chunk in qa_chunks:
                    chunk.metadata.update({
                        'chunking_method': 'qa',
                        'contains_tables': False
                    });
                chunks.extend(qa_chunks);
        
        return chunks;

    def _chunk_resume_table_aware(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Resume chunking with table detection and special handling"""
        chunks = []
        
        for doc in documents:
            content = doc.page_content;
            
            # Check if content contains table-like structures
            if detect_table_in_text(content):
                logger.info(f"Detected table-like content in resume, applying mixed chunking")
                
                # For resumes with tables (like skills matrices), preserve structure
                resume_sections = self._split_resume_sections(content);
                
                for section_name, section_content in resume_sections:
                    if detect_table_in_text(section_content):
                        # This resume section contains a table (e.g., skills matrix)
                        table_chunks = self._extract_table_sections_from_text(section_content, config);
                        for chunk in table_chunks:
                            chunk.metadata.update({
                                'chunking_method': 'resume_table_aware',
                                'resume_section': section_name,
                                'contains_tables': True
                            });
                        chunks.extend(table_chunks);
                    else:
                        # Regular resume section processing
                        section_doc = Document(page_content=section_content, metadata=doc.metadata.copy());
                        resume_chunks = self._chunk_resume([section_doc], config);
                        for chunk in resume_chunks:
                            chunk.metadata.update({
                                'chunking_method': 'resume',
                                'resume_section': section_name,
                                'contains_tables': False
                            });
                        chunks.extend(resume_chunks);
            else:
                # Regular resume chunking
                resume_chunks = self._chunk_resume([doc], config);
                for chunk in resume_chunks:
                    chunk.metadata.update({
                        'chunking_method': 'resume',
                        'contains_tables': False
                    });
                chunks.extend(resume_chunks);
        
        return chunks;

    def _extract_table_sections_from_text(self, content: str, config: ChunkingConfig) -> List[Document]:
        """Extract and chunk table sections from text content"""
        chunks = [];
        lines = content.split('\n');
        
        current_section = [];
        current_section_type = 'text';  # 'text' or 'table'
        
        i = 0;
        while i < len(lines):
            line = lines[i].strip();
            
            # Look ahead to detect table start
            table_start = self._detect_table_start(lines[i:i+5]);
            
            if table_start and current_section_type == 'text':
                # End current text section
                if current_section:
                    text_content = '\n'.join(current_section).strip();
                    if text_content:
                        chunks.append(Document(
                            page_content=text_content,
                            metadata={'chunk_type': 'text', 'extraction_method': 'text_section'}
                        ));
                
                # Start table section
                current_section = [line];
                current_section_type = 'table';
                
            elif not table_start and current_section_type == 'table':
                # End current table section
                if current_section:
                    table_content = '\n'.join(current_section).strip();
                    if table_content:
                        formatted_table = self._format_detected_table(table_content);
                        chunks.append(Document(
                            page_content=formatted_table,
                            metadata={
                                'chunk_type': 'table',
                                'extraction_method': 'text_detection',
                                'confidence': 0.8
                            }
                        ));
                
                # Start text section
                current_section = [line] if line else [];
                current_section_type = 'text';
                
            else:
                # Continue current section
                if line or current_section:  # Include empty lines within sections
                    current_section.append(line);
            
            i += 1;
        
        # Handle remaining section
        if current_section:
            section_content = '\n'.join(current_section).strip();
            if section_content:
                if current_section_type == 'table':
                    formatted_table = self._format_detected_table(section_content);
                    chunks.append(Document(
                        page_content=formatted_table,
                        metadata={
                            'chunk_type': 'table',
                            'extraction_method': 'text_detection',
                            'confidence': 0.8
                        }
                    ));
                else:
                    chunks.append(Document(
                        page_content=section_content,
                        metadata={'chunk_type': 'text', 'extraction_method': 'text_section'}
                    ));
        
        return chunks;

    def _detect_table_start(self, lines: List[str]) -> bool:
        """Detect if the next few lines start a table"""
        if len(lines) < 2:
            return False;
        
        # Look for table indicators in the next few lines
        table_lines = 0;
        for line in lines[:5]:
            if line.strip():
                if (line.count('|') >= 2 or 
                    line.count('\t') >= 1 or 
                    re.search(r'\s{3,}', line) or  # Multiple spaces
                    line.count(',') >= 2):
                    table_lines += 1;
        
        return table_lines >= 2;

    def _format_detected_table(self, table_content: str) -> str:
        """Format detected table content for better readability"""
        lines = table_content.split('\n');
        
        # Try to clean up and format the table
        formatted_lines = [];
        for line in lines:
            if line.strip():
                # Clean up excessive spaces and pipes
                cleaned_line = re.sub(r'\s*\|\s*', ' | ', line);
                cleaned_line = re.sub(r'\s+', ' ', cleaned_line).strip();
                formatted_lines.append(cleaned_line);
        
        return '\n'.join(formatted_lines);

    def _split_qa_sections(self, content: str) -> List[str]:
        """Split content into Q&A sections"""
        qa_patterns = [r'Q:', r'A:', r'Question:', r'Answer:', r'\d+\.', r'Q\d+', r'A\d+'];
        
        boundaries = [0];
        for pattern in qa_patterns:
            for match in re.finditer(pattern, content, re.IGNORECASE):
                boundaries.append(match.start());
        
        boundaries = sorted(set(boundaries));
        boundaries.append(len(content));
        
        sections = [];
        for i in range(len(boundaries) - 1):
            start, end = boundaries[i], boundaries[i + 1];
            section = content[start:end].strip();
            if section:
                sections.append(section);
        
        return sections;

    def _split_resume_sections(self, content: str) -> List[Tuple[str, str]]:
        """Split resume content into named sections"""
        section_patterns = {
            'experience': r'EXPERIENCE|WORK EXPERIENCE|EMPLOYMENT',
            'education': r'EDUCATION|ACADEMIC',
            'skills': r'SKILLS|TECHNICAL SKILLS|COMPETENCIES',
            'projects': r'PROJECTS|PROJECT EXPERIENCE',
            'certifications': r'CERTIFICATIONS|CERTIFICATES',
            'summary': r'SUMMARY|OBJECTIVE|PROFILE'
        };
        
        boundaries = [(0, 'header')];
        for section_name, pattern in section_patterns.items():
            for match in re.finditer(pattern, content, re.IGNORECASE):
                boundaries.append((match.start(), section_name));
        
        boundaries.sort();
        boundaries.append((len(content), 'end'));
        
        sections = [];
        for i in range(len(boundaries) - 1):
            start_pos, section_name = boundaries[i];
            end_pos, _ = boundaries[i + 1];
            
            section_content = content[start_pos:end_pos].strip();
            if section_content:
                sections.append((section_name, section_content));
        
        return sections;
    
    def _chunk_table_like_text(self, content: str, config: ChunkingConfig) -> List[Document]:
        """Chunk text that has table-like structure but isn't CSV/TSV"""
        chunks = []
        lines = content.split('\n')
        
        # Group lines that appear to be part of the same table
        current_table_lines = []
        
        for line in lines:
            if line.strip():
                # Check if this line looks like part of a table
                if ('|' in line or 
                    re.search(r'\s{3,}', line) or  # Multiple spaces (column separation)
                    re.search(r'^\s*[\w\s]+\s+[\w\s]+\s+[\w\s]+', line)):  # Multiple words separated by spaces
                    current_table_lines.append(line)
                else:
                    # End of current table - process it
                    if current_table_lines:
                        table_text = '\n'.join(current_table_lines)
                        chunk = Document(
                            page_content=table_text,
                            metadata={
                                'chunk_type': 'table',
                                'table_chunk': True,
                                'extraction_method': 'structure_detection'
                            }
                        )
                        chunks.append(chunk)
                        current_table_lines = []
                    
                    # Process non-table line normally
                    if line.strip():
                        chunk = Document(
                            page_content=line,
                            metadata={'chunk_type': 'text'}
                        )
                        chunks.append(chunk)
            else:
                # Empty line - might be table separator
                if current_table_lines:
                    current_table_lines.append(line)
        
        # Handle remaining table lines
        if current_table_lines:
            table_text = '\n'.join(current_table_lines)
            chunk = Document(
                page_content=table_text,
                metadata={
                    'chunk_type': 'table',
                    'table_chunk': True,
                    'extraction_method': 'structure_detection'
                }
            )
            chunks.append(chunk)
        
        return chunks
    
    def _chunk_presentation(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Presentation chunking - typically one chunk per slide, with table awareness"""
        chunks = []
        table_extractor = get_table_extractor()
        
        for doc in documents:
            content = doc.page_content
            source_file = doc.metadata.get('source', '')
            
            # First, try to extract tables from the presentation if we have a file path
            if source_file and Path(source_file).exists() and source_file.lower().endswith(('.pptx', '.ppt')):
                try:
                    extracted_tables = table_extractor.extract_tables(str(source_file))
                    
                    if extracted_tables:
                        # Group tables by slide
                        slide_tables = {}
                        for table_info in extracted_tables:
                            slide_num = table_info.get('slide_number', 1)
                            if slide_num not in slide_tables:
                                slide_tables[slide_num] = []
                            slide_tables[slide_num].append(table_info)
                        
                        # Process each slide's tables using proper row-based chunking
                        for slide_num, tables in slide_tables.items():
                            for table_info in tables:
                                table_data = table_info['data']
                                
                                if hasattr(table_data, 'shape') and not table_data.empty:
                                    # Use proper row-based chunking for tables
                                    if len(table_data) <= 20:  # Small tables: one chunk per row
                                        table_chunks = create_row_based_chunks(table_data, table_info)
                                    else:  # Large tables: semantic grouping
                                        table_chunks = create_semantic_table_chunks(table_data, table_info, max_rows_per_chunk=10)
                                    
                                    # Convert to Document objects
                                    for chunk_data in table_chunks:
                                        chunk_metadata = chunk_data['metadata']
                                        chunk_metadata.update({
                                            'slide_chunk': True,
                                            'content_type': 'table'
                                        })
                                        chunk_metadata.update(doc.metadata)  # Include original document metadata
                                        
                                        chunk = Document(
                                            page_content=chunk_data['content'],
                                            metadata=chunk_metadata
                                        )
                                        chunks.append(chunk)
                        
                        # If we successfully extracted tables, continue with regular slide processing for non-table content
                        # This allows both table and text content to be processed from presentations
                        
                except Exception as e:
                    logger.warning(f"Table extraction failed for presentation {source_file}: {e}")
            
            # Process slide boundaries for remaining content
            slide_patterns = [r'Slide\s+\d+', r'Page\s+\d+', r'^\d+\s*$']
            boundaries = [0]
            
            for pattern in slide_patterns:
                for match in re.finditer(pattern, content, re.MULTILINE):
                    boundaries.append(match.start())
            
            # If no clear slide boundaries, split by double newlines
            if len(boundaries) <= 1:
                slide_breaks = [m.start() for m in re.finditer(r'\n\s*\n\s*\n', content)]
                boundaries.extend(slide_breaks)
            
            boundaries = sorted(set(boundaries))
            boundaries.append(len(content))
            
            # Create chunks for each slide
            for i in range(len(boundaries) - 1):
                start, end = boundaries[i], boundaries[i + 1]
                slide_content = content[start:end].strip()
                
                if slide_content:
                    # Check if this slide content contains table-like structures
                    if detect_table_in_text(slide_content):
                        # Process as table-aware content
                        table_chunks = self._chunk_table_like_text(slide_content, config)
                        for chunk in table_chunks:
                            chunk.metadata.update({
                                'slide_chunk': True, 
                                'slide_number': i,
                                'content_type': 'mixed'
                            })
                        chunks.extend(table_chunks)
                    else:
                        # Regular slide content
                        chunk = Document(
                            page_content=slide_content,
                            metadata={
                                'slide_chunk': True, 
                                'slide_number': i,
                                'content_type': 'text'
                            }
                        )
                        chunks.append(chunk)
        
        return chunks
    
    def _chunk_picture(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Image content chunking - extract text via OCR if available"""
        chunks = []
        
        for doc in documents:
            # For image documents, the content might be OCR text
            content = doc.page_content
            
            if content and len(content.strip()) > 0:
                # Text was extracted from image
                if len(content) > config.chunk_token_num and config.chunk_token_num > 0:
                    splitter = RecursiveCharacterTextSplitter(
                        chunk_size=config.chunk_token_num,
                        chunk_overlap=config.chunk_overlap
                    )
                    image_chunks = splitter.split_documents([doc])
                    
                    # Distribute OCR bounding boxes among chunks if available
                    if 'ocr_bounding_boxes' in doc.metadata:
                        bounding_boxes = doc.metadata['ocr_bounding_boxes']
                        self._distribute_bounding_boxes_to_chunks(image_chunks, bounding_boxes)
                    
                    chunks.extend(image_chunks)
                else:
                    # Single chunk - include all bounding boxes
                    if 'ocr_bounding_boxes' in doc.metadata:
                        doc.metadata['chunk_bounding_boxes'] = doc.metadata['ocr_bounding_boxes']
                    chunks.append(doc)
            else:
                # No text content - create placeholder chunk
                chunk = Document(
                    page_content="[Image content - no text extracted]",
                    metadata={'image_only': True}  # Only simple metadata
                )
                chunks.append(chunk)
        
        return chunks

    def _chunk_email(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Email content chunking with table awareness"""
        chunks = []
        
        for doc in documents:
            content = doc.page_content
            
            # Check if email contains table-like structures
            if detect_table_in_text(content):
                # Process tables in email
                table_chunks = self._extract_table_sections_from_text(content, config)
                for chunk in table_chunks:
                    chunk.metadata.update({
                        'email_chunk': True,
                        'content_type': chunk.metadata.get('chunk_type', 'mixed')
                    })
                chunks.extend(table_chunks)
            else:
                # Regular email processing - split by email sections
                email_sections = self._split_email_sections(content)
                
                for section_name, section_content in email_sections:
                    if len(section_content.strip()) > 50:  # Minimum section size
                        chunk = Document(
                            page_content=section_content,
                            metadata={
                                'email_chunk': True,
                                'email_section': section_name,
                                'content_type': 'text'
                            }
                        )
                        chunks.append(chunk)
        
        return chunks

    def _split_email_sections(self, content: str) -> List[Tuple[str, str]]:
        """Split email content into sections (header, body, signature)"""
        sections = []
        
        # Simple email parsing
        lines = content.split('\n')
        
        header_end = 0
        signature_start = len(lines)
        
        # Find end of headers (first empty line or start of body content)
        for i, line in enumerate(lines):
            if not line.strip() or not (':' in line and i < 10):
                header_end = i
                break
        
        # Find start of signature (common signature indicators)
        signature_patterns = [r'--', r'Best regards', r'Sincerely', r'Thanks', r'Regards']
        for i in range(len(lines) - 1, max(len(lines) - 10, header_end), -1):
            line = lines[i].strip()
            if any(re.search(pattern, line, re.IGNORECASE) for pattern in signature_patterns):
                signature_start = i
                break
        
        # Extract sections
        if header_end > 0:
            header_content = '\n'.join(lines[:header_end]).strip()
            if header_content:
                sections.append(('header', header_content))
        
        body_content = '\n'.join(lines[header_end:signature_start]).strip()
        if body_content:
            sections.append(('body', body_content))
        
        if signature_start < len(lines):
            signature_content = '\n'.join(lines[signature_start:]).strip()
            if signature_content:
                sections.append(('signature', signature_content))
        
        return sections

    def _distribute_bounding_boxes_to_chunks(self, chunks: List[Document], bounding_boxes: List[dict]):
        """Distribute OCR bounding boxes among text chunks based on text content matching"""
        if not bounding_boxes:
            return
            
        # Create a mapping of words to their bounding boxes
        word_to_bbox = {}
        for bbox in bounding_boxes:
            word_key = bbox['text'].lower().strip()
            if word_key:
                if word_key not in word_to_bbox:
                    word_to_bbox[word_key] = []
                word_to_bbox[word_key].append(bbox)
        
        # For each chunk, find matching bounding boxes
        for chunk in chunks:
            chunk_bboxes = []
            chunk_text = chunk.page_content.lower()
            chunk_words = chunk_text.split()
            
            # Try to match words from chunk to bounding boxes
            used_bboxes = set()
            for word in chunk_words:
                clean_word = ''.join(c for c in word if c.isalnum()).lower()
                if clean_word in word_to_bbox:
                    # Find an unused bounding box for this word
                    for bbox in word_to_bbox[clean_word]:
                        bbox_id = f"{bbox['left']},{bbox['top']},{bbox['width']},{bbox['height']}"
                        if bbox_id not in used_bboxes:
                            chunk_bboxes.append(bbox)
                            used_bboxes.add(bbox_id)
                            break
            
            # Add bounding boxes to chunk metadata
            chunk.metadata['chunk_bounding_boxes'] = chunk_bboxes
    
    def _chunk_email(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Email chunking based on email structure"""
        chunks = []
        email_patterns = [
            r'^From:', r'^To:', r'^Subject:', r'^Date:',
            r'^Reply-To:', r'^CC:', r'^BCC:'
        ]
        
        for doc in documents:
            content = doc.page_content
            
            # Find email header boundaries
            boundaries = [0]
            for pattern in email_patterns:
                for match in re.finditer(pattern, content, re.MULTILINE | re.IGNORECASE):
                    boundaries.append(match.start())
            
            # Also split on email separators
            email_separators = [
                r'^>{1,}\s*',  # Quote markers
                r'^-{3,}',     # Dash separators
                r'^={3,}',     # Equal separators
                r'On .* wrote:'  # Email thread markers
            ]
            
            for pattern in email_separators:
                for match in re.finditer(pattern, content, re.MULTILINE):
                    boundaries.append(match.start())
            
            boundaries = sorted(set(boundaries))
            boundaries.append(len(content))
            
            # Create chunks for email sections
            for i in range(len(boundaries) - 1):
                start, end = boundaries[i], boundaries[i + 1]
                email_section = content[start:end].strip()
                
                if len(email_section) > 50:  # Minimum meaningful email section
                    chunk = Document(
                        page_content=email_section,
                        metadata={'email_section': True}  # Only simple metadata
                    )
                    chunks.append(chunk)
        
        return chunks
    

    
    # Document loading methods
    def _process_pdf(self, file_path: str) -> List[Document]:
        """Process PDF files with table extraction"""
        documents = []
        table_extractor = get_table_extractor()
        
        try:
            # First load the PDF with standard text extraction
            loader = PyPDFLoader(file_path)
            text_documents = loader.load()
            
            # Extract tables from the PDF
            try:
                extracted_tables = table_extractor.extract_tables(file_path)
                
                if extracted_tables:
                    # Create separate documents for tables using proper row-based chunking
                    for table_info in extracted_tables:
                        table_data = table_info['data']
                        
                        if hasattr(table_data, 'shape') and not table_data.empty:
                            # Use proper row-based chunking for tables
                            if len(table_data) <= 20:  # Small tables: one chunk per row
                                table_chunks = create_row_based_chunks(table_data, table_info)
                            else:  # Large tables: semantic grouping
                                table_chunks = create_semantic_table_chunks(table_data, table_info, max_rows_per_chunk=10)
                            
                            # Convert to Document objects
                            for chunk_data in table_chunks:
                                chunk_metadata = chunk_data['metadata']
                                chunk_metadata.update({
                                    'source': Path(file_path).name,
                                    'content_type': 'table',
                                    'extracted_from': 'pdf'
                                })
                                
                                table_doc = Document(
                                    page_content=chunk_data['content'],
                                    metadata=chunk_metadata
                                )
                                documents.append(table_doc)
                
            except Exception as e:
                logger.warning(f"Table extraction failed for PDF {file_path}: {e}")
            
            # Process text documents and check for table-like content
            for doc in text_documents:
                content = doc.page_content
                
                # Check if this page contains table-like structures
                if detect_table_in_text(content):
                    # Mark as containing tables
                    doc.metadata.update({
                        'content_type': 'mixed',
                        'contains_tables': True
                    })
                else:
                    doc.metadata.update({
                        'content_type': 'text',
                        'contains_tables': False
                    })
                
                documents.append(doc)
            
            return documents
            
        except Exception as e:
            logger.error(f"Failed to process PDF {file_path}: {e}")
            return []
    
    def _process_docx(self, file_path: str) -> List[Document]:
        """Process DOCX files"""
        if ADVANCED_LOADERS_AVAILABLE:
            try:
                loader = Docx2txtLoader(file_path)
                return loader.load()
            except Exception as e:
                logger.error(f"Failed to process DOCX {file_path}: {e}")
        
        # Fallback to simple text reading
        return self._process_text(file_path)
    
    def _process_text(self, file_path: str) -> List[Document]:
        """Process text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            return [Document(
                page_content=content,
                metadata={'source': Path(file_path).name}  # Use filename only
            )]
        except Exception as e:
            logger.error(f"Failed to process text file {file_path}: {e}")
            return []
    
    def _process_csv(self, file_path: str) -> List[Document]:
        """Process CSV files with proper table extraction"""
        documents = []
        table_extractor = get_table_extractor()
        
        try:
            # Use the table extractor for proper CSV processing
            extracted_tables = table_extractor.extract_tables(file_path)
            
            if extracted_tables:
                for table_info in extracted_tables:
                    table_data = table_info['data']
                    
                    if hasattr(table_data, 'shape') and not table_data.empty:
                        # Use proper row-based chunking for CSV tables
                        if len(table_data) <= 20:  # Small tables: one chunk per row
                            table_chunks = create_row_based_chunks(table_data, table_info)
                        else:  # Large tables: semantic grouping
                            table_chunks = create_semantic_table_chunks(table_data, table_info, max_rows_per_chunk=10)
                        
                        # Convert to Document objects
                        for chunk_data in table_chunks:
                            chunk_metadata = chunk_data['metadata']
                            chunk_metadata.update({
                                'source': Path(file_path).name,
                                'type': 'csv',
                                'content_type': 'table',
                                'extracted_from': 'csv'
                            })
                            
                            documents.append(Document(
                                page_content=chunk_data['content'],
                                metadata=table_metadata
                            ))
            else:
                # Fallback to LangChain CSV loader
                if ADVANCED_LOADERS_AVAILABLE:
                    try:
                        loader = CSVLoader(file_path)
                        documents = loader.load()
                        
                        # Mark as table content
                        for doc in documents:
                            doc.metadata.update({
                                'content_type': 'table',
                                'extraction_method': 'langchain_csv'
                            })
                    except Exception as e:
                        logger.error(f"Failed to process CSV {file_path}: {e}")
                        # Fallback to simple text reading
                        documents = self._process_text(file_path)
                        if documents:
                            documents[0].metadata.update({
                                'content_type': 'table',
                                'extraction_method': 'text_fallback'
                            })
        
        except Exception as e:
            logger.error(f"Failed to process CSV file {file_path}: {e}")
            # Fallback to simple text reading
            documents = self._process_text(file_path)
            if documents:
                documents[0].metadata.update({
                    'content_type': 'table',
                    'extraction_method': 'text_fallback'
                })
        
        return documents
    
    def _process_excel(self, file_path: str) -> List[Document]:
        """Process Excel files with proper table extraction"""
        documents = []
        table_extractor = get_table_extractor()
        
        try:
            # Use the table extractor for proper Excel processing
            extracted_tables = table_extractor.extract_tables(file_path)
            
            if extracted_tables:
                for table_info in extracted_tables:
                    table_data = table_info['data']
                    
                    if hasattr(table_data, 'shape') and not table_data.empty:
                        # Use proper row-based chunking for Excel tables
                        if len(table_data) <= 20:  # Small tables: one chunk per row
                            table_chunks = create_row_based_chunks(table_data, table_info)
                        else:  # Large tables: semantic grouping
                            table_chunks = create_semantic_table_chunks(table_data, table_info, max_rows_per_chunk=10)
                        
                        # Convert to Document objects
                        for chunk_data in table_chunks:
                            chunk_metadata = chunk_data['metadata']
                            chunk_metadata.update({
                                'source': Path(file_path).name,
                                'type': 'excel',
                                'content_type': 'table',
                                'extracted_from': 'excel'
                            })
                            
                            documents.append(Document(
                                page_content=chunk_data['content'],
                                metadata=chunk_metadata
                            ))
            else:
                # Fallback to pandas if table extractor doesn't work
                if PANDAS_AVAILABLE:
                    try:
                        df = pd.read_excel(file_path)
                        content = df.to_string()
                        
                        documents.append(Document(
                            page_content=content,
                            metadata={
                                'source': Path(file_path).name, 
                                'type': 'excel',
                                'content_type': 'table',
                                'extraction_method': 'pandas_fallback'
                            }
                        ))
                    except Exception as e:
                        logger.error(f"Failed to process Excel file with pandas {file_path}: {e}")
        
        except Exception as e:
            logger.error(f"Failed to process Excel file {file_path}: {e}")
        
        return documents
    
    def _process_presentation(self, file_path: str) -> List[Document]:
        """Process PowerPoint files"""
        try:
            # Try using python-pptx first
            try:
                from pptx import Presentation
                return self._process_pptx_with_python_pptx(file_path)
            except ImportError:
                logger.warning("python-pptx not available, trying unstructured")
            
            # Fallback to unstructured if available
            if ADVANCED_LOADERS_AVAILABLE:
                try:
                    loader = UnstructuredPowerPointLoader(file_path)
                    return loader.load()
                except Exception as e:
                    logger.error(f"Failed to process presentation with unstructured {file_path}: {e}")
            
        except Exception as e:
            logger.error(f"Failed to process presentation {file_path}: {e}")
        
        return []
    
    def _process_pptx_with_python_pptx(self, file_path: str) -> List[Document]:
        """Process PPTX files using python-pptx library with table extraction"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            documents = []
            table_extractor = get_table_extractor()
            
            # Extract tables using the table extractor
            extracted_tables = {}  # slide_number -> [table_info]
            try:
                tables = table_extractor.extract_tables(file_path)
                for table_info in tables:
                    slide_num = table_info.get('slide_number', 1)
                    if slide_num not in extracted_tables:
                        extracted_tables[slide_num] = []
                    extracted_tables[slide_num].append(table_info)
            except Exception as e:
                logger.warning(f"Table extraction failed for presentation {file_path}: {e}")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_tables = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Handle tables - check if shape actually contains a table
                    if hasattr(shape, 'table') and shape.has_table:
                        try:
                            table_text = []
                            for row in shape.table.rows:
                                row_text = []
                                for cell in row.cells:
                                    if cell.text.strip():
                                        row_text.append(cell.text.strip())
                                if row_text:
                                    table_text.append(' | '.join(row_text))
                            if table_text:
                                table_content = '\n'.join(table_text)
                                slide_tables.append(table_content)
                        except Exception as e:
                            logger.warning(f"Could not process table in slide {slide_num}: {e}")
                
                # Create document for slide text content
                if slide_text:
                    slide_content = '\n\n'.join(slide_text)
                    
                    # Check if slide contains table-like structures
                    contains_tables = bool(slide_tables) or detect_table_in_text(slide_content)
                    
                    documents.append(Document(
                        page_content=slide_content,
                        metadata={
                            'source': Path(file_path).name,
                            'slide_number': slide_num,
                            'type': 'presentation',
                            'content_type': 'mixed' if contains_tables else 'text',
                            'contains_tables': contains_tables,
                            'total_slides': len(prs.slides)
                        }
                    ))
                
                # Create separate documents for extracted tables on this slide using proper row-based chunking
                if slide_num in extracted_tables:
                    for table_info in extracted_tables[slide_num]:
                        table_data = table_info['data']
                        
                        if hasattr(table_data, 'shape') and not table_data.empty:
                            # Use proper row-based chunking for presentation tables
                            if len(table_data) <= 20:  # Small tables: one chunk per row
                                table_chunks = create_row_based_chunks(table_data, table_info)
                            else:  # Large tables: semantic grouping
                                table_chunks = create_semantic_table_chunks(table_data, table_info, max_rows_per_chunk=10)
                            
                            # Convert to Document objects
                            for chunk_data in table_chunks:
                                chunk_metadata = chunk_data['metadata']
                                chunk_metadata.update({
                                    'source': Path(file_path).name,
                                    'type': 'presentation',
                                    'content_type': 'table',
                                    'extracted_from': 'presentation',
                                    'total_slides': len(prs.slides)
                                })
                                
                                documents.append(Document(
                                    page_content=chunk_data['content'],
                                    metadata=chunk_metadata
                                ))
                
                # Also add simple table documents for inline tables
                for table_content in slide_tables:
                    documents.append(Document(
                        page_content=table_content,
                        metadata={
                            'source': Path(file_path).name,
                            'slide_number': slide_num,
                            'type': 'presentation',
                            'content_type': 'table',
                            'extraction_method': 'inline_table',
                            'total_slides': len(prs.slides)
                        }
                    ))
            
            return documents
            
        except Exception as e:
            logger.error(f"Failed to process PPTX with python-pptx {file_path}: {e}")
            return []
    
    def _process_html(self, file_path: str) -> List[Document]:
        """Process HTML files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Simple HTML tag removal
            import re
            clean_content = re.sub(r'<[^>]+>', '', content)
            clean_content = re.sub(r'\s+', ' ', clean_content).strip()
            
            return [Document(
                page_content=clean_content,
                metadata={'source': Path(file_path).name, 'type': 'html'}  # Use filename only
            )]
        except Exception as e:
            logger.error(f"Failed to process HTML file {file_path}: {e}")
            return []
    
    def _process_json(self, file_path: str) -> List[Document]:
        """Process JSON files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Convert JSON to readable text
            content = json.dumps(data, indent=2)
            
            return [Document(
                page_content=content,
                metadata={'source': Path(file_path).name, 'type': 'json'}  # Use filename only
            )]
        except Exception as e:
            logger.error(f"Failed to process JSON file {file_path}: {e}")
            return []
    
    def _process_email(self, file_path: str) -> List[Document]:
        """Process email files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            return [Document(
                page_content=content,
                metadata={'source': Path(file_path).name, 'type': 'email'}  # Use filename only
            )]
        except Exception as e:
            logger.error(f"Failed to process email file {file_path}: {e}")
            return []
    
    def _process_image(self, file_path: str) -> List[Document]:
        """Process image files with OCR if available"""
        if OCR_AVAILABLE:
            try:
                logger.info(f"Processing image with OCR: {file_path}")
                image = Image.open(file_path)
                
                # Enhance image for better OCR results
                # Convert to RGB if needed
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                
                # Get OCR text with better configuration
                custom_config = r'--oem 3 --psm 6'
                text = pytesseract.image_to_string(image, config=custom_config)
                
                # Get OCR data with bounding boxes
                ocr_data = pytesseract.image_to_data(image, config=custom_config, output_type=pytesseract.Output.DICT)
                
                # Process bounding box data with lower confidence threshold
                bounding_boxes = []
                image_width, image_height = image.size
                
                for i in range(len(ocr_data['text'])):
                    confidence = int(ocr_data['conf'][i])
                    if confidence > 10:  # Lower confidence threshold to capture more text
                        word_text = ocr_data['text'][i].strip()
                        if word_text and len(word_text) > 0:  # Only include non-empty text
                            bbox = {
                                'text': word_text,
                                'left': int(ocr_data['left'][i]),
                                'top': int(ocr_data['top'][i]),
                                'width': int(ocr_data['width'][i]),
                                'height': int(ocr_data['height'][i]),
                                'confidence': confidence,
                                'block_num': int(ocr_data['block_num'][i]),
                                'par_num': int(ocr_data['par_num'][i]),
                                'line_num': int(ocr_data['line_num'][i]),
                                'word_num': int(ocr_data['word_num'][i])
                            }
                            bounding_boxes.append(bbox)
                
                # Log OCR output
                logger.info(f"OCR extracted text from {Path(file_path).name}:")
                logger.info(f"--- OCR OUTPUT START ---")
                logger.info(text.strip())
                logger.info(f"--- OCR OUTPUT END ---")
                logger.info(f"OCR extracted {len(text.strip())} characters and {len(bounding_boxes)} bounding boxes from {Path(file_path).name}")
                
                # If no text was extracted, try different OCR settings
                if len(text.strip()) < 10 and len(bounding_boxes) < 5:
                    logger.warning(f"Poor OCR results, trying alternative settings for {Path(file_path).name}")
                    
                    # Try different PSM mode for single block of text
                    alt_config = r'--oem 3 --psm 8'
                    alt_text = pytesseract.image_to_string(image, config=alt_config)
                    alt_ocr_data = pytesseract.image_to_data(image, config=alt_config, output_type=pytesseract.Output.DICT)
                    
                    if len(alt_text.strip()) > len(text.strip()):
                        logger.info(f"Alternative OCR settings produced better results for {Path(file_path).name}")
                        text = alt_text
                        ocr_data = alt_ocr_data
                        
                        # Reprocess bounding boxes with alternative data
                        bounding_boxes = []
                        for i in range(len(ocr_data['text'])):
                            confidence = int(ocr_data['conf'][i])
                            if confidence > 10:
                                word_text = ocr_data['text'][i].strip()
                                if word_text and len(word_text) > 0:
                                    bbox = {
                                        'text': word_text,
                                        'left': int(ocr_data['left'][i]),
                                        'top': int(ocr_data['top'][i]),
                                        'width': int(ocr_data['width'][i]),
                                        'height': int(ocr_data['height'][i]),
                                        'confidence': confidence,
                                        'block_num': int(ocr_data['block_num'][i]),
                                        'par_num': int(ocr_data['par_num'][i]),
                                        'line_num': int(ocr_data['line_num'][i]),
                                        'word_num': int(ocr_data['word_num'][i])
                                    }
                                    bounding_boxes.append(bbox)
                
                return [Document(
                    page_content=text,
                    metadata={
                        'source': Path(file_path).name, 
                        'type': 'image', 
                        'extracted_via': 'ocr',
                        'image_width': image_width,
                        'image_height': image_height,
                        'ocr_bounding_boxes': bounding_boxes
                    }
                )]
            except Exception as e:
                logger.error(f"Failed to process image {file_path}: {e}")
        else:
            logger.warning(f"OCR not available for image processing: {file_path}")
        
        # Return placeholder if OCR not available
        return [Document(
            page_content="[Image file - OCR not available]",
            metadata={'source': Path(file_path).name, 'type': 'image', 'extracted_via': 'none'}  # Use filename only
        )]

# Global instance
_document_processor = None

def get_document_processor() -> EnhancedDocumentProcessor:
    """Get global document processor instance"""
    global _document_processor
    if _document_processor is None:
        _document_processor = EnhancedDocumentProcessor()
    return _document_processor
