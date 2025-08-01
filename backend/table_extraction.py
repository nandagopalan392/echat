"""
Table Extraction and Processing for Document Chunking
Handles table detection and extraction from various document formats
"""
import pandas as pd
import re
import logging
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
from collections import Counter

logger = logging.getLogger(__name__)

# Try to import table extraction libraries
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    logger.warning("pdfplumber not available. PDF table extraction will be limited.")
    PDFPLUMBER_AVAILABLE = False

try:
    import tabula
    TABULA_AVAILABLE = True
except ImportError:
    logger.warning("tabula-py not available. Advanced PDF table extraction disabled.")
    TABULA_AVAILABLE = False

try:
    import camelot
    CAMELOT_AVAILABLE = True
except ImportError:
    logger.warning("camelot-py not available. Advanced PDF table extraction disabled.")
    CAMELOT_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    logger.warning("python-pptx not available. PowerPoint table extraction disabled.")
    PPTX_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    logger.warning("python-docx not available. Word table extraction disabled.")
    DOCX_AVAILABLE = False


class TableExtractor:
    """Extract tables from various document formats"""
    
    def __init__(self):
        self.supported_formats = {
            '.pdf': self.extract_pdf_tables,
            '.pptx': self.extract_pptx_tables,
            '.ppt': self.extract_pptx_tables,
            '.docx': self.extract_docx_tables,
            '.doc': self.extract_docx_tables,
            '.csv': self.extract_csv_tables,
            '.xlsx': self.extract_excel_tables,
            '.xls': self.extract_excel_tables
        }
    
    def extract_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract tables from a document file
        
        Returns:
            List of table dictionaries with structure:
            {
                'data': pandas.DataFrame,
                'page_number': int or None,
                'slide_number': int or None,
                'table_number': int,
                'bbox': tuple or None (x1, y1, x2, y2),
                'extraction_method': str,
                'confidence': float
            }
        """
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext not in self.supported_formats:
            logger.warning(f"File format {file_ext} not supported for table extraction")
            return []
        
        try:
            return self.supported_formats[file_ext](file_path)
        except Exception as e:
            logger.error(f"Error extracting tables from {file_path}: {e}")
            return []
    
    def extract_pdf_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract tables from PDF files using multiple methods"""
        tables = []
        
        # Method 1: pdfplumber (best for simple tables)
        if PDFPLUMBER_AVAILABLE:
            tables.extend(self._extract_pdf_tables_pdfplumber(file_path))
        
        # Method 2: tabula-py (good for complex tables) - only if pdfplumber failed
        if TABULA_AVAILABLE and len(tables) == 0:
            tables.extend(self._extract_pdf_tables_tabula(file_path))
        
        # Method 3: camelot (most accurate but slower) - only if others failed
        if CAMELOT_AVAILABLE and len(tables) == 0:
            tables.extend(self._extract_pdf_tables_camelot(file_path))
        
        return tables
    
    def _extract_pdf_tables_pdfplumber(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract tables using pdfplumber with improved structure preservation"""
        print(f"\n🔍 DEBUG: Starting pdfplumber extraction for {file_path}")
        tables = []
        
        try:
            with pdfplumber.open(file_path) as pdf:
                print(f"🔍 DEBUG: PDF has {len(pdf.pages)} pages")
                
                for page_num, page in enumerate(pdf.pages, 1):
                    print(f"\n🔍 DEBUG: Processing page {page_num}")
                    
                    # Try different table extraction settings
                    extraction_settings = [
                        {"vertical_strategy": "lines", "horizontal_strategy": "lines"},
                        {"vertical_strategy": "text", "horizontal_strategy": "text"},
                        {"vertical_strategy": "lines_strict", "horizontal_strategy": "lines_strict"},
                    ]
                    
                    page_tables_found = False
                    
                    for settings in extraction_settings:
                        print(f"🔍 DEBUG: Trying extraction with settings: {settings}")
                        try:
                            page_tables = page.extract_tables(table_settings=settings)
                            print(f"🔍 DEBUG: Found {len(page_tables)} raw tables with current settings")
                            
                            for table_num, table_data in enumerate(page_tables, 1):
                                print(f"🔍 DEBUG: Processing table {table_num} - rows: {len(table_data) if table_data else 0}")
                                
                                if table_data and len(table_data) > 1:  # Must have header + data
                                    print(f"🔍 DEBUG: Raw table data preview (first 3 rows):")
                                    for i, row in enumerate(table_data[:3]):
                                        print(f"  Row {i}: {row}")
                                    
                                    # Quick validation: check if this looks like a real table
                                    if not self._quick_table_validation(table_data):
                                        print(f"🔍 DEBUG: Quick validation failed for table {table_num} - skipping")
                                        continue
                                    
                                    # Use direct table processing (no complex structure recognition)
                                    structured_table = self._process_table_data(
                                        table_data, page_num, table_num, settings["vertical_strategy"]
                                    )
                                    
                                    if structured_table:
                                        print(f"🔍 DEBUG: Successfully processed table {table_num}")
                                        print(f"  DataFrame shape: {structured_table['data'].shape}")
                                        print(f"  DataFrame columns: {list(structured_table['data'].columns)}")
                                        print(f"  DataFrame head:")
                                        print(structured_table['data'].head(3))
                                        
                                        tables.append(structured_table)
                                        page_tables_found = True
                                        logger.info(f"Successfully extracted table {table_num} from page {page_num} using {settings}")
                                    else:
                                        print(f"🔍 DEBUG: Failed to process table {table_num}")
                            
                            if page_tables_found:
                                print(f"🔍 DEBUG: Found tables on page {page_num}, stopping extraction attempts")
                                break  # Use the first successful extraction method
                                
                        except Exception as e:
                            print(f"🔍 DEBUG: Table extraction failed with settings {settings}: {e}")
                            logger.warning(f"Table extraction failed with settings {settings} on page {page_num}: {e}")
                            continue
                    
                    if not page_tables_found:
                        print(f"🔍 DEBUG: No tables found on page {page_num} with pdfplumber")
                        logger.info(f"No tables found on page {page_num} with pdfplumber")
                        
        except Exception as e:
            print(f"🔍 DEBUG: pdfplumber extraction failed: {e}")
            logger.error(f"pdfplumber table extraction failed: {e}")
        
        print(f"🔍 DEBUG: Total tables extracted: {len(tables)}")
        return tables

    def _quick_table_validation(self, table_data: List[List]) -> bool:
        """Quick validation to filter out obvious non-tables before processing"""
        print(f"\n🔍 DEBUG: Quick table validation for {len(table_data)} rows")
        
        if not table_data or len(table_data) < 2:
            print("🔍 DEBUG: Too few rows for a table")
            return False
        
        # Check if we have a reasonable number of columns consistently
        row_lengths = [len(row) for row in table_data if row]
        if not row_lengths:
            print("🔍 DEBUG: No valid rows found")
            return False
        
        max_cols = max(row_lengths)
        min_cols = min(row_lengths)
        
        print(f"  Column count range: {min_cols} to {max_cols}")
        
        # If columns vary too much, it's likely not a structured table
        if max_cols > 2 and (max_cols - min_cols) > max_cols * 0.5:
            print("🔍 DEBUG: Too much variation in column count - likely not a table")
            return False
        
        # Check for extremely long cells (indicating paragraph text)
        very_long_cells = 0
        total_non_empty_cells = 0
        
        for row in table_data[:5]:  # Check first 5 rows
            for cell in row:
                if cell and str(cell).strip():
                    cell_str = str(cell).strip()
                    total_non_empty_cells += 1
                    if len(cell_str) > 200:  # Very long cell
                        very_long_cells += 1
                        print(f"    Found very long cell ({len(cell_str)} chars): {cell_str[:100]}...")
        
        if total_non_empty_cells > 0 and (very_long_cells / total_non_empty_cells) > 0.3:
            print(f"🔍 DEBUG: Too many very long cells ({very_long_cells}/{total_non_empty_cells}) - likely text, not table")
            return False
        
        # Check if we have at least some structure (not all cells are huge text blocks)
        reasonable_cells = 0
        for row in table_data[:5]:
            for cell in row:
                if cell and str(cell).strip():
                    cell_str = str(cell).strip()
                    if 2 <= len(cell_str) <= 100:  # Reasonable cell size
                        reasonable_cells += 1
        
        print(f"  Reasonable cells in first 5 rows: {reasonable_cells}")
        
        if reasonable_cells < 3:  # Need at least a few reasonable cells
            print("🔍 DEBUG: Not enough reasonable-sized cells - likely not a table")
            return False
        
        print("🔍 DEBUG: Quick validation passed")
        return True

    def _process_table_data(self, table_data: List[List], page_num: int, 
                           table_num: int, method: str) -> Optional[Dict[str, Any]]:
        """Process table data directly without complex structure recognition"""
        print(f"\n🔍 DEBUG: _process_table_data called for table {table_num} on page {page_num}")
        print(f"  Method: {method}")
        print(f"  Raw table_data length: {len(table_data) if table_data else 0}")
        
        try:
            # Clean the table data
            cleaned_table = self._clean_table_data(table_data)
            print(f"🔍 DEBUG: Cleaned table data: {len(cleaned_table) if cleaned_table else 0} rows")
            
            if not cleaned_table or len(cleaned_table) < 2:
                print("🔍 DEBUG: Insufficient table data after cleaning")
                return None
            
            # Detect and merge multi-row headers
            headers, data_start_idx = self._detect_multi_row_headers(cleaned_table)
            data_rows = cleaned_table[data_start_idx:]
            
            print(f"🔍 DEBUG: Merged headers: {headers}")
            print(f"🔍 DEBUG: Data starts at row {data_start_idx}, {len(data_rows)} data rows")
            print(f"🔍 DEBUG: First data row: {data_rows[0] if data_rows else 'N/A'}")
            
            if not data_rows:
                print("🔍 DEBUG: No data rows found after header detection")
                return None
            
            # Create DataFrame with original headers
            df = pd.DataFrame(data_rows, columns=headers)
            
            # Store the original meaningful headers mapping before cleaning
            original_headers_map = {}
            for i, header in enumerate(headers):
                if not str(header).startswith('Column_') and str(header).strip() not in ['', 'None', 'nan']:
                    original_headers_map[i] = str(header).strip()
            
            print(f"🔍 DEBUG: Original meaningful headers map: {original_headers_map}")
            
            df = self._clean_dataframe(df, original_headers_map)
            
            print(f"🔍 DEBUG: Final DataFrame shape: {df.shape}")
            print(f"  Columns: {list(df.columns)}")
            print(f"  Non-empty columns: {sum(1 for col in df.columns if not df[col].isna().all())}")
            
            if df.empty:
                print("🔍 DEBUG: DataFrame is empty after cleaning")
                return None
            
            # Check if we have meaningful data
            non_empty_cols = sum(1 for col in df.columns if not df[col].isna().all())
            if non_empty_cols == 0:
                print("🔍 DEBUG: No non-empty columns found")
                return None
            
            # Validate that this is actually a structured table, not just text
            if not self._validate_table_structure(df):
                print("🔍 DEBUG: Failed table structure validation - likely not a real table")
                return None
            
            result = {
                'data': df,
                'page_number': page_num,
                'slide_number': None,
                'table_number': table_num,
                'bbox': None,
                'extraction_method': f'pdfplumber_{method}_direct',
                'confidence': 0.9 if non_empty_cols >= 2 else 0.7
            }
            
            print(f"🔍 DEBUG: Returning successful result with {non_empty_cols} meaningful columns")
            return result
            
        except Exception as e:
            print(f"🔍 DEBUG: Error in _process_table_data: {e}")
            logger.warning(f"Table processing failed: {e}")
            return None

    def _detect_multi_row_headers(self, cleaned_table: List[List]) -> Tuple[List[str], int]:
        """Detect and merge multi-row headers from a table"""
        if not cleaned_table or len(cleaned_table) < 2:
            return [], 0
        
        # Initialize header accumulator
        num_cols = len(cleaned_table[0])
        merged_headers = [''] * num_cols
        
        # Look at the first several rows to find header information
        data_start_idx = 1  # Default: assume first row is header
        header_rows_processed = 0
        
        print(f"🔍 DEBUG: Detecting headers from {len(cleaned_table)} rows, {num_cols} columns")
        
        for row_idx, row in enumerate(cleaned_table[:8]):  # Look at first 8 rows max
            # Ensure row has same number of columns
            while len(row) < num_cols:
                row.append('')
            
            # Check if this row contains header-like content
            meaningful_cells = 0
            header_like_words = 0
            
            for col_idx, cell in enumerate(row[:num_cols]):
                cell_str = str(cell).strip()
                if cell_str:
                    meaningful_cells += 1
                    
                    # Check if cell contains header-like words
                    if (cell_str and len(cell_str) > 1 and 
                        not cell_str.replace('.', '').replace(',', '').replace('%', '').replace(' ', '').isdigit() and
                        cell_str not in ['', '-', 'n/a', 'N/A']):
                        header_like_words += 1
                        
                        # Merge header information
                        if merged_headers[col_idx]:
                            merged_headers[col_idx] += f" {cell_str}"
                        else:
                            merged_headers[col_idx] = cell_str
            
            print(f"🔍 DEBUG: Row {row_idx}: {meaningful_cells} meaningful cells, {header_like_words} header-like words")
            print(f"  Row content: {row[:num_cols]}")
            
            # If this row has mostly data-like content (numbers, specific patterns), it's likely data
            numeric_like = 0
            for cell in row[:num_cols]:
                cell_str = str(cell).strip()
                if (cell_str and 
                    (cell_str.replace('.', '').replace(',', '').replace('%', '').replace(' ', '').isdigit() or
                     'sec' in cell_str.lower() or
                     cell_str.lower() in ['blind', 'low vision', 'dexterity', 'mobility'])):  # Common data patterns
                    numeric_like += 1
            
            print(f"  Numeric/data-like cells: {numeric_like}")
            
            # If we have many numeric/data-like cells, this is probably the start of data
            if numeric_like >= 2 and row_idx > 0:
                data_start_idx = row_idx
                print(f"🔍 DEBUG: Data appears to start at row {row_idx}")
                break
            
            header_rows_processed += 1
            
            # If we've processed several rows and haven't found clear data, use conservative approach
            if row_idx >= 6:
                data_start_idx = min(row_idx + 1, len(cleaned_table) - 1)
                break
        
        # Clean up merged headers
        final_headers = []
        for i, header in enumerate(merged_headers):
            header_str = str(header).strip()
            if not header_str:
                final_headers.append(f'Column_{i+1}')
            else:
                # Clean up header
                cleaned_header = re.sub(r'\s+', ' ', header_str)
                cleaned_header = cleaned_header.replace('\n', ' ').strip()
                final_headers.append(cleaned_header)
        
        print(f"🔍 DEBUG: Final merged headers: {final_headers}")
        print(f"🔍 DEBUG: Data starts at row {data_start_idx}")
        
        return final_headers, data_start_idx

    def _clean_table_data(self, table_data: List[List]) -> List[List]:
        """Clean raw table data from PDF extraction"""
        if not table_data:
            return []
        
        print(f"🔍 DEBUG: Cleaning table data with {len(table_data)} rows")
        cleaned_table = []
        
        for row_idx, row in enumerate(table_data):
            cleaned_row = []
            for cell in row:
                if cell is None:
                    cleaned_row.append("")
                else:
                    # Clean the cell content
                    cleaned_cell = str(cell).strip()
                    # Remove excessive whitespace
                    cleaned_cell = re.sub(r'\s+', ' ', cleaned_cell)
                    cleaned_row.append(cleaned_cell)
            
            # Only add row if it has some content
            if any(cell.strip() for cell in cleaned_row):
                cleaned_table.append(cleaned_row)
                print(f"  Cleaned row {row_idx}: {cleaned_row}")
        
        print(f"🔍 DEBUG: Cleaned to {len(cleaned_table)} rows with content")
        return cleaned_table

    def _clean_dataframe(self, df: pd.DataFrame, original_headers_map=None) -> pd.DataFrame:
        """Clean DataFrame by removing mostly empty columns and fixing headers"""
        if not isinstance(df, pd.DataFrame):
            logger.warning(f"Expected pandas DataFrame, got {type(df)}. Attempting to convert.")
            try:
                df = pd.DataFrame(df)
            except Exception as e:
                logger.error(f"Failed to convert to DataFrame: {e}")
                return pd.DataFrame()
        
        if df.empty:
            return df
        
        print(f"🔍 DEBUG: Cleaning DataFrame with shape {df.shape}")
        
        # Store original column names for mapping
        original_columns = list(df.columns)
        
        # Clean column names but preserve meaningful ones
        new_columns = []
        for i, col in enumerate(df.columns):
            col_str = str(col).strip()
            if pd.isna(col) or col_str == '' or col_str.lower() in ['none', 'nan'] or col_str.startswith('Column_'):
                new_columns.append(f'Column_{i+1}')
            else:
                # Clean column name but keep it meaningful
                cleaned_name = col_str
                cleaned_name = re.sub(r'\s+', ' ', cleaned_name)
                cleaned_name = cleaned_name.replace('\n', ' ').replace('\r', ' ')
                new_columns.append(cleaned_name)
        
        df.columns = new_columns
        print(f"🔍 DEBUG: Cleaned column names: {list(df.columns)}")
        
        # Clean cell values
        for col in df.columns:
            df[col] = df[col].astype(str).str.strip()
            df[col] = df[col].replace(['None', 'nan', ''], pd.NA)
        
        # Remove completely empty rows first
        before_empty_removal = len(df)
        df = df.dropna(how='all').reset_index(drop=True)
        print(f"🔍 DEBUG: Removed {before_empty_removal - len(df)} completely empty rows")
        
        # Check column quality and remove truly empty columns, but preserve meaningful names
        columns_to_keep = []
        original_indices_to_keep = []  # Track original indices of kept columns
        
        print(f"🔍 DEBUG: original_headers_map provided: {original_headers_map}")
        
        for i, col in enumerate(df.columns):
            non_empty_count = df[col].notna().sum()
            if non_empty_count > 0:
                columns_to_keep.append(col)
                original_indices_to_keep.append(i)  # Store the original index
                print(f"  Keeping column '{col}' (original index {i}): {non_empty_count} non-empty values")
        
        print(f"🔍 DEBUG: Original indices of kept columns: {original_indices_to_keep}")
        
        # Filter the DataFrame to keep only non-empty columns
        if columns_to_keep:
            df = df[columns_to_keep]
            print(f"🔍 DEBUG: After filtering, DataFrame columns: {list(df.columns)}")
            
            # Now create mapping from new filtered positions to meaningful names
            meaningful_column_mapping = {}
            for new_index, original_index in enumerate(original_indices_to_keep):
                current_col_name = df.columns[new_index]  # Current column name in filtered df
                
                # Check if we have a meaningful name for this original column position
                if original_headers_map and original_index in original_headers_map:
                    meaningful_name = original_headers_map[original_index]
                    meaningful_column_mapping[current_col_name] = meaningful_name
                    print(f"  📝 Will rename column at new index {new_index} ('{current_col_name}') to '{meaningful_name}' (from original index {original_index})")
                else:
                    print(f"  ❌ No meaningful name found for column at new index {new_index} (original index {original_index})")
            
            print(f"🔍 DEBUG: Column mapping to apply: {meaningful_column_mapping}")
            
            # Apply meaningful column names
            if meaningful_column_mapping:
                print(f"🔍 DEBUG: Before renaming: {list(df.columns)}")
                df = df.rename(columns=meaningful_column_mapping)
                print(f"🔍 DEBUG: After renaming: {list(df.columns)}")
                print(f"🔍 DEBUG: Applied meaningful column names: {meaningful_column_mapping}")
            else:
                print(f"🔍 DEBUG: No meaningful column mapping to apply")
            
            print(f"🔍 DEBUG: After filtering empty columns, kept {len(df.columns)} columns")
            print(f"  Final column names: {list(df.columns)}")
        else:
            print("🔍 DEBUG: No columns to keep - all columns are empty")
        
        # Additional cleaning: Remove rows that are mostly headers/junk
        meaningful_rows = []
        for idx, row in df.iterrows():
            # Count meaningful values (not just header words)
            meaningful_count = 0
            for val in row.values:
                if pd.notna(val) and str(val).strip():
                    val_str = str(val).strip()
                    # Any non-empty value is considered meaningful data
                    if len(val_str) > 0:
                        meaningful_count += 1
            
            # Keep rows with meaningful data (disability names, numbers, percentages, etc.)
            if meaningful_count >= 2:  # At least 2 meaningful values
                meaningful_rows.append(idx)
                print(f"  Keeping row {idx} with {meaningful_count} meaningful values: {[str(v)[:20] for v in row.values if pd.notna(v) and str(v).strip()]}")
        
        if meaningful_rows:
            df = df.iloc[meaningful_rows].reset_index(drop=True)
            print(f"🔍 DEBUG: Kept {len(meaningful_rows)} meaningful rows")
        
        print(f"🔍 DEBUG: Final cleaned DataFrame shape: {df.shape}")
        return df

    def _validate_table_structure(self, df: pd.DataFrame) -> bool:
        """Validate that the DataFrame represents a real table structure, not just text fragments"""
        print(f"\n🔍 DEBUG: Validating table structure for DataFrame with shape {df.shape}")
        
        if df.empty or len(df) < 2:
            print("🔍 DEBUG: Table too small to validate")
            return False
        
        # Check 1: Are cells too long? (indicates paragraph text, not table data)
        max_cell_length = 0
        avg_cell_length = 0
        total_cells = 0
        very_long_cells = 0
        
        for col in df.columns:
            for val in df[col].values:
                if pd.notna(val):
                    cell_str = str(val).strip()
                    if cell_str:
                        cell_len = len(cell_str)
                        max_cell_length = max(max_cell_length, cell_len)
                        avg_cell_length += cell_len
                        total_cells += 1
                        
                        # Count cells that are suspiciously long (paragraph-like)
                        if cell_len > 150:  # More than 150 chars is likely paragraph text
                            very_long_cells += 1
        
        if total_cells > 0:
            avg_cell_length /= total_cells
        
        print(f"  Max cell length: {max_cell_length}")
        print(f"  Avg cell length: {avg_cell_length:.1f}")
        print(f"  Very long cells (>150 chars): {very_long_cells}/{total_cells}")
        
        # If too many cells are very long, this is likely text, not a table
        if very_long_cells > 0 and (very_long_cells / total_cells) > 0.3:
            print("🔍 DEBUG: Too many very long cells - likely text paragraphs, not table")
            return False
        
        if max_cell_length > 500:  # Single cell over 500 chars is definitely paragraph text
            print("🔍 DEBUG: Cell too long - likely paragraph text, not table data")
            return False
        
        # Check 2: Do we have consistent data types/patterns across columns?
        structured_columns = 0
        for col in df.columns:
            non_empty_values = [str(val).strip() for val in df[col].values if pd.notna(val) and str(val).strip()]
            if len(non_empty_values) < 2:
                continue
                
            # Check if column has consistent short values (typical of table data)
            short_values = sum(1 for val in non_empty_values if len(val) <= 50)
            if short_values >= len(non_empty_values) * 0.7:  # 70% of values are short
                structured_columns += 1
                print(f"  Column '{col}' appears structured: {short_values}/{len(non_empty_values)} short values")
            else:
                print(f"  Column '{col}' appears unstructured: {short_values}/{len(non_empty_values)} short values")
        
        print(f"  Structured columns: {structured_columns}/{len(df.columns)}")
        
        # Need at least 2 columns that look structured
        if structured_columns < 2:
            print("🔍 DEBUG: Not enough structured columns - likely not a real table")
            return False
        
        # Check 3: Look for table-like patterns
        has_numeric_data = False
        has_short_categorical_data = False
        
        for col in df.columns:
            non_empty_values = [str(val).strip() for val in df[col].values if pd.notna(val) and str(val).strip()]
            if not non_empty_values:
                continue
            
            # Check for numeric patterns
            numeric_count = 0
            for val in non_empty_values:
                # Remove common formatting and check if numeric
                clean_val = re.sub(r'[,$%\s]', '', val)
                if clean_val and (clean_val.replace('.', '').isdigit() or 
                                 re.match(r'^\d+(\.\d+)?$', clean_val)):
                    numeric_count += 1
            
            if numeric_count >= len(non_empty_values) * 0.5:  # 50% numeric
                has_numeric_data = True
                print(f"  Column '{col}' has numeric data: {numeric_count}/{len(non_empty_values)}")
            
            # Check for short categorical data
            short_categorical = sum(1 for val in non_empty_values 
                                  if len(val) <= 30 and not val.replace('.', '').replace(',', '').isdigit())
            if short_categorical >= len(non_empty_values) * 0.7:
                has_short_categorical_data = True
                print(f"  Column '{col}' has categorical data: {short_categorical}/{len(non_empty_values)}")
        
        # A real table should have either numeric data or short categorical data
        if not (has_numeric_data or has_short_categorical_data):
            print("🔍 DEBUG: No numeric or categorical patterns found - likely not a table")
            return False
        
        print("🔍 DEBUG: Table structure validation passed")
        return True
    
    def _extract_pdf_tables_tabula(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract tables using tabula-py with multiple strategies"""
        tables = []
        
        try:
            # Strategy 1: Try lattice method (for tables with clear borders)
            try:
                dfs_lattice = tabula.read_pdf(
                    file_path, 
                    pages='all', 
                    multiple_tables=True,
                    lattice=True,
                    pandas_options={'header': 0}
                )
                
                for table_num, df in enumerate(dfs_lattice, 1):
                    if not df.empty:
                        df = self._clean_dataframe(df)
                        df = df.dropna(how='all').dropna(axis=1, how='all')
                        
                        if not df.empty and len(df.columns) >= 2:
                            tables.append({
                                'data': df,
                                'page_number': None,
                                'slide_number': None,
                                'table_number': table_num,
                                'bbox': None,
                                'extraction_method': 'tabula_lattice',
                                'confidence': 0.8
                            })
                
                if tables:
                    return tables
                    
            except Exception as e:
                logger.warning(f"tabula lattice method failed: {e}")
            
            # Strategy 2: Try stream method (for tables without clear borders)
            try:
                dfs_stream = tabula.read_pdf(
                    file_path, 
                    pages='all', 
                    multiple_tables=True,
                    stream=True,
                    pandas_options={'header': 0}
                )
                
                for table_num, df in enumerate(dfs_stream, 1):
                    if not df.empty:
                        df = self._clean_dataframe(df)
                        df = df.dropna(how='all').dropna(axis=1, how='all')
                        
                        if not df.empty and len(df.columns) >= 2:
                            tables.append({
                                'data': df,
                                'page_number': None,
                                'slide_number': None,
                                'table_number': table_num,
                                'bbox': None,
                                'extraction_method': 'tabula_stream',
                                'confidence': 0.7
                            })
                            
            except Exception as e:
                logger.warning(f"tabula stream method failed: {e}")
        
        except Exception as e:
            logger.error(f"tabula table extraction failed: {e}")
        
        return tables
    
    def _extract_pdf_tables_camelot(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract tables using camelot-py"""
        tables = []
        
        try:
            # Extract tables with lattice method (for tables with lines)
            lattice_tables = camelot.read_pdf(file_path, pages='all', flavor='lattice')
            
            for table in lattice_tables:
                if table.parsing_report['accuracy'] > 50:
                    df = table.df
                    if not df.empty:
                        tables.append({
                            'data': df,
                            'page_number': table.page,
                            'slide_number': None,
                            'table_number': len(tables) + 1,
                            'bbox': table._bbox,
                            'extraction_method': 'camelot_lattice',
                            'confidence': table.parsing_report['accuracy'] / 100
                        })
        
        except Exception as e:
            logger.error(f"camelot table extraction failed: {e}")
        
        return tables
    
    def extract_pptx_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract tables from PowerPoint files"""
        tables = []
        
        if not PPTX_AVAILABLE:
            return tables
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_tables = []
                
                for shape in slide.shapes:
                    if hasattr(shape, 'table') and shape.has_table:
                        table_data = []
                        
                        for row in shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text.strip())
                            table_data.append(row_data)
                        
                        if table_data and len(table_data) > 1:
                            try:
                                df = pd.DataFrame(table_data[1:], columns=table_data[0])
                                df = self._clean_dataframe(df)
                                
                                if not df.empty:
                                    tables.append({
                                        'data': df,
                                        'page_number': None,
                                        'slide_number': slide_num,
                                        'table_number': len(slide_tables) + 1,
                                        'bbox': None,
                                        'extraction_method': 'python_pptx',
                                        'confidence': 0.9
                                    })
                                    slide_tables.append(df)
                            except Exception as e:
                                logger.warning(f"Failed to process table in slide {slide_num}: {e}")
        
        except Exception as e:
            logger.error(f"PowerPoint table extraction failed: {e}")
        
        return tables
    
    def extract_docx_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract tables from Word documents"""
        tables = []
        
        if not DOCX_AVAILABLE:
            return tables
        
        try:
            doc = DocxDocument(file_path)
            
            for table_num, table in enumerate(doc.tables, 1):
                table_data = []
                
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                
                if table_data and len(table_data) > 1:
                    try:
                        df = pd.DataFrame(table_data[1:], columns=table_data[0])
                        df = self._clean_dataframe(df)
                        
                        if not df.empty:
                            tables.append({
                                'data': df,
                                'page_number': None,
                                'slide_number': None,
                                'table_number': table_num,
                                'bbox': None,
                                'extraction_method': 'python_docx',
                                'confidence': 0.9
                            })
                    except Exception as e:
                        logger.warning(f"Failed to process table {table_num}: {e}")
        
        except Exception as e:
            logger.error(f"Word document table extraction failed: {e}")
        
        return tables
    
    def extract_csv_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract tables from CSV files"""
        tables = []
        
        try:
            df = pd.read_csv(file_path)
            
            if not df.empty:
                tables.append({
                    'data': df,
                    'page_number': None,
                    'slide_number': None,
                    'table_number': 1,
                    'bbox': None,
                    'extraction_method': 'pandas_csv',
                    'confidence': 1.0
                })
        
        except Exception as e:
            logger.error(f"CSV table extraction failed: {e}")
        
        return tables
    
    def extract_excel_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract tables from Excel files"""
        tables = []
        
        try:
            excel_file = pd.ExcelFile(file_path)
            
            for sheet_num, sheet_name in enumerate(excel_file.sheet_names, 1):
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                if not df.empty:
                    tables.append({
                        'data': df,
                        'page_number': None,
                        'slide_number': sheet_num,
                        'table_number': 1,
                        'bbox': None,
                        'extraction_method': 'pandas_excel',
                        'confidence': 1.0,
                        'sheet_name': sheet_name
                    })
        
        except Exception as e:
            logger.error(f"Excel table extraction failed: {e}")
        
        return tables


def detect_table_in_text(text: str) -> bool:
    """
    Detect if text contains table-like structures
    """
    lines = text.strip().split('\n')
    
    if len(lines) < 2:
        return False
    
    # Check for common table indicators
    table_indicators = [
        # Multiple consecutive lines with delimiters
        lambda: sum(1 for line in lines if '|' in line or '\t' in line) >= 2,
        
        # Lines with consistent column structure
        lambda: _has_consistent_columns(lines),
        
        # Table headers pattern
        lambda: any(re.search(r'^[\w\s]+\s*\|\s*[\w\s]+\s*\|', line) for line in lines[:3]),
        
        # CSV-like structure
        lambda: sum(1 for line in lines if line.count(',') >= 2) >= 2,
        
        # Tab-separated values
        lambda: sum(1 for line in lines if line.count('\t') >= 1) >= 2,
    ]
    
    return any(indicator() for indicator in table_indicators)


def _has_consistent_columns(lines: List[str]) -> bool:
    """Check if lines have consistent column structure"""
    if len(lines) < 3:
        return False
    
    # Check for pipe-separated columns
    pipe_counts = [line.count('|') for line in lines[:5] if line.strip()]
    if len(set(pipe_counts)) <= 2 and pipe_counts and max(pipe_counts) >= 1:
        return True
    
    # Check for tab-separated columns
    tab_counts = [line.count('\t') for line in lines[:5] if line.strip()]
    if len(set(tab_counts)) <= 2 and tab_counts and max(tab_counts) >= 1:
        return True
    
    return False


def create_adaptive_table_chunks(table_data: pd.DataFrame, table_info: Dict[str, Any], 
                                max_chunk_size: int = 7800, preserve_context: bool = True) -> List[Dict[str, Any]]:
    """
    Create adaptive table chunks based on table size and content
    """
    print(f"\n🔍 DEBUG: create_adaptive_table_chunks called")
    print(f"  table_data shape: {table_data.shape if not table_data.empty else 'EMPTY'}")
    print(f"  table_info: {table_info}")
    print(f"  max_chunk_size: {max_chunk_size}")
    print(f"  preserve_context: {preserve_context}")
    
    if table_data.empty:
        print("🔍 DEBUG: Table data is empty, returning empty list")
        return []
    
    # Estimate size per row
    sample_row = table_data.iloc[0] if len(table_data) > 0 else pd.Series()
    estimated_row_size = sum(len(str(val)) for val in sample_row.values) + 20  # 20 for formatting
    estimated_row_size = max(estimated_row_size, 30)  # Minimum estimate
    
    # Calculate optimal chunk size - keep small to medium tables as single chunks
    if len(table_data) <= 10:  # Increased threshold for single chunk
        max_rows_per_chunk = len(table_data)
        print("🔍 DEBUG: Small/medium table detected, using single chunk")
    else:
        max_rows_per_chunk = max(5, min(20, max_chunk_size // estimated_row_size))
        print(f"🔍 DEBUG: Large table detected, using {max_rows_per_chunk} rows per chunk")
        print(f"  Estimated size per row: {estimated_row_size}")
    
    print("🔍 DEBUG: Calling create_semantic_table_chunks")
    return create_semantic_table_chunks(table_data, table_info, max_rows_per_chunk)


def create_semantic_table_chunks(table_data: pd.DataFrame, table_info: Dict[str, Any], 
                                max_rows_per_chunk: int = 10) -> List[Dict[str, Any]]:
    """
    Create semantically grouped table chunks
    """
    print(f"\n🔍 DEBUG: create_semantic_table_chunks called")
    print(f"  table_data shape: {table_data.shape if not table_data.empty else 'EMPTY'}")
    print(f"  table_info: {table_info}")
    print(f"  max_rows_per_chunk: {max_rows_per_chunk}")
    
    if table_data.empty:
        print("🔍 DEBUG: Table data is empty, returning empty list")
        return []
    
    chunks = []
    
    # Clean the DataFrame
    print("🔍 DEBUG: Cleaning DataFrame")
    cleaned_df = table_data.copy()
    cleaned_df = cleaned_df.fillna('')
    
    print(f"  After fillna: {cleaned_df.shape}")
    print(f"  Columns: {list(cleaned_df.columns)}")
    print(f"  Data sample:")
    
    # Print actual data rows to see what we're working with
    for i, (idx, row) in enumerate(cleaned_df.iterrows()):
        if i < 5:  # Show first 5 rows
            non_empty_values = [str(val) for val in row.values if val and str(val).strip()]
            print(f"    Row {i}: {non_empty_values}")
    
    print(f"  Total non-empty data rows: {len([idx for idx, row in cleaned_df.iterrows() if any(val and str(val).strip() for val in row.values)])}")
    
    # Convert all columns to string
    print("🔍 DEBUG: Converting columns to string")
    for col in cleaned_df.columns:
        cleaned_df[col] = cleaned_df[col].astype(str).str.strip()
        print(f"  Converted column '{col}'")
    
    print(f"🔍 DEBUG: Final cleaned DataFrame:")
    print(cleaned_df.head(3))
    
    # Get table title/identifier
    table_identifier = f"Table {table_info.get('table_number', 1)}"
    if table_info.get('page_number'):
        table_identifier += f" (Page {table_info['page_number']})"
    elif table_info.get('slide_number'):
        table_identifier += f" (Slide {table_info['slide_number']})"
    
    print(f"🔍 DEBUG: Table identifier: {table_identifier}")
    
    # Split into chunks
    total_rows = len(cleaned_df)
    print(f"🔍 DEBUG: Total rows: {total_rows}, max_rows_per_chunk: {max_rows_per_chunk}")
    
    # For small tables, use single chunk
    if max_rows_per_chunk >= total_rows:
        print("🔍 DEBUG: Creating single chunk for small table")
        chunk_content = create_single_table_chunk(cleaned_df, table_identifier)
        print(f"  Single chunk content length: {len(chunk_content)}")
        print(f"  Content preview: {chunk_content[:200]}...")
        
        chunk_metadata = create_table_chunk_metadata(table_info, 0)
        print(f"  Chunk metadata: {chunk_metadata}")
        
        chunks.append({
            'content': chunk_content,
            'metadata': chunk_metadata
        })
        
        print(f"🔍 DEBUG: Created 1 single chunk, returning")
        return chunks
    
    # Multiple chunks needed
    print("🔍 DEBUG: Creating multiple chunks")
    for chunk_start in range(0, total_rows, max_rows_per_chunk):
        chunk_end = min(chunk_start + max_rows_per_chunk, total_rows)
        chunk_df = cleaned_df.iloc[chunk_start:chunk_end]
        
        print(f"  Creating chunk for rows {chunk_start}-{chunk_end-1}")
        print(f"  Chunk df shape: {chunk_df.shape}")
        
        chunk_content = create_multi_table_chunk(
            chunk_df, table_identifier,
            chunk_start, chunk_end, total_rows
        )
        
        print(f"  Multi-chunk content length: {len(chunk_content)}")
        print(f"  Content preview: {chunk_content[:200]}...")
        
        # Create metadata
        chunk_metadata = create_table_chunk_metadata(table_info, chunk_index=chunk_start // max_rows_per_chunk)
        chunk_metadata.update({
            'chunk_type': 'table_group',
            'row_start': chunk_start,
            'row_end': chunk_end - 1,
            'row_count': len(chunk_df),
            'chunking_strategy': 'semantic',
            'formatting_style': 'structured'
        })
        
        print(f"  Multi-chunk metadata: {chunk_metadata}")
        
        chunks.append({
            'content': chunk_content,
            'metadata': chunk_metadata
        })
    
    print(f"🔍 DEBUG: Created {len(chunks)} chunks total")
    return chunks


def create_single_table_chunk(df: pd.DataFrame, table_identifier: str) -> str:
    """Create optimized single chunk for small tables"""
    print(f"\n🔍 DEBUG: create_single_table_chunk called")
    print(f"  df shape: {df.shape}")
    print(f"  table_identifier: {table_identifier}")
    
    lines = []
    
    # Add title and context
    lines.append(f"# {table_identifier}")
    lines.append("")
    lines.append("This is a data table extracted from the document.")
    lines.append("")
    
    # Use markdown table format for better readability
    if len(df.columns) <= 10:  # Most tables should fit this
        # Create header row with proper column names
        header_cells = []
        for col in df.columns:
            # Clean column names for better display
            clean_col = str(col).replace('\n', ' ').replace('\r', ' ').strip()
            header_cells.append(clean_col)
        
        header_row = "| " + " | ".join(header_cells) + " |"
        lines.append(header_row)
        
        # Create separator row
        separator_row = "|" + "|".join("-" * (len(cell) + 2) for cell in header_cells) + "|"
        lines.append(separator_row)
        
        # Add data rows with proper cell formatting
        for _, row in df.iterrows():
            data_cells = []
            for val in row.values:
                # Clean cell values and handle multi-line content
                clean_val = str(val).replace('\n', ' ').replace('\r', ' ').strip()
                if not clean_val or clean_val.lower() in ['nan', 'none', '']:
                    clean_val = "-"
                data_cells.append(clean_val)
            
            data_row = "| " + " | ".join(data_cells) + " |"
            lines.append(data_row)
    else:
        # For very wide tables, use key-value format
        lines.append("**Table Data:**")
        lines.append("")
        for row_idx, (_, row) in enumerate(df.iterrows(), 1):
            lines.append(f"**Row {row_idx}:**")
            for col, val in row.items():
                clean_val = str(val).strip()
                if clean_val and clean_val.lower() not in ['nan', 'none', '']:
                    lines.append(f"- {col}: {clean_val}")
            lines.append("")
    
    result = '\n'.join(lines)
    print(f"🔍 DEBUG: Single chunk result length: {len(result)}")
    print(f"  Preview: {result[:300]}...")
    
    return result


def create_multi_table_chunk(chunk_df: pd.DataFrame, table_identifier: str,
                           chunk_start: int, chunk_end: int, total_rows: int) -> str:
    """Create chunk for multi-chunk tables with context preservation"""
    print(f"\n🔍 DEBUG: create_multi_table_chunk called")
    print(f"  chunk_df shape: {chunk_df.shape}")
    print(f"  chunk_start: {chunk_start}, chunk_end: {chunk_end}, total_rows: {total_rows}")
    
    lines = []
    
    # Chunk title with context
    chunk_title = f"═══ {table_identifier} - Rows {chunk_start + 1} to {chunk_end} ═══"
    lines.append(chunk_title)
    lines.append("")
    
    # Add table context
    lines.append(f"📋 Table Context: {len(chunk_df.columns)} columns, {total_rows} total rows")
    lines.append(f"📊 This Chunk: {len(chunk_df)} rows ({chunk_start + 1}-{chunk_end})")
    lines.append("")
    
    # Use markdown table format
    if len(chunk_df.columns) <= 8:
        # Create header row
        header_row = "| " + " | ".join(str(col) for col in chunk_df.columns) + " |"
        lines.append(header_row)
        
        # Create separator row
        separator_row = "|" + "|".join("--------" for _ in chunk_df.columns) + "|"
        lines.append(separator_row)
        
        # Add data rows
        for _, row in chunk_df.iterrows():
            data_row = "| " + " | ".join(str(val) for val in row.values) + " |"
            lines.append(data_row)
    else:
        # For wide tables, use key-value format
        for row_idx, (_, row) in enumerate(chunk_df.iterrows()):
            lines.append(f"Row {chunk_start + row_idx + 1}:")
            for col, val in row.items():
                if str(val).strip():
                    lines.append(f"  {col}: {val}")
            lines.append("")
    
    result = '\n'.join(lines)
    print(f"🔍 DEBUG: Multi-chunk result length: {len(result)}")
    
    return result


def create_table_chunk_metadata(table_info: Dict[str, Any], chunk_index: int = 0) -> Dict[str, Any]:
    """Create metadata for table chunks"""
    metadata = {
        'chunk_type': 'table',
        'table_number': table_info.get('table_number', 1),
        'extraction_method': table_info.get('extraction_method', 'unknown'),
        'confidence': table_info.get('confidence', 0.8),
        'chunk_index': chunk_index
    }
    
    if table_info.get('page_number'):
        metadata['page_number'] = table_info['page_number']
    
    if table_info.get('slide_number'):
        metadata['slide_number'] = table_info['slide_number']
    
    return metadata


def identify_key_columns(df: pd.DataFrame) -> List[str]:
    """Identify key columns in a table"""
    key_columns = []
    
    for col in df.columns:
        # Check if column looks like an identifier or primary key
        col_lower = str(col).lower()
        if any(keyword in col_lower for keyword in ['id', 'name', 'title', 'category', 'type']):
            key_columns.append(col)
        
        # Check if column has unique or mostly unique values
        try:
            unique_ratio = len(df[col].unique()) / len(df)
            if unique_ratio > 0.8:  # More than 80% unique values
                key_columns.append(col)
        except:
            pass
    
    return key_columns[:3]  # Limit to top 3 key columns


# Global table extractor instance
_table_extractor_instance = None


def get_table_extractor() -> TableExtractor:
    """Get singleton table extractor instance"""
    global _table_extractor_instance
    if _table_extractor_instance is None:
        _table_extractor_instance = TableExtractor()
    return _table_extractor_instance


def create_contextual_table_chunks(table_data: pd.DataFrame, table_info: Dict[str, Any], 
                                  max_chunk_size: int = 7800, preserve_context: bool = True,
                                  surrounding_text: str = None) -> List[Dict[str, Any]]:
    """
    Create contextual table chunks (alias for create_adaptive_table_chunks for backward compatibility)
    
    Args:
        table_data: DataFrame containing the table data
        table_info: Metadata about the table
        max_chunk_size: Maximum size for each chunk
        preserve_context: Whether to preserve context across chunks
        surrounding_text: Text context surrounding the table (unused but kept for compatibility)
    """
    print(f"\n🔍 DEBUG: create_contextual_table_chunks called with surrounding_text: {bool(surrounding_text)}")
    return create_adaptive_table_chunks(table_data, table_info, max_chunk_size, preserve_context)
