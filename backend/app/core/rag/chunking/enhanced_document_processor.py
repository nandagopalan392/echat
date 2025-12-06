"""
Enhanced Document Processor with Docling Integration
Handles different document types and chunking strategies with advanced document parsing
"""

import os
import logging
import time
import traceback

from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import json
import re
from dataclasses import dataclass

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter, TokenTextSplitter

from app.config.chunking import ChunkingMethod, ChunkingConfig, get_chunking_config_manager
from app.core.rag.chunking.table_extraction import (
    get_table_extractor, detect_table_in_text, 
    create_table_chunk_metadata, create_adaptive_table_chunks, 
    create_contextual_table_chunks
)
from app.utils.gpu_utils import configure_docling_device, is_cuda_memory_error, clear_gpu_memory

logger = logging.getLogger(__name__)

# Try to import Docling for advanced document processing
try:
    from docling.document_converter import DocumentConverter
    from docling.datamodel.pipeline_options import PipelineOptions
    from docling.backend.pypdfium2_backend import PyPdfiumDocumentBackend
    DOCLING_AVAILABLE = True
    logger.info("Docling available for advanced document processing")
except ImportError as e:
    logger.warning(f"Docling not available. Falling back to basic document loaders. Error: {e}")
    DOCLING_AVAILABLE = False

# Fallback document loaders
try:
    from langchain_community.document_loaders import PyPDFLoader
    from langchain_community.document_loaders import Docx2txtLoader, CSVLoader, TextLoader
    from langchain_community.document_loaders import UnstructuredPowerPointLoader
    ADVANCED_LOADERS_AVAILABLE = True
except ImportError:
    logger.warning("Advanced document loaders not available. Some features may be limited.")
    ADVANCED_LOADERS_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    logger.warning("Pandas not available. Table processing may be limited.")
    PANDAS_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    logger.warning("OCR libraries not available. Image text extraction disabled.")
    OCR_AVAILABLE = False

@dataclass
class ChunkingResult:
    """Result of document chunking operation"""
    chunks: List[Document]
    metadata: Dict[str, Any]
    method_used: ChunkingMethod
    config_used: ChunkingConfig
    warnings: Optional[List[str]] = None
    
    def __post_init__(self):
        if self.warnings is None:
            self.warnings = []

class DoclingDocumentProcessor:
    """Docling-based document processor for advanced document parsing"""
    
    def __init__(self):
        if DOCLING_AVAILABLE:
            try:
                start_time = time.time()
                logger.info("🚀 DEBUG: Starting Docling DocumentConverter initialization...")
                
                # Smart GPU/CPU configuration based on memory availability
                use_gpu, device_status = configure_docling_device()
                logger.info(f"🚀 DEBUG: Device configuration: {device_status}")
                
                # Create DocumentConverter with appropriate configuration
                logger.info("🚀 DEBUG: Creating DocumentConverter...")
                init_start = time.time()
                
                try:
                    # For Docling 2.43.0, use simple initialization
                    # VLM pipeline has CUDA device issues, so we'll use standard pipeline
                    # Images will fall back to OCR which works reliably
                    if not use_gpu:
                        # Force CPU if GPU memory is insufficient
                        import os
                        os.environ['CUDA_VISIBLE_DEVICES'] = ''
                        logger.info("🚀 DEBUG: Forcing CPU mode via environment variables")
                    else:
                        logger.info("🚀 DEBUG: Using GPU acceleration for standard pipeline")
                    
                    logger.info("🚀 DEBUG: Creating DocumentConverter with standard pipeline...")
                    self.converter = DocumentConverter()
                    device_type = "GPU" if use_gpu else "CPU"
                    logger.info(f"🚀 DEBUG: DocumentConverter created with standard pipeline ({device_type})")
                    
                except ImportError:
                    # Fallback to simple initialization if VLM imports fail
                    logger.info("🚀 DEBUG: VLM configuration not available, using simple initialization...")
                    if not use_gpu:
                        import os
                        os.environ['CUDA_VISIBLE_DEVICES'] = ''
                        logger.info("🚀 DEBUG: Forcing CPU mode via environment variables")
                    
                    self.converter = DocumentConverter()
                    device_type = "GPU" if use_gpu else "CPU"
                    logger.info(f"🚀 DEBUG: DocumentConverter created with default settings ({device_type})")
                
                init_time = time.time() - init_start
                logger.info(f"🚀 DEBUG: DocumentConverter created in {init_time:.2f} seconds")
                
                # Supported formats by Docling (new API uses file extensions directly)
                # Excluding image formats due to VLM CUDA device issues - images will use OCR fallback
                self.docling_formats = {'.pdf', '.docx', '.pptx', '.html', '.md', '.txt'}
                
                total_time = time.time() - start_time
                device_type = "GPU" if use_gpu else "CPU"
                logger.info(f"Docling DocumentConverter initialized with {device_type} in {total_time:.2f} seconds")
                
                # PRODUCTION-GRADE: Pre-load models during initialization (eager loading)
                logger.info("🚀 WARMUP: Pre-loading Docling models (downloading if needed)...")
                self._warmup_models()
                
            except Exception as e:
                logger.warning(f"Failed to initialize Docling converter: {e}")
                logger.warning(f"🚀 DEBUG: Full traceback: {traceback.format_exc()}")
                self.converter = None
                self.docling_formats = set()
        else:
            self.converter = None
            self.docling_formats = set()
    
    def _warmup_models(self):
        """
        Warm up Docling models by processing a minimal dummy document.
        This triggers model downloads during initialization instead of first request.
        PRODUCTION-GRADE PATTERN: Eager loading for predictable performance.
        """
        if not self.converter:
            return
            
        try:
            warmup_start = time.time()
            
            # Create a minimal PDF in memory to trigger model loading
            import io
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            # Generate minimal PDF
            pdf_buffer = io.BytesIO()
            c = canvas.Canvas(pdf_buffer, pagesize=letter)
            c.drawString(100, 750, "Warmup")
            c.save()
            pdf_buffer.seek(0)
            
            # Save to temp file for processing
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
                temp_pdf.write(pdf_buffer.getvalue())
                temp_path = temp_pdf.name
            
            try:
                # Process dummy document - this triggers model download
                logger.info("🚀 WARMUP: Processing warmup document to download models...")
                from docling.datamodel.base_models import InputFormat
                from docling.document_converter import DocumentConverter, PdfFormatOption
                from docling.datamodel.pipeline_options import PdfPipelineOptions
                
                # Convert the warmup PDF
                result = self.converter.convert(temp_path)
                
                warmup_time = time.time() - warmup_start
                logger.info(f"✅ WARMUP COMPLETE: Docling models pre-loaded in {warmup_time:.2f} seconds")
                logger.info("✅ All subsequent document processing will be fast (models already loaded)")
                
            finally:
                # Clean up temp file
                import os
                try:
                    os.unlink(temp_path)
                except:
                    pass
                    
        except ImportError as ie:
            logger.warning(f"⚠️ WARMUP: reportlab not available, skipping model warmup: {ie}")
            logger.warning("⚠️ Models will download on first document processing (7+ minute delay)")
        except Exception as e:
            logger.warning(f"⚠️ WARMUP: Failed to warm up models: {e}")
            logger.warning("⚠️ Models will download on first document processing (7+ minute delay)")
    
    def can_process(self, file_path: str) -> bool:
        """Check if Docling can process this file format"""
        if not DOCLING_AVAILABLE or not self.converter:
            return False
        
        ext = Path(file_path).suffix.lower()
        return ext in self.docling_formats
    
    def process_document(self, file_path: str) -> List[Document]:
        """Process document using Docling for comprehensive content extraction"""
        
        if not self.can_process(file_path):
            raise ValueError(f"Cannot process {file_path} with Docling")
        
        logger.info(f"🚀 DEBUG: Starting Docling processing for: {file_path}")
        total_start = time.time()
        
        # First attempt with current configuration
        try:
            return self._process_document_internal(file_path, total_start)
        except Exception as e:
            # Check if it's a CUDA memory error
            if is_cuda_memory_error(e):
                logger.warning(f"🔄 CUDA memory error detected, switching to CPU-only mode: {e}")
                
                # Clear GPU memory cache
                clear_gpu_memory()
                
                # Force CPU mode and reinitialize converter
                self._switch_to_cpu_mode()
                
                # Retry with CPU
                try:
                    return self._process_document_internal(file_path, total_start, retry=True)
                except Exception as retry_error:
                    total_time = time.time() - total_start
                    logger.error(f"🚀 DEBUG: Docling processing failed even with CPU mode after {total_time:.2f} seconds: {retry_error}")
                    raise retry_error
            else:
                # Re-raise non-memory errors
                raise e
    
    def _switch_to_cpu_mode(self):
        """Switch Docling converter to CPU-only mode"""
        try:
            logger.info("🔄 Switching Docling to CPU-only mode...")
            
            # Force CPU configuration
            os.environ['CUDA_VISIBLE_DEVICES'] = ''
            
            # Reinitialize converter with CPU settings
            init_start = time.time()
            
            try:
                # Simple CPU-only reinitialization (avoiding VLM CUDA device issues)
                logger.info("🔄 Reinitializing DocumentConverter with CPU-only mode...")
                self.converter = DocumentConverter()
                logger.info("🔄 DocumentConverter reinitialized with standard pipeline (CPU-only)")
                
                init_time = time.time() - init_start
                logger.info(f"🔄 CPU-only converter ready in {init_time:.2f} seconds")
                
            except Exception as reinit_error:
                logger.error(f"Failed to reinitialize converter in CPU mode: {reinit_error}")
                raise
            
        except Exception as e:
            logger.error(f"Failed to switch to CPU mode: {e}")
            raise
    
    def _process_document_internal(self, file_path: str, total_start: float, retry: bool = False) -> List[Document]:
        """Internal document processing logic"""
        
        try:
            # Convert document using Docling (v2.43.0 simplified API)
            logger.info("🚀 DEBUG: Starting document conversion...")
            convert_start = time.time()
            conversion_results = list(self.converter.convert(file_path))
            convert_time = time.time() - convert_start
            logger.info(f"🚀 DEBUG: Document conversion completed in {convert_time:.2f} seconds")
            
            if not conversion_results:
                raise ValueError(f"No conversion results for {file_path}")
            
            # Get the first (and typically only) result
            logger.info("🚀 DEBUG: Extracting conversion result...")
            result_start = time.time()
            documents = []
            
            # Extract content from Docling result (v2.43.0 returns tuples)
            # Find the document tuple in the results
            doc_content = None
            for result_tuple in conversion_results:
                if isinstance(result_tuple, tuple) and len(result_tuple) == 2:
                    key, value = result_tuple
                    if key == 'document':
                        doc_content = value
                        break
            
            if not doc_content:
                raise ValueError("Document content not found in conversion results")
            
            result_time = time.time() - result_start
            logger.info(f"🚀 DEBUG: Result extraction completed in {result_time:.2f} seconds")
            
            # Process tables first (if any)
            logger.info("🚀 DEBUG: Starting table extraction...")
            table_start = time.time()
            table_docs = self._extract_tables_from_docling(doc_content, file_path)
            table_time = time.time() - table_start
            logger.info(f"🚀 DEBUG: Table extraction completed in {table_time:.2f} seconds, found {len(table_docs)} tables")
            documents.extend(table_docs)
            
            # Process text content
            logger.info("🚀 DEBUG: Starting text extraction...")
            text_start = time.time()
            text_docs = self._extract_text_from_docling(doc_content, file_path)
            text_time = time.time() - text_start
            logger.info(f"🚀 DEBUG: Text extraction completed in {text_time:.2f} seconds, found {len(text_docs)} text sections")
            documents.extend(text_docs)
            
            # Process images with OCR (if any)
            logger.info("🚀 DEBUG: Starting image extraction...")
            image_start = time.time()
            image_docs = self._extract_images_from_docling(doc_content, file_path)
            image_time = time.time() - image_start
            logger.info(f"🚀 DEBUG: Image extraction completed in {image_time:.2f} seconds, found {len(image_docs)} image sections")
            documents.extend(image_docs)
            
            total_time = time.time() - total_start
            logger.info(f"🚀 DEBUG: Total processing time: {total_time:.2f} seconds")
            logger.info(f"🚀 DEBUG: Breakdown - Convert: {convert_time:.2f}s, Extract: {result_time:.2f}s, Tables: {table_time:.2f}s, Text: {text_time:.2f}s, Images: {image_time:.2f}s")
            logger.info(f"Docling processed {file_path}: {len(documents)} document sections extracted")
            return documents
            
        except Exception as e:
            total_time = time.time() - total_start
            logger.error(f"🚀 DEBUG: Docling processing failed after {total_time:.2f} seconds for {file_path}: {e}")
            logger.error(f"🚀 DEBUG: Full processing traceback: {traceback.format_exc()}")
            raise
    
    def _extract_tables_from_docling(self, doc_content, file_path: str) -> List[Document]:
        """Extract tables from Docling document"""
        documents = []
        
        try:
            # Access tables from Docling document structure (v2.43.0)
            if hasattr(doc_content, 'tables') and doc_content.tables:
                for i, table in enumerate(doc_content.tables):
                    try:
                        # Use export_to_dataframe method for Docling v2.43.0
                        if hasattr(table, 'export_to_dataframe'):
                            df = table.export_to_dataframe()
                            
                            if df is not None and not df.empty:
                                # Create table document with enhanced metadata
                                table_content = df.to_string(index=False)
                                
                                # Add table context and structure information
                                context_info = f"Table {i+1} from document"
                                
                                metadata = {
                                    'source': Path(file_path).name,
                                    'content_type': 'table',
                                    'extracted_from': f'docling_{Path(file_path).suffix[1:]}',
                                    'table_index': i,
                                    'table_rows': len(df),
                                    'table_cols': len(df.columns),
                                    'context_info': context_info,
                                    'extraction_method': 'docling_v2.43.0',
                                    'confidence_score': 0.95,  # High confidence for Docling extraction
                                    'docling_version': '2.43.0'
                                }
                                
                                doc = Document(
                                    page_content=table_content,
                                    metadata=metadata
                                )
                                documents.append(doc)
                                
                        # Fallback: use text representation if available
                        elif hasattr(table, 'text') and table.text:
                            table_content = table.text
                            
                            metadata = {
                                'source': Path(file_path).name,
                                'content_type': 'table',
                                'extracted_from': f'docling_{Path(file_path).suffix[1:]}',
                                'table_index': i,
                                'context_info': f"Table {i+1} from document (text format)",
                                'extraction_method': 'docling_v2.43.0_text',
                                'confidence_score': 0.8,
                                'docling_version': '2.43.0'
                            }
                            
                            doc = Document(
                                page_content=table_content,
                                metadata=metadata
                            )
                            documents.append(doc)
                            
                    except Exception as e:
                        logger.warning(f"Failed to process table {i} from Docling: {e}")
                        
        except Exception as e:
            logger.warning(f"Table extraction from Docling failed: {e}")
        
        return documents
    
    def _extract_text_from_docling(self, doc_content, file_path: str) -> List[Document]:
        """Extract text content from Docling document"""
        documents = []
        
        try:
            # Use export_to_markdown for full document text (Docling v2.43.0)
            if hasattr(doc_content, 'export_to_markdown'):
                try:
                    text_content = doc_content.export_to_markdown()
                    
                    if text_content and text_content.strip():
                        metadata = {
                            'source': Path(file_path).name,
                            'content_type': 'text',
                            'extracted_from': f'docling_{Path(file_path).suffix[1:]}',
                            'extraction_method': 'docling_markdown',
                            'docling_version': '2.43.0'
                        }
                        
                        doc = Document(
                            page_content=text_content,
                            metadata=metadata
                        )
                        documents.append(doc)
                        
                except Exception as e:
                    logger.warning(f"Markdown export failed: {e}")
            
            # Fallback: Extract from main_text list (Docling v2.43.0 structure)
            elif hasattr(doc_content, 'main_text') and doc_content.main_text:
                text_parts = []
                for text_obj in doc_content.main_text:
                    if hasattr(text_obj, 'text') and text_obj.text:
                        text_parts.append(text_obj.text)
                
                if text_parts:
                    text_content = '\n'.join(text_parts)
                    
                    metadata = {
                        'source': Path(file_path).name,
                        'content_type': 'text',
                        'extracted_from': f'docling_{Path(file_path).suffix[1:]}',
                        'extraction_method': 'docling_main_text',
                        'docling_version': '2.43.0',
                        'text_segments': len(text_parts)
                    }
                    
                    doc = Document(
                        page_content=text_content,
                        metadata=metadata
                    )
                    documents.append(doc)
                        
        except Exception as e:
            logger.warning(f"Text extraction from Docling failed: {e}")
        
        return documents
    
    def _extract_images_from_docling(self, doc_content, file_path: str) -> List[Document]:
        """Extract and OCR images from Docling document"""
        documents = []
        
        try:
            if hasattr(doc_content, 'images') and doc_content.images and OCR_AVAILABLE:
                for i, image in enumerate(doc_content.images):
                    try:
                        # OCR image content if available
                        if hasattr(image, 'pil_image'):
                            ocr_text = pytesseract.image_to_string(image.pil_image)
                            
                            if ocr_text.strip():
                                metadata = {
                                    'source': Path(file_path).name,
                                    'content_type': 'image_text',
                                    'extracted_from': f'docling_{Path(file_path).suffix[1:]}',
                                    'extraction_method': 'docling_ocr',
                                    'image_index': i
                                }
                                
                                doc = Document(
                                    page_content=ocr_text,
                                    metadata=metadata
                                )
                                documents.append(doc)
                                
                    except Exception as e:
                        logger.warning(f"OCR failed for image {i}: {e}")
                        
        except Exception as e:
            logger.warning(f"Image extraction from Docling failed: {e}")
        
        return documents

class EnhancedDocumentProcessor:
    """Enhanced document processor with Docling integration and multiple chunking methods"""
    
    def __init__(self):
        self.config_manager = get_chunking_config_manager()
        self.docling_processor = DoclingDocumentProcessor()
        
        # Fallback processors for formats not supported by Docling
        self.fallback_processors = {
            '.pdf': self._process_pdf_fallback,
            '.docx': self._process_docx_fallback,
            '.doc': self._process_docx_fallback,
            '.txt': self._process_text_fallback,
            '.md': self._process_text_fallback,
            '.csv': self._process_csv,
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,
            '.ppt': self._process_presentation,
            '.pptx': self._process_pptx_fallback,
            '.html': self._process_html_fallback,
            '.json': self._process_json,
            '.eml': self._process_email,
            '.jpg': self._process_image,
            '.jpeg': self._process_image,
            '.png': self._process_image,
            '.gif': self._process_image,
            '.tif': self._process_image,
            '.tiff': self._process_image
        }
        
        # All supported extensions (Docling + fallback)
        self.supported_extensions = set(self.fallback_processors.keys())
        if DOCLING_AVAILABLE:
            self.supported_extensions.update(self.docling_processor.docling_formats)
    
    def process_document(self, file_path: str, method: ChunkingMethod = None, 
                        config: ChunkingConfig = None, user_id: str = None, 
                        original_filename: str = None, document_id: int = None) -> ChunkingResult:
        """
        Process document with Docling-first approach and specified chunking method
        """
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext not in self.supported_extensions:
            raise ValueError(f"Unsupported file format: {file_ext}")
        
        # Use original filename if provided, otherwise use the file path name
        # Always extract just the filename (no path) for consistent metadata storage
        if original_filename:
            source_filename = Path(original_filename).name
        else:
            source_filename = Path(file_path).name
            
        if method is None:
            from app.config.chunking import FileFormatSupport
            method = FileFormatSupport.get_optimal_method(file_ext[1:])  # Remove dot
        
        # Get configuration if not specified
        if config is None:
            config = self.config_manager.get_config(method, user_id)
        
        # Validate configuration
        warnings = self.config_manager.validate_config(config)
        
        logger.info(f"Processing {file_path} with method {method.value}")
        
        # Try Docling first for supported formats
        raw_documents = []
        docling_used = False
        
        if self.docling_processor.can_process(file_path):
            try:
                raw_documents = self.docling_processor.process_document(file_path)
                docling_used = True
                logger.info(f"Successfully processed {file_path} with Docling")
            except Exception as e:
                logger.warning(f"Docling processing failed for {file_path}: {e}. Falling back to traditional methods.")
        
        # Fallback to traditional processors if Docling failed or not available
        if not raw_documents:
            if file_ext in self.fallback_processors:
                loader_func = self.fallback_processors[file_ext]
                raw_documents = loader_func(file_path)
            else:
                raise ValueError(f"No processor available for {file_ext}")
        
        if not raw_documents:
            raise ValueError(f"No content extracted from {file_path}")
        
        # Apply chunking method
        chunks = self._apply_chunking_method(raw_documents, method, config, docling_used)
        
        # Add metadata to chunks - clear existing metadata first to avoid complex objects
        for i, chunk in enumerate(chunks):
            # Preserve existing metadata and add processing info
            existing_metadata = chunk.metadata.copy() if chunk.metadata else {}
            chunk.metadata = {  # Replace instead of update to avoid complex objects
                'source_file': source_filename,
                'chunk_index': i,
                'chunking_method': method.value,
                'total_chunks': len(chunks),
                'processed_with_docling': docling_used
            }
            # Add back important existing metadata
            for key in ['content_type', 'extracted_from', 'table_index', 'page_number', 'extraction_method']:
                if key in existing_metadata:
                    chunk.metadata[key] = existing_metadata[key]
            
            # Add document_id if provided (needed for deletion functionality)
            if document_id is not None:
                chunk.metadata['document_id'] = str(document_id)
        
        # Debug: Log metadata after setting
        if chunks:
            logger.info(f"Enhanced processor - chunk metadata set: {chunks[0].metadata}")
        
        # Create processing metadata
        metadata = {
            'source_file': source_filename,
            'file_size': os.path.getsize(file_path),
            'chunk_count': len(chunks),
            'method_used': method.value,
            'config_used': config.to_dict(),
            'total_tokens': sum(len(chunk.page_content.split()) for chunk in chunks),
            'processed_with_docling': docling_used,
            'processor_version': 'docling_enhanced_v2.0'
        }
        
        return ChunkingResult(
            chunks=chunks,
            metadata=metadata,
            method_used=method,
            config_used=config,
            warnings=warnings
        )
    
    def _apply_chunking_method(self, documents: List[Document], method: ChunkingMethod, 
                              config: ChunkingConfig, docling_used: bool = False) -> List[Document]:
        """Apply specific chunking method to documents with Docling-aware processing"""
        
        all_chunks = []
        
        # If documents were processed with Docling, they already have proper structure
        if docling_used:
            logger.info("Processing Docling-extracted documents with enhanced chunking")
            
            # Separate tables and text for different processing
            table_docs = [doc for doc in documents if doc.metadata.get('content_type') == 'table']
            text_docs = [doc for doc in documents if doc.metadata.get('content_type') in ['text', 'image_text']]
            
            # Process tables with table-specific chunking
            for table_doc in table_docs:
                # Tables from Docling are already well-structured, just apply method-specific processing
                if method in [ChunkingMethod.TABLE, ChunkingMethod.GENERAL]:
                    # Keep table structure intact for table-focused methods
                    all_chunks.append(table_doc)
                else:
                    # For other methods, apply standard chunking to table content
                    table_chunks = self._apply_method_to_documents([table_doc], method, config)
                    all_chunks.extend(table_chunks)
            
            # Process text documents
            for text_doc in text_docs:
                text_chunks = self._apply_method_to_documents([text_doc], method, config)
                all_chunks.extend(text_chunks)
                
        else:
            # Fallback processing - use existing table-aware logic
            table_extractor = get_table_extractor()
            
            # Track which source files already had tables extracted during document loading
            sources_with_tables = set()
            table_documents = []
            text_documents = []
            
            # Separate table and text documents, and track sources with tables
            for doc in documents:
                if doc.metadata.get('content_type') == 'table' and doc.metadata.get('extracted_from'):
                    table_documents.append(doc)
                    source_file = doc.metadata.get('source', '')
                    if source_file:
                        sources_with_tables.add(source_file)
                else:
                    text_documents.append(doc)
            
            # Process table documents directly (they're already processed)
            logger.info(f"Found {len(table_documents)} pre-processed table documents")
            for table_doc in table_documents:
                all_chunks.append(table_doc)
            
            # Process text documents, but skip table extraction if tables were already extracted for this file
            logger.info(f"Processing {len(text_documents)} text documents")
            for doc in text_documents:
                source_file = doc.metadata.get('source', '')
                
                # Skip table extraction if we already have tables from this source file
                if source_file in sources_with_tables:
                    logger.info(f"Skipping table extraction for {source_file} - tables already extracted during document loading")
                    # Process only as text content
                    method_chunks = self._apply_method_to_documents([doc], method, config)
                    all_chunks.extend(method_chunks)
                    continue
                
                # For sources without pre-extracted tables, attempt table extraction
                doc_chunks = []
                extracted_tables = []
                if source_file and Path(source_file).exists():
                    try:
                        extracted_tables = table_extractor.extract_tables(str(source_file))
                        if extracted_tables:
                            logger.info(f"Extracted {len(extracted_tables)} tables from {source_file} during chunking")
                            # Mark this source as having tables to avoid re-extraction
                            sources_with_tables.add(source_file)
                    except Exception as e:
                        logger.warning(f"Table extraction failed for {source_file}: {e}")
                
                # Process extracted tables using enhanced table-aware chunking
                table_chunks_created = False
                for table_info in extracted_tables:
                    table_data = table_info['data']
                    
                    try:
                        if (PANDAS_AVAILABLE and 
                            isinstance(table_data, pd.DataFrame) and 
                            not table_data.empty):
                            
                            max_chunk_size = config.chunk_token_num if config.chunk_token_num else 1000
                            surrounding_text = doc.page_content[:500] + doc.page_content[-500:]
                            
                            table_chunks = create_contextual_table_chunks(
                                table_data, 
                                table_info, 
                                surrounding_text=surrounding_text,
                                max_chunk_size=max_chunk_size
                            )
                            
                            for chunk_data in table_chunks:
                                chunk_metadata = chunk_data['metadata']
                                chunk_metadata['chunking_method'] = f"{method.value}_table_aware"
                                chunk_metadata.update(doc.metadata)
                                
                                chunk = Document(
                                    page_content=chunk_data['content'],
                                    metadata=chunk_metadata
                                )
                                doc_chunks.append(chunk)
                            
                            table_chunks_created = True
                    except Exception as e:
                        logger.warning(f"Failed to process table data: {e}")
                
                # Add table chunks if they were created
                if table_chunks_created:
                    logger.info(f"Detected table content in document, adding {len(doc_chunks)} table chunks")
                    all_chunks.extend(doc_chunks)
                
                # Always process text content as well (unless document is purely tables)
                if doc.metadata.get('content_type') != 'table':
                    method_chunks = self._apply_method_to_documents([doc], method, config)
                    all_chunks.extend(method_chunks)
                else:
                    logger.info(f"Skipping text processing for pure table document")
        
        return all_chunks
    
    def _apply_method_to_documents(self, documents: List[Document], method: ChunkingMethod, 
                                  config: ChunkingConfig) -> List[Document]:
        """Apply specific chunking method to documents"""
        
        if method == ChunkingMethod.GENERAL:
            return self._chunk_general_table_aware(documents, config)
        elif method == ChunkingMethod.QA:
            return self._chunk_qa_table_aware(documents, config)
        elif method == ChunkingMethod.RESUME:
            return self._chunk_resume_table_aware(documents, config)
        elif method == ChunkingMethod.TABLE:
            return self._chunk_table(documents, config)
        elif method == ChunkingMethod.PRESENTATION:
            return self._chunk_presentation(documents, config)
        elif method == ChunkingMethod.PICTURE:
            return self._chunk_picture(documents, config)
        elif method == ChunkingMethod.EMAIL:
            return self._chunk_email(documents, config)
        else:
            # Default to general chunking with table awareness
            return self._chunk_general_table_aware(documents, config)
    
    def _chunk_general(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Standard recursive character text splitting"""
        # Parse delimiter string into list of separators
        separators = config.delimiter.split('|') if '|' in config.delimiter else [config.delimiter]
        separators = [sep.replace('\\n', '\n').replace('\\t', '\t') for sep in separators]
        
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=config.chunk_token_num,
            chunk_overlap=config.chunk_overlap,
            separators=separators,
            length_function=len
        )
        
        return splitter.split_documents(documents)
    
    def _chunk_qa(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Q&A format chunking - split on question/answer patterns"""
        chunks = []
        qa_patterns = [r'Q:', r'A:', r'Question:', r'Answer:', r'\d+\.', r'Q\d+', r'A\d+']
        
        for doc in documents:
            content = doc.page_content
            
            # Find Q&A boundaries
            boundaries = []
            for pattern in qa_patterns:
                for match in re.finditer(pattern, content, re.IGNORECASE):
                    boundaries.append(match.start())
            
            boundaries = sorted(set(boundaries))
            boundaries.append(len(content))
            
            # Create chunks based on Q&A boundaries
            for i in range(len(boundaries) - 1):
                start, end = boundaries[i], boundaries[i + 1]
                chunk_text = content[start:end].strip()
                
                if len(chunk_text) > 50:  # Minimum chunk size
                    chunk = Document(
                        page_content=chunk_text,
                        metadata={'qa_chunk': True}  # Only simple metadata, no complex inheritance
                    )
                    chunks.append(chunk)
        
        return chunks
    
    def _chunk_resume(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Resume-specific chunking based on sections"""
        chunks = []
        section_patterns = [
            r'EXPERIENCE|WORK EXPERIENCE|EMPLOYMENT',
            r'EDUCATION|ACADEMIC',
            r'SKILLS|TECHNICAL SKILLS|COMPETENCIES',
            r'PROJECTS|PROJECT EXPERIENCE',
            r'CERTIFICATIONS|CERTIFICATES',
            r'SUMMARY|OBJECTIVE|PROFILE'
        ]
        
        for doc in documents:
            content = doc.page_content
            
            # Find section boundaries
            boundaries = [0]  # Start with beginning
            for pattern in section_patterns:
                for match in re.finditer(pattern, content, re.IGNORECASE):
                    boundaries.append(match.start())
            
            boundaries = sorted(set(boundaries))
            boundaries.append(len(content))
            
            # Create chunks for each section
            for i in range(len(boundaries) - 1):
                start, end = boundaries[i], boundaries[i + 1]
                chunk_text = content[start:end].strip()
                
                if len(chunk_text) > 100:  # Minimum meaningful section size
                    chunk = Document(
                        page_content=chunk_text,
                        metadata={'resume_section': True}  # Only simple metadata
                    )
                    chunks.append(chunk)
        
        return chunks
    

    
    def _chunk_table(self, documents: List[Document], config: ChunkingConfig, 
                     pre_extracted_tables: List[Dict] = None) -> List[Document]:
        """Table-aware chunking for structured data"""
        chunks = []
        table_extractor = get_table_extractor()
        
        for doc in documents:
            content = doc.page_content
            source_file = doc.metadata.get('source', '')
            
            # Check if this is pre-extracted table data
            if doc.metadata.get('content_type') == 'table' and doc.metadata.get('table_data'):
                # This is already extracted table data - format as table chunks
                table_data = doc.metadata['table_data']
                
                try:
                    if isinstance(table_data, dict):
                        # Single table - create contextual chunks
                        max_chunk_size = config.chunk_token_num if config.chunk_token_num else 1000
                        surrounding_text = doc.page_content[:500] + doc.page_content[-500:]
                        
                        table_chunks = create_contextual_table_chunks(
                            table_data.get('data'), 
                            table_data, 
                            surrounding_text=surrounding_text,
                            max_chunk_size=max_chunk_size
                        )
                        
                        for chunk_data in table_chunks:
                            chunk_metadata = chunk_data['metadata']
                            chunk_metadata.update(doc.metadata)
                            
                            chunk = Document(
                                page_content=chunk_data['content'],
                                metadata=chunk_metadata
                            )
                            chunks.append(chunk)
                            
                    elif isinstance(table_data, list):
                        # Multiple tables
                        for table in table_data:
                            max_chunk_size = config.chunk_token_num if config.chunk_token_num else 1000
                            surrounding_text = doc.page_content[:500] + doc.page_content[-500:]
                            
                            table_chunks = create_contextual_table_chunks(
                                table.get('data'), 
                                table, 
                                surrounding_text=surrounding_text,
                                max_chunk_size=max_chunk_size
                            )
                            
                            for chunk_data in table_chunks:
                                chunk_metadata = chunk_data['metadata']
                                chunk_metadata.update(doc.metadata)
                                
                                chunk = Document(
                                    page_content=chunk_data['content'],
                                    metadata=chunk_metadata
                                )
                                chunks.append(chunk)
                                
                except Exception as e:
                    logger.warning(f"Failed to process pre-extracted table data: {e}")
                    
                continue  # Move to next document
            
            # Check if tables were already extracted at file level
            if doc.metadata.get('file_tables_extracted'):
                logger.info(f"Skipping table extraction for {doc.metadata.get('source_file_path', 'unknown')} - tables already extracted at file level")
                continue  # Skip table extraction for text docs when tables already extracted
            
            # Use pre-extracted tables if available to avoid duplicate extraction
            extracted_tables = pre_extracted_tables if pre_extracted_tables else []
            
            # Only extract tables if not already provided
            if not extracted_tables and source_file and Path(source_file).exists():
                try:
                    extracted_tables = table_extractor.extract_tables(str(source_file))
                    
                    if extracted_tables:
                        logger.info(f"Extracted {len(extracted_tables)} tables from {source_file} in _chunk_table")
                except Exception as e:
                    logger.warning(f"Table extraction failed for {source_file}: {e}")
            
            # Process extracted tables 
            if extracted_tables:
                for table_info in extracted_tables:
                    table_data = table_info['data']
                    
                    try:
                        if (PANDAS_AVAILABLE and 
                            isinstance(table_data, pd.DataFrame) and 
                            not table_data.empty):  # pandas DataFrame
                            # Use enhanced adaptive table chunking
                            max_chunk_size = config.chunk_token_num if config.chunk_token_num else 1000
                            
                            # Get surrounding text context
                            surrounding_text = doc.page_content[:500] + doc.page_content[-500:]
                            
                            # Create contextual table chunks (includes adaptive chunking internally)
                            table_chunks = create_contextual_table_chunks(
                                table_data, 
                                table_info, 
                                surrounding_text=surrounding_text,
                                max_chunk_size=max_chunk_size
                            )
                            
                            # Convert to Document objects
                            for chunk_data in table_chunks:
                                chunk_metadata = chunk_data['metadata']
                                chunk_metadata.update(doc.metadata)  # Include original document metadata
                                
                                chunk = Document(
                                    page_content=chunk_data['content'],
                                    metadata=chunk_metadata
                                )
                                chunks.append(chunk)
                    except Exception as e:
                        logger.warning(f"Failed to process table data: {e}")
                        # Skip this table and continue with others
                
                # If we successfully extracted tables, return them
                if chunks:
                    return chunks
            
            # Fallback: detect table structure in text content
            if detect_table_in_text(content):
                lines = content.split('\n')
                
                # Check if content looks like CSV/TSV
                if any('\t' in line or ',' in line for line in lines[:5]):
                    # Process as tabular data
                    delimiter = '\t' if '\t' in content else ','
                    
                    # Group rows into chunks
                    header = lines[0] if lines else ""
                    data_lines = lines[1:] if len(lines) > 1 else []
                    
                    chunk_size = max(10, config.chunk_token_num // 50)  # Rough estimate of rows per chunk
                    
                    for i in range(0, len(data_lines), chunk_size):
                        chunk_lines = [header] + data_lines[i:i + chunk_size]
                        chunk_content = '\n'.join(chunk_lines)
                        
                        chunk = Document(
                            page_content=chunk_content,
                            metadata={
                                'chunk_type': 'table',
                                'table_chunk': True, 
                                'chunk_rows': len(chunk_lines) - 1,
                                'extraction_method': 'text_detection'
                            }
                        )
                        chunks.append(chunk)
                else:
                    # Table-like structure but not CSV/TSV - preserve structure
                    table_chunks = self._chunk_table_like_text(content, config)
                    chunks.extend(table_chunks)
            else:
                # Not clearly tabular, use regular chunking
                chunks.extend(self._chunk_general([doc], config))
        
        return chunks

    def _chunk_general_table_aware(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """General chunking with table detection and special handling"""
        chunks = []
        
        for doc in documents:
            content = doc.page_content;
            
            # Check if content contains table-like structures
            if detect_table_in_text(content):
                logger.info(f"Detected table-like content in document, applying table-aware chunking")
                
                # Split content into table and non-table sections
                table_chunks = self._extract_table_sections_from_text(content, config)
                for chunk in table_chunks:
                    chunk.metadata.update({
                        'chunking_method': 'general_table_aware',
                        'contains_tables': True
                    })
                chunks.extend(table_chunks)
            else:
                # Regular general chunking for non-table content
                regular_chunks = self._chunk_general([doc], config)
                for chunk in regular_chunks:
                    chunk.metadata.update({
                        'chunking_method': 'general',
                        'contains_tables': False
                    })
                chunks.extend(regular_chunks);
        
        return chunks

    def _chunk_qa_table_aware(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Q&A chunking with table detection and special handling"""
        chunks = []
        
        for doc in documents:
            content = doc.page_content;
            
            # Check if content contains table-like structures
            if detect_table_in_text(content):
                logger.info(f"Detected table-like content in Q&A document, applying mixed chunking")
                
                # For Q&A documents with tables, we need to be more careful
                # Split by Q&A patterns first, then check each section for tables
                qa_sections = self._split_qa_sections(content);
                
                for i, section in enumerate(qa_sections):
                    if detect_table_in_text(section):
                        # This Q&A section contains a table
                        table_chunks = self._extract_table_sections_from_text(section, config);
                        for chunk in table_chunks:
                            chunk.metadata.update({
                                'chunking_method': 'qa_table_aware',
                                'qa_section': i,
                                'contains_tables': True
                            });
                        chunks.extend(table_chunks);
                    else:
                        # Regular Q&A processing
                        section_doc = Document(page_content=section, metadata=doc.metadata.copy());
                        qa_chunks = self._chunk_qa([section_doc], config);
                        for chunk in qa_chunks:
                            chunk.metadata.update({
                                'chunking_method': 'qa',
                                'qa_section': i,
                                'contains_tables': False
                            });
                        chunks.extend(qa_chunks);
            else:
                # Regular Q&A chunking
                qa_chunks = self._chunk_qa([doc], config);
                for chunk in qa_chunks:
                    chunk.metadata.update({
                        'chunking_method': 'qa',
                        'contains_tables': False
                    });
                chunks.extend(qa_chunks);
        
        return chunks;

    def _chunk_resume_table_aware(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Resume chunking with table detection and special handling"""
        chunks = []
        
        for doc in documents:
            content = doc.page_content;
            
            # Check if content contains table-like structures
            if detect_table_in_text(content):
                logger.info(f"Detected table-like content in resume, applying mixed chunking")
                
                # For resumes with tables (like skills matrices), preserve structure
                resume_sections = self._split_resume_sections(content);
                
                for section_name, section_content in resume_sections:
                    if detect_table_in_text(section_content):
                        # This resume section contains a table (e.g., skills matrix)
                        table_chunks = self._extract_table_sections_from_text(section_content, config);
                        for chunk in table_chunks:
                            chunk.metadata.update({
                                'chunking_method': 'resume_table_aware',
                                'resume_section': section_name,
                                'contains_tables': True
                            });
                        chunks.extend(table_chunks);
                    else:
                        # Regular resume section processing
                        section_doc = Document(page_content=section_content, metadata=doc.metadata.copy());
                        resume_chunks = self._chunk_resume([section_doc], config);
                        for chunk in resume_chunks:
                            chunk.metadata.update({
                                'chunking_method': 'resume',
                                'resume_section': section_name,
                                'contains_tables': False
                            });
                        chunks.extend(resume_chunks);
            else:
                # Regular resume chunking
                resume_chunks = self._chunk_resume([doc], config);
                for chunk in resume_chunks:
                    chunk.metadata.update({
                        'chunking_method': 'resume',
                        'contains_tables': False
                    });
                chunks.extend(resume_chunks);
        
        return chunks;

    def _extract_table_sections_from_text(self, content: str, config: ChunkingConfig) -> List[Document]:
        """Extract and chunk table sections from text content"""
        chunks = [];
        lines = content.split('\n');
        
        current_section = [];
        current_section_type = 'text';  # 'text' or 'table'
        
        i = 0;
        while i < len(lines):
            line = lines[i].strip();
            
            # Look ahead to detect table start
            table_start = self._detect_table_start(lines[i:i+5]);
            
            if table_start and current_section_type == 'text':
                # End current text section
                if current_section:
                    text_content = '\n'.join(current_section).strip();
                    if text_content:
                        chunks.append(Document(
                            page_content=text_content,
                            metadata={'chunk_type': 'text', 'extraction_method': 'text_section'}
                        ));
                
                # Start table section
                current_section = [line];
                current_section_type = 'table';
                
            elif not table_start and current_section_type == 'table':
                # End current table section
                if current_section:
                    table_content = '\n'.join(current_section).strip();
                    if table_content:
                        formatted_table = self._format_detected_table(table_content);
                        chunks.append(Document(
                            page_content=formatted_table,
                            metadata={
                                'chunk_type': 'table',
                                'extraction_method': 'text_detection',
                                'confidence': 0.8
                            }
                        ));
                
                # Start text section
                current_section = [line] if line else [];
                current_section_type = 'text';
                
            else:
                # Continue current section
                if line or current_section:  # Include empty lines within sections
                    current_section.append(line);
            
            i += 1;
        
        # Handle remaining section
        if current_section:
            section_content = '\n'.join(current_section).strip();
            if section_content:
                if current_section_type == 'table':
                    formatted_table = self._format_detected_table(section_content);
                    chunks.append(Document(
                        page_content=formatted_table,
                        metadata={
                            'chunk_type': 'table',
                            'extraction_method': 'text_detection',
                            'confidence': 0.8
                        }
                    ));
                else:
                    chunks.append(Document(
                        page_content=section_content,
                        metadata={'chunk_type': 'text', 'extraction_method': 'text_section'}
                    ));
        
        return chunks;

    def _detect_table_start(self, lines: List[str]) -> bool:
        """Detect if the next few lines start a table"""
        if len(lines) < 2:
            return False;
        
        # Look for table indicators in the next few lines
        table_lines = 0;
        for line in lines[:5]:
            if line.strip():
                if (line.count('|') >= 2 or 
                    line.count('\t') >= 1 or 
                    re.search(r'\s{3,}', line) or  # Multiple spaces
                    line.count(',') >= 2):
                    table_lines += 1;
        
        return table_lines >= 2;

    def _format_detected_table(self, table_content: str) -> str:
        """Format detected table content for better readability"""
        lines = table_content.split('\n');
        
        # Try to clean up and format the table
        formatted_lines = [];
        for line in lines:
            if line.strip():
                # Clean up excessive spaces and pipes
                cleaned_line = re.sub(r'\s*\|\s*', ' | ', line);
                cleaned_line = re.sub(r'\s+', ' ', cleaned_line).strip();
                formatted_lines.append(cleaned_line);
        
        return '\n'.join(formatted_lines);

    def _split_qa_sections(self, content: str) -> List[str]:
        """Split content into Q&A sections"""
        qa_patterns = [r'Q:', r'A:', r'Question:', r'Answer:', r'\d+\.', r'Q\d+', r'A\d+'];
        
        boundaries = [0];
        for pattern in qa_patterns:
            for match in re.finditer(pattern, content, re.IGNORECASE):
                boundaries.append(match.start());
        
        boundaries = sorted(set(boundaries));
        boundaries.append(len(content));
        
        sections = [];
        for i in range(len(boundaries) - 1):
            start, end = boundaries[i], boundaries[i + 1];
            section = content[start:end].strip();
            if section:
                sections.append(section);
        
        return sections;

    def _split_resume_sections(self, content: str) -> List[Tuple[str, str]]:
        """Split resume content into named sections"""
        section_patterns = {
            'experience': r'EXPERIENCE|WORK EXPERIENCE|EMPLOYMENT',
            'education': r'EDUCATION|ACADEMIC',
            'skills': r'SKILLS|TECHNICAL SKILLS|COMPETENCIES',
            'projects': r'PROJECTS|PROJECT EXPERIENCE',
            'certifications': r'CERTIFICATIONS|CERTIFICATES',
            'summary': r'SUMMARY|OBJECTIVE|PROFILE'
        };
        
        boundaries = [(0, 'header')];
        for section_name, pattern in section_patterns.items():
            for match in re.finditer(pattern, content, re.IGNORECASE):
                boundaries.append((match.start(), section_name));
        
        boundaries.sort();
        boundaries.append((len(content), 'end'));
        
        sections = [];
        for i in range(len(boundaries) - 1):
            start_pos, section_name = boundaries[i];
            end_pos, _ = boundaries[i + 1];
            
            section_content = content[start_pos:end_pos].strip();
            if section_content:
                sections.append((section_name, section_content));
        
        return sections;
    
    def _chunk_table_like_text(self, content: str, config: ChunkingConfig) -> List[Document]:
        """Chunk text that has table-like structure but isn't CSV/TSV"""
        chunks = []
        lines = content.split('\n')
        
        # Group lines that appear to be part of the same table
        current_table_lines = []
        
        for line in lines:
            if line.strip():
                # Check if this line looks like part of a table
                if ('|' in line or 
                    re.search(r'\s{3,}', line) or  # Multiple spaces (column separation)
                    re.search(r'^\s*[\w\s]+\s+[\w\s]+\s+[\w\s]+', line)):  # Multiple words separated by spaces
                    current_table_lines.append(line)
                else:
                    # End of current table - process it
                    if current_table_lines:
                        table_text = '\n'.join(current_table_lines)
                        chunk = Document(
                            page_content=table_text,
                            metadata={
                                'chunk_type': 'table',
                                'table_chunk': True,
                                'extraction_method': 'structure_detection'
                            }
                        )
                        chunks.append(chunk)
                        current_table_lines = []
                    
                    # Process non-table line normally
                    if line.strip():
                        chunk = Document(
                            page_content=line,
                            metadata={'chunk_type': 'text'}
                        )
                        chunks.append(chunk)
            else:
                # Empty line - might be table separator
                if current_table_lines:
                    current_table_lines.append(line)
        
        # Handle remaining table lines
        if current_table_lines:
            table_text = '\n'.join(current_table_lines)
            chunk = Document(
                page_content=table_text,
                metadata={
                    'chunk_type': 'table',
                    'table_chunk': True,
                    'extraction_method': 'structure_detection'
                }
            )
            chunks.append(chunk)
        
        return chunks
    
    def _chunk_presentation(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Presentation chunking - typically one chunk per slide, with table awareness"""
        chunks = []
        table_extractor = get_table_extractor()
        
        for doc in documents:
            content = doc.page_content
            source_file = doc.metadata.get('source', '')
            
            # First, try to extract tables from the presentation if we have a file path
            if source_file and Path(source_file).exists() and source_file.lower().endswith(('.pptx', '.ppt')):
                try:
                    extracted_tables = table_extractor.extract_tables(str(source_file))
                    
                    if extracted_tables:
                        # Group tables by slide
                        slide_tables = {}
                        for table_info in extracted_tables:
                            slide_num = table_info.get('slide_number', 1)
                            if slide_num not in slide_tables:
                                slide_tables[slide_num] = []
                            slide_tables[slide_num].append(table_info)
                        
                        # Process each slide's tables using proper row-based chunking
                        for slide_num, tables in slide_tables.items():
                            for table_info in tables:
                                table_data = table_info['data']
                                
                                try:
                                    if (PANDAS_AVAILABLE and 
                                        isinstance(table_data, pd.DataFrame) and 
                                        not table_data.empty):
                                        # Use enhanced adaptive table chunking
                                        max_chunk_size = config.chunk_token_num if config.chunk_token_num else 1000
                                        
                                        # Get surrounding text context from doc
                                        surrounding_text = doc.page_content[:500] + doc.page_content[-500:]
                                        
                                        # Create contextual table chunks (includes adaptive chunking internally)
                                        table_chunks = create_contextual_table_chunks(
                                            table_data, 
                                            table_info, 
                                            surrounding_text=surrounding_text,
                                            max_chunk_size=max_chunk_size
                                        )
                                        
                                        # Convert to Document objects
                                        for chunk_data in table_chunks:
                                            chunk_metadata = chunk_data['metadata']
                                            chunk_metadata.update({
                                                'slide_chunk': True,
                                                'content_type': 'table'
                                            })
                                            chunk_metadata.update(doc.metadata)  # Include original document metadata
                                            
                                            chunk = Document(
                                                page_content=chunk_data['content'],
                                                metadata=chunk_metadata
                                            )
                                            chunks.append(chunk)
                                except Exception as e:
                                    logger.warning(f"Failed to process table data: {e}")
                                    # Skip this table and continue with others
                        
                        # If we successfully extracted tables, continue with regular slide processing for non-table content
                        # This allows both table and text content to be processed from presentations
                        
                except Exception as e:
                    logger.warning(f"Table extraction failed for presentation {source_file}: {e}")
            
            # Process slide boundaries for remaining content
            slide_patterns = [r'Slide\s+\d+', r'Page\s+\d+', r'^\d+\s*$']
            boundaries = [0]
            
            for pattern in slide_patterns:
                for match in re.finditer(pattern, content, re.MULTILINE):
                    boundaries.append(match.start())
            
            # If no clear slide boundaries, split by double newlines
            if len(boundaries) <= 1:
                slide_breaks = [m.start() for m in re.finditer(r'\n\s*\n\s*\n', content)]
                boundaries.extend(slide_breaks)
            
            boundaries = sorted(set(boundaries))
            boundaries.append(len(content))
            
            # Create chunks for each slide
            for i in range(len(boundaries) - 1):
                start, end = boundaries[i], boundaries[i + 1]
                slide_content = content[start:end].strip()
                
                if slide_content:
                    # Check if this slide content contains table-like structures
                    if detect_table_in_text(slide_content):
                        # Process as table-aware content
                        table_chunks = self._chunk_table_like_text(slide_content, config)
                        for chunk in table_chunks:
                            chunk.metadata.update({
                                'slide_chunk': True, 
                                'slide_number': i,
                                'content_type': 'mixed'
                            })
                        chunks.extend(table_chunks)
                    else:
                        # Regular slide content
                        chunk = Document(
                            page_content=slide_content,
                            metadata={
                                'slide_chunk': True, 
                                'slide_number': i,
                                'content_type': 'text'
                            }
                        )
                        chunks.append(chunk)
        
        return chunks
    
    def _chunk_picture(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Image content chunking - extract text via OCR if available"""
        chunks = []
        
        for doc in documents:
            # For image documents, the content might be OCR text
            content = doc.page_content
            
            if content and len(content.strip()) > 0:
                # Text was extracted from image
                if len(content) > config.chunk_token_num and config.chunk_token_num > 0:
                    splitter = RecursiveCharacterTextSplitter(
                        chunk_size=config.chunk_token_num,
                        chunk_overlap=config.chunk_overlap
                    )
                    image_chunks = splitter.split_documents([doc])
                    
                    # Distribute OCR bounding boxes among chunks if available
                    if 'ocr_bounding_boxes' in doc.metadata:
                        bounding_boxes = doc.metadata['ocr_bounding_boxes']
                        self._distribute_bounding_boxes_to_chunks(image_chunks, bounding_boxes)
                    
                    chunks.extend(image_chunks)
                else:
                    # Single chunk - include all bounding boxes
                    if 'ocr_bounding_boxes' in doc.metadata:
                        doc.metadata['chunk_bounding_boxes'] = doc.metadata['ocr_bounding_boxes']
                    chunks.append(doc)
            else:
                # No text content - create placeholder chunk
                chunk = Document(
                    page_content="[Image content - no text extracted]",
                    metadata={'image_only': True}  # Only simple metadata
                )
                chunks.append(chunk)
        
        return chunks

    def _chunk_email(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Email content chunking with table awareness"""
        chunks = []
        
        for doc in documents:
            content = doc.page_content
            
            # Check if email contains table-like structures
            if detect_table_in_text(content):
                # Process tables in email
                table_chunks = self._extract_table_sections_from_text(content, config)
                for chunk in table_chunks:
                    chunk.metadata.update({
                        'email_chunk': True,
                        'content_type': chunk.metadata.get('chunk_type', 'mixed')
                    })
                chunks.extend(table_chunks)
            else:
                # Regular email processing - split by email sections
                email_sections = self._split_email_sections(content)
                
                for section_name, section_content in email_sections:
                    if len(section_content.strip()) > 50:  # Minimum section size
                        chunk = Document(
                            page_content=section_content,
                            metadata={
                                'email_chunk': True,
                                'email_section': section_name,
                                'content_type': 'text'
                            }
                        )
                        chunks.append(chunk)
        
        return chunks

    def _split_email_sections(self, content: str) -> List[Tuple[str, str]]:
        """Split email content into sections (header, body, signature)"""
        sections = []
        
        # Simple email parsing
        lines = content.split('\n')
        
        header_end = 0
        signature_start = len(lines)
        
        # Find end of headers (first empty line or start of body content)
        for i, line in enumerate(lines):
            if not line.strip() or not (':' in line and i < 10):
                header_end = i
                break
        
        # Find start of signature (common signature indicators)
        signature_patterns = [r'--', r'Best regards', r'Sincerely', r'Thanks', r'Regards']
        for i in range(len(lines) - 1, max(len(lines) - 10, header_end), -1):
            line = lines[i].strip()
            if any(re.search(pattern, line, re.IGNORECASE) for pattern in signature_patterns):
                signature_start = i
                break
        
        # Extract sections
        if header_end > 0:
            header_content = '\n'.join(lines[:header_end]).strip()
            if header_content:
                sections.append(('header', header_content))
        
        body_content = '\n'.join(lines[header_end:signature_start]).strip()
        if body_content:
            sections.append(('body', body_content))
        
        if signature_start < len(lines):
            signature_content = '\n'.join(lines[signature_start:]).strip()
            if signature_content:
                sections.append(('signature', signature_content))
        
        return sections

    def _distribute_bounding_boxes_to_chunks(self, chunks: List[Document], bounding_boxes: List[dict]):
        """Distribute OCR bounding boxes among text chunks based on text content matching"""
        if not bounding_boxes:
            return
            
        # Create a mapping of words to their bounding boxes
        word_to_bbox = {}
        for bbox in bounding_boxes:
            word_key = bbox['text'].lower().strip()
            if word_key:
                if word_key not in word_to_bbox:
                    word_to_bbox[word_key] = []
                word_to_bbox[word_key].append(bbox)
        
        # For each chunk, find matching bounding boxes
        for chunk in chunks:
            chunk_bboxes = []
            chunk_text = chunk.page_content.lower()
            chunk_words = chunk_text.split()
            
            # Try to match words from chunk to bounding boxes
            used_bboxes = set()
            for word in chunk_words:
                clean_word = ''.join(c for c in word if c.isalnum()).lower()
                if clean_word in word_to_bbox:
                    # Find an unused bounding box for this word
                    for bbox in word_to_bbox[clean_word]:
                        bbox_id = f"{bbox['left']},{bbox['top']},{bbox['width']},{bbox['height']}"
                        if bbox_id not in used_bboxes:
                            chunk_bboxes.append(bbox)
                            used_bboxes.add(bbox_id)
                            break
            
            # Add bounding boxes to chunk metadata
            chunk.metadata['chunk_bounding_boxes'] = chunk_bboxes
    
    def _chunk_email(self, documents: List[Document], config: ChunkingConfig) -> List[Document]:
        """Email chunking based on email structure"""
        chunks = []
        email_patterns = [
            r'^From:', r'^To:', r'^Subject:', r'^Date:',
            r'^Reply-To:', r'^CC:', r'^BCC:'
        ]
        
        for doc in documents:
            content = doc.page_content
            
            # Find email header boundaries
            boundaries = [0]
            for pattern in email_patterns:
                for match in re.finditer(pattern, content, re.MULTILINE | re.IGNORECASE):
                    boundaries.append(match.start())
            
            # Also split on email separators
            email_separators = [
                r'^>{1,}\s*',  # Quote markers
                r'^-{3,}',     # Dash separators
                r'^={3,}',     # Equal separators
                r'On .* wrote:'  # Email thread markers
            ]
            
            for pattern in email_separators:
                for match in re.finditer(pattern, content, re.MULTILINE):
                    boundaries.append(match.start())
            
            boundaries = sorted(set(boundaries))
            boundaries.append(len(content))
            
            # Create chunks for email sections
            for i in range(len(boundaries) - 1):
                start, end = boundaries[i], boundaries[i + 1]
                email_section = content[start:end].strip()
                
                if len(email_section) > 50:  # Minimum meaningful email section
                    chunk = Document(
                        page_content=email_section,
                        metadata={'email_section': True}  # Only simple metadata
                    )
                    chunks.append(chunk)
        
        return chunks
    

    
    # Fallback document loading methods (used when Docling is not available or fails)
    def _process_pdf_fallback(self, file_path: str) -> List[Document]:
        """Fallback PDF processing when Docling is not available"""
        documents = []
        table_extractor = get_table_extractor()
        
        try:
            if not ADVANCED_LOADERS_AVAILABLE:
                logger.error("No PDF processing capability available")
                return []
                
            # First load the PDF with standard text extraction
            loader = PyPDFLoader(file_path)
            text_documents = loader.load()
            
            # Extract tables from the PDF
            table_documents_created = False
            try:
                extracted_tables = table_extractor.extract_tables(file_path)
                
                if extracted_tables:
                    # Create separate documents for tables using proper row-based chunking
                    for table_info in extracted_tables:
                        table_data = table_info['data']
                        
                        try:
                            if (PANDAS_AVAILABLE and isinstance(table_data, pd.DataFrame) and not table_data.empty):
                                # Use enhanced adaptive table chunking with default chunk size
                                max_chunk_size = 1000  # Default chunk size for document loading
                                
                                # Get surrounding text context for better understanding
                                surrounding_text = ""
                                if text_documents:
                                    # Get context from nearest text document
                                    for text_doc in text_documents:
                                        if text_doc.page_content:
                                            surrounding_text = text_doc.page_content[:500] + text_doc.page_content[-500:]
                                            break
                                
                                # Create contextual table chunks (includes adaptive chunking internally)
                                table_chunks = create_contextual_table_chunks(
                                    table_data, 
                                    table_info, 
                                    surrounding_text=surrounding_text,
                                    max_chunk_size=max_chunk_size
                                )
                            
                                # Convert to Document objects
                                for chunk_data in table_chunks:
                                    chunk_metadata = {
                                        'source': Path(file_path).name,
                                        'content_type': 'table',
                                        'extracted_from': 'pdf_fallback',
                                        'extraction_method': 'fallback_table_extractor',
                                        'chunk_type': chunk_data.get('chunk_type', 'adaptive'),
                                        'chunk_index': chunk_data.get('chunk_index', 0),
                                        'original_table_rows': chunk_data.get('original_table_rows', len(table_data)),
                                        'original_table_cols': chunk_data.get('original_table_cols', len(table_data.columns)),
                                        'context_info': chunk_data.get('context_info', ''),
                                        'key_columns': chunk_data.get('key_columns', []),
                                        'row_range': chunk_data.get('row_range', ''),
                                        'confidence_score': chunk_data.get('confidence_score', 0.0)
                                    }
                                    
                                    table_doc = Document(
                                        page_content=chunk_data['content'],
                                        metadata=chunk_metadata
                                    )
                                    documents.append(table_doc)
                                    table_documents_created = True
                        except Exception as e:
                            logger.warning(f"Failed to process table data: {e}")
                            
            except Exception as e:
                logger.warning(f"Table extraction failed for PDF {file_path}: {e}")
            
            # Always add text documents in addition to tables
            # Process text documents and check for table-like content
            for doc in text_documents:
                content = doc.page_content
                
                # Mark that tables have been extracted at file level (if any were found)
                doc.metadata.update({
                    'file_tables_extracted': table_documents_created,
                    'source_file_path': file_path  # Keep track of source file path
                })
                
                # Check if this page contains table-like structures
                if detect_table_in_text(content):
                    # Mark as containing tables
                    doc.metadata.update({
                        'content_type': 'mixed',
                        'contains_tables': True,
                        'extraction_method': 'fallback_text'
                    })
                else:
                    doc.metadata.update({
                        'content_type': 'text',
                        'contains_tables': False,
                        'extraction_method': 'fallback_text'
                    })
                
                documents.append(doc)
            
            logger.info(f"PDF processing result: {len(documents)} documents total "
                       f"(tables: {sum(1 for d in documents if d.metadata.get('content_type') == 'table')}, "
                       f"text: {sum(1 for d in documents if d.metadata.get('content_type') in ['text', 'mixed'])})")
            
            return documents
            
        except Exception as e:
            logger.error(f"Failed to process PDF {file_path}: {e}")
            return []
    
    def _process_docx_fallback(self, file_path: str) -> List[Document]:
        """Fallback DOCX processing when Docling is not available"""
        if ADVANCED_LOADERS_AVAILABLE:
            try:
                loader = Docx2txtLoader(file_path)
                docs = loader.load()
                # Add extraction method metadata
                for doc in docs:
                    doc.metadata.update({
                        'extraction_method': 'fallback_docx',
                        'content_type': 'text'
                    })
                return docs
            except Exception as e:
                logger.error(f"Failed to process DOCX {file_path}: {e}")
        
        # Fallback to simple text reading
        return self._process_text_fallback(file_path)
    
    def _process_text_fallback(self, file_path: str) -> List[Document]:
        """Fallback text processing when Docling is not available"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            return [Document(
                page_content=content,
                metadata={
                    'source': Path(file_path).name,
                    'extraction_method': 'fallback_text',
                    'content_type': 'text'
                }
            )]
        except Exception as e:
            logger.error(f"Failed to process text file {file_path}: {e}")
            return []
    
    def _process_html_fallback(self, file_path: str) -> List[Document]:
        """Fallback HTML processing when Docling is not available"""
        return self._process_text_fallback(file_path)  # Simple text fallback
    
    def _process_pptx_fallback(self, file_path: str) -> List[Document]:
        """Fallback PPTX processing when Docling is not available"""
        if ADVANCED_LOADERS_AVAILABLE:
            try:
                loader = UnstructuredPowerPointLoader(file_path)
                docs = loader.load()
                # Add extraction method metadata
                for doc in docs:
                    doc.metadata.update({
                        'extraction_method': 'fallback_pptx',
                        'content_type': 'presentation'
                    })
                return docs
            except Exception as e:
                logger.error(f"Failed to process PPTX {file_path}: {e}")
        
        return []

    def _process_csv(self, file_path: str) -> List[Document]:
        """Process CSV files with proper table extraction"""
        documents = []
        table_extractor = get_table_extractor()
        
        try:
            # Use the table extractor for proper CSV processing
            extracted_tables = table_extractor.extract_tables(file_path)
            
            if extracted_tables:
                for table_info in extracted_tables:
                    table_data = table_info['data']
                    
                    try:
                        if (PANDAS_AVAILABLE and isinstance(table_data, pd.DataFrame) and not table_data.empty):
                            # Use enhanced adaptive table chunking for CSV tables
                            max_chunk_size = 1000  # Default for CSV files
                            
                            # CSV files typically don't have surrounding text context
                            surrounding_text = f"CSV file: {Path(file_path).name}"
                            
                            # Create contextual table chunks (includes adaptive chunking internally)
                            table_chunks = create_contextual_table_chunks(
                                table_data, 
                                table_info, 
                                surrounding_text=surrounding_text,
                                max_chunk_size=max_chunk_size
                            )
                            
                            # Convert to Document objects
                            for chunk_data in table_chunks:
                                chunk_metadata = chunk_data['metadata']
                                chunk_metadata.update({
                                    'source': Path(file_path).name,
                                    'type': 'csv',
                                    'content_type': 'table',
                                    'extracted_from': 'csv'
                                })
                                
                                documents.append(Document(
                                    page_content=chunk_data['content'],
                                    metadata=chunk_metadata
                                ))
                    except Exception as e:
                        logger.warning(f"Failed to process CSV table data: {e}")
                        
            else:
                # Fallback to LangChain CSV loader
                if ADVANCED_LOADERS_AVAILABLE:
                    try:
                        loader = CSVLoader(file_path)
                        documents = loader.load()
                        
                        # Mark as table content
                        for doc in documents:
                            doc.metadata.update({
                                'content_type': 'table',
                                'extraction_method': 'langchain_csv'
                            })
                    except Exception as e:
                        logger.error(f"Failed to process CSV {file_path}: {e}")
                        # Fallback to simple text reading
                        documents = self._process_text(file_path)
                        if documents:
                            documents[0].metadata.update({
                                'content_type': 'table',
                                'extraction_method': 'text_fallback'
                            })
                            
        except Exception as e:
            logger.error(f"CSV processing failed for {file_path}: {e}")
            # Final fallback to text processing
            documents = self._process_text(file_path)
        
        except Exception as e:
            logger.error(f"Failed to process CSV file {file_path}: {e}")
            # Fallback to simple text reading
            documents = self._process_text(file_path)
            if documents:
                documents[0].metadata.update({
                    'content_type': 'table',
                    'extraction_method': 'text_fallback'
                })
        
        return documents
    
    def _process_excel(self, file_path: str) -> List[Document]:
        """Process Excel files with proper table extraction"""
        documents = []
        table_extractor = get_table_extractor()
        
        try:
            # Use the table extractor for proper Excel processing
            extracted_tables = table_extractor.extract_tables(file_path)
            
            if extracted_tables:
                for table_info in extracted_tables:
                    table_data = table_info['data']
                    
                    try:
                        if (PANDAS_AVAILABLE and isinstance(table_data, pd.DataFrame) and not table_data.empty):
                            # Use enhanced adaptive table chunking for Excel tables
                            max_chunk_size = 1000  # Default for Excel files
                            
                            # Excel files may have sheet context
                            sheet_name = table_info.get('sheet_name', 'Unknown Sheet')
                            surrounding_text = f"Excel file: {Path(file_path).name}, Sheet: {sheet_name}"
                            
                            # Create contextual table chunks (includes adaptive chunking internally)
                            table_chunks = create_contextual_table_chunks(
                                table_data, 
                                table_info, 
                                surrounding_text=surrounding_text,
                                max_chunk_size=max_chunk_size
                            )
                            
                            # Convert to Document objects
                            for chunk_data in table_chunks:
                                chunk_metadata = chunk_data['metadata']
                                chunk_metadata.update({
                                    'source': Path(file_path).name,
                                    'type': 'excel',
                                    'content_type': 'table',
                                    'extracted_from': 'excel'
                                })
                                
                                documents.append(Document(
                                    page_content=chunk_data['content'],
                                    metadata=chunk_metadata
                                ))
                    except Exception as e:
                        logger.warning(f"Failed to process Excel table data: {e}")
                        
            else:
                # Fallback to pandas if table extractor doesn't work
                if PANDAS_AVAILABLE:
                    try:
                        df = pd.read_excel(file_path)
                        content = df.to_string()
                        
                        documents.append(Document(
                            page_content=content,
                            metadata={
                                'source': Path(file_path).name, 
                                'type': 'excel',
                                'content_type': 'table',
                                'extraction_method': 'pandas_fallback'
                            }
                        ))
                    except Exception as e:
                        logger.error(f"Failed to process Excel file with pandas {file_path}: {e}")
        
        except Exception as e:
            logger.error(f"Failed to process Excel file {file_path}: {e}")
        
        return documents
    
    def _process_presentation(self, file_path: str) -> List[Document]:
        """Process PowerPoint files"""
        try:
            # Try using python-pptx first
            try:
                from pptx import Presentation
                return self._process_pptx_with_python_pptx(file_path)
            except ImportError:
                logger.warning("python-pptx not available, trying unstructured")
            
            # Fallback to unstructured if available
            if ADVANCED_LOADERS_AVAILABLE:
                try:
                    loader = UnstructuredPowerPointLoader(file_path)
                    return loader.load()
                except Exception as e:
                    logger.error(f"Failed to process presentation with unstructured {file_path}: {e}")
            
        except Exception as e:
            logger.error(f"Failed to process presentation {file_path}: {e}")
        
        return []
    
    def _process_pptx_with_python_pptx(self, file_path: str) -> List[Document]:
        """Process PPTX files using python-pptx library with table extraction"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            documents = []
            table_extractor = get_table_extractor()
            
            # Extract tables using the table extractor
            extracted_tables = {}  # slide_number -> [table_info]
            try:
                tables = table_extractor.extract_tables(file_path)
                for table_info in tables:
                    slide_num = table_info.get('slide_number', 1)
                    if slide_num not in extracted_tables:
                        extracted_tables[slide_num] = []
                    extracted_tables[slide_num].append(table_info)
            except Exception as e:
                logger.warning(f"Table extraction failed for presentation {file_path}: {e}")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_tables = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Handle tables - check if shape actually contains a table
                    if hasattr(shape, 'table') and shape.has_table:
                        try:
                            table_text = []
                            for row in shape.table.rows:
                                row_text = []
                                for cell in row.cells:
                                    if cell.text.strip():
                                        row_text.append(cell.text.strip())
                                if row_text:
                                    table_text.append(' | '.join(row_text))
                            if table_text:
                                table_content = '\n'.join(table_text)
                                slide_tables.append(table_content)
                        except Exception as e:
                            logger.warning(f"Could not process table in slide {slide_num}: {e}")
                
                # Create document for slide text content
                if slide_text:
                    slide_content = '\n\n'.join(slide_text)
                    
                    # Check if slide contains table-like structures
                    contains_tables = bool(slide_tables) or detect_table_in_text(slide_content)
                    
                    documents.append(Document(
                        page_content=slide_content,
                        metadata={
                            'source': Path(file_path).name,
                            'slide_number': slide_num,
                            'type': 'presentation',
                            'content_type': 'mixed' if contains_tables else 'text',
                            'contains_tables': contains_tables,
                            'total_slides': len(prs.slides)
                        }
                    ))
                
                # Create separate documents for extracted tables on this slide using proper row-based chunking
                if slide_num in extracted_tables:
                    for table_info in extracted_tables[slide_num]:
                        table_data = table_info['data']
                        
                        try:
                            if (PANDAS_AVAILABLE and isinstance(table_data, pd.DataFrame) and not table_data.empty):
                                # Use enhanced adaptive table chunking for presentation tables
                                max_chunk_size = 1000  # Default for presentation files
                                
                                # Get slide context for better understanding
                                slide_content = '\n\n'.join(slide_text) if slide_text else ""
                                surrounding_text = f"Presentation slide {slide_num}: {slide_content[:300]}"
                                
                                # Create contextual table chunks (includes adaptive chunking internally)
                                table_chunks = create_contextual_table_chunks(
                                    table_data, 
                                    table_info, 
                                    surrounding_text=surrounding_text,
                                    max_chunk_size=max_chunk_size
                                )
                                
                                # Convert to Document objects
                                for chunk_data in table_chunks:
                                    chunk_metadata = chunk_data['metadata']
                                    chunk_metadata.update({
                                        'source': Path(file_path).name,
                                        'type': 'presentation',
                                        'content_type': 'table',
                                        'extracted_from': 'presentation',
                                        'total_slides': len(prs.slides)
                                    })
                                
                                documents.append(Document(
                                    page_content=chunk_data['content'],
                                    metadata=chunk_metadata
                                ))
                        except Exception as e:
                            logger.warning(f"Failed to process presentation table data: {e}")
                
                # Also add simple table documents for inline tables
                for table_content in slide_tables:
                    documents.append(Document(
                        page_content=table_content,
                        metadata={
                            'source': Path(file_path).name,
                            'slide_number': slide_num,
                            'type': 'presentation',
                            'content_type': 'table',
                            'extraction_method': 'inline_table',
                            'total_slides': len(prs.slides)
                        }
                    ))
            
            return documents
            
        except Exception as e:
            logger.error(f"Failed to process PPTX with python-pptx {file_path}: {e}")
            return []
    
    def _process_html(self, file_path: str) -> List[Document]:
        """Process HTML files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Simple HTML tag removal
            import re
            clean_content = re.sub(r'<[^>]+>', '', content)
            clean_content = re.sub(r'\s+', ' ', clean_content).strip()
            
            return [Document(
                page_content=clean_content,
                metadata={'source': Path(file_path).name, 'type': 'html'}  # Use filename only
            )]
        except Exception as e:
            logger.error(f"Failed to process HTML file {file_path}: {e}")
            return []
    
    def _process_json(self, file_path: str) -> List[Document]:
        """Process JSON files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Convert JSON to readable text
            content = json.dumps(data, indent=2)
            
            return [Document(
                page_content=content,
                metadata={'source': Path(file_path).name, 'type': 'json'}  # Use filename only
            )]
        except Exception as e:
            logger.error(f"Failed to process JSON file {file_path}: {e}")
            return []
    
    def _process_email(self, file_path: str) -> List[Document]:
        """Process email files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            return [Document(
                page_content=content,
                metadata={'source': Path(file_path).name, 'type': 'email'}  # Use filename only
            )]
        except Exception as e:
            logger.error(f"Failed to process email file {file_path}: {e}")
            return []
    
    def _process_image(self, file_path: str) -> List[Document]:
        """Process image files with OCR if available"""
        if OCR_AVAILABLE:
            try:
                logger.info(f"Processing image with OCR: {file_path}")
                image = Image.open(file_path)
                
                # Enhance image for better OCR results
                # Convert to RGB if needed
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                
                # Get OCR text with better configuration
                custom_config = r'--oem 3 --psm 6'
                text = pytesseract.image_to_string(image, config=custom_config)
                
                # Get OCR data with bounding boxes
                ocr_data = pytesseract.image_to_data(image, config=custom_config, output_type=pytesseract.Output.DICT)
                
                # Process bounding box data with lower confidence threshold
                bounding_boxes = []
                image_width, image_height = image.size
                
                for i in range(len(ocr_data['text'])):
                    confidence = int(ocr_data['conf'][i])
                    if confidence > 10:  # Lower confidence threshold to capture more text
                        word_text = ocr_data['text'][i].strip()
                        if word_text and len(word_text) > 0:  # Only include non-empty text
                            bbox = {
                                'text': word_text,
                                'left': int(ocr_data['left'][i]),
                                'top': int(ocr_data['top'][i]),
                                'width': int(ocr_data['width'][i]),
                                'height': int(ocr_data['height'][i]),
                                'confidence': confidence,
                                'block_num': int(ocr_data['block_num'][i]),
                                'par_num': int(ocr_data['par_num'][i]),
                                'line_num': int(ocr_data['line_num'][i]),
                                'word_num': int(ocr_data['word_num'][i])
                            }
                            bounding_boxes.append(bbox)
                
                # Log OCR output
                logger.info(f"OCR extracted text from {Path(file_path).name}:")
                logger.info(f"--- OCR OUTPUT START ---")
                logger.info(text.strip())
                logger.info(f"--- OCR OUTPUT END ---")
                logger.info(f"OCR extracted {len(text.strip())} characters and {len(bounding_boxes)} bounding boxes from {Path(file_path).name}")
                
                # If no text was extracted, try different OCR settings
                if len(text.strip()) < 10 and len(bounding_boxes) < 5:
                    logger.warning(f"Poor OCR results, trying alternative settings for {Path(file_path).name}")
                    
                    # Try different PSM mode for single block of text
                    alt_config = r'--oem 3 --psm 8'
                    alt_text = pytesseract.image_to_string(image, config=alt_config)
                    alt_ocr_data = pytesseract.image_to_data(image, config=alt_config, output_type=pytesseract.Output.DICT)
                    
                    if len(alt_text.strip()) > len(text.strip()):
                        logger.info(f"Alternative OCR settings produced better results for {Path(file_path).name}")
                        text = alt_text
                        ocr_data = alt_ocr_data
                        
                        # Reprocess bounding boxes with alternative data
                        bounding_boxes = []
                        for i in range(len(ocr_data['text'])):
                            confidence = int(ocr_data['conf'][i])
                            if confidence > 10:
                                word_text = ocr_data['text'][i].strip()
                                if word_text and len(word_text) > 0:
                                    bbox = {
                                        'text': word_text,
                                        'left': int(ocr_data['left'][i]),
                                        'top': int(ocr_data['top'][i]),
                                        'width': int(ocr_data['width'][i]),
                                        'height': int(ocr_data['height'][i]),
                                        'confidence': confidence,
                                        'block_num': int(ocr_data['block_num'][i]),
                                        'par_num': int(ocr_data['par_num'][i]),
                                        'line_num': int(ocr_data['line_num'][i]),
                                        'word_num': int(ocr_data['word_num'][i])
                                    }
                                    bounding_boxes.append(bbox)
                
                return [Document(
                    page_content=text,
                    metadata={
                        'source': Path(file_path).name, 
                        'type': 'image', 
                        'extracted_via': 'ocr',
                        'image_width': image_width,
                        'image_height': image_height,
                        'ocr_bounding_boxes': bounding_boxes
                    }
                )]
            except Exception as e:
                logger.error(f"Failed to process image {file_path}: {e}")
        else:
            logger.warning(f"OCR not available for image processing: {file_path}")
        
        # Return placeholder if OCR not available
        return [Document(
            page_content="[Image file - OCR not available]",
            metadata={'source': Path(file_path).name, 'type': 'image', 'extracted_via': 'none'}  # Use filename only
        )]

# Global instance
_document_processor = None

def get_document_processor() -> EnhancedDocumentProcessor:
    """Get global document processor instance"""
    global _document_processor
    if _document_processor is None:
        _document_processor = EnhancedDocumentProcessor()
    return _document_processor
