"""
Table Extraction and Processing for Document Chunking
Handles table detection and extraction from various document formats using Docling
"""
import os
import pandas as pd
import re
import logging

from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
from collections import Counter

logger = logging.getLogger(__name__)

# Try to import Docling for advanced document parsing
try:
    from docling.document_converter import DocumentConverter
    from docling.datamodel.pipeline_options import PipelineOptions
    DOCLING_AVAILABLE = True
    logger.info("Docling is available for advanced document parsing")
except ImportError as e:
    logger.warning(f"Docling not available. Falling back to basic extraction methods. Error: {e}")
    DOCLING_AVAILABLE = False

# Fallback libraries for non-PDF formats
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    logger.warning("python-pptx not available. PowerPoint table extraction disabled.")
    PPTX_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    logger.warning("python-docx not available. Word table extraction disabled.")
    DOCX_AVAILABLE = False

# Legacy libraries for fallback
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    logger.warning("pdfplumber not available. PDF table extraction will be limited.")
    PDFPLUMBER_AVAILABLE = False


class DoclingTableExtractor:
    """Extract tables using Docling's advanced document parsing capabilities"""
    
    def __init__(self):
        import time
        start_time = time.time()
        logger.info("🚀 DEBUG: Starting DoclingTableExtractor initialization...")
        
        self.converter = None
        if DOCLING_AVAILABLE:
            try:
                # Import GPU check utility
                from gpu_utils import configure_docling_device
                
                # Smart GPU/CPU configuration based on memory availability
                use_gpu, device_status = configure_docling_device()
                logger.info(f"🚀 DEBUG: Table extractor device configuration: {device_status}")
                
                # Initialize Docling converter with appropriate configuration
                logger.info("🚀 DEBUG: Creating DocumentConverter for table extraction...")
                converter_start = time.time()
                
                try:
                    # For Docling 2.43.0, use simple initialization for table extraction
                    # GPU/CPU control is handled via environment variables
                    if not use_gpu:
                        # Force CPU if GPU memory is insufficient
                        import os
                        os.environ['CUDA_VISIBLE_DEVICES'] = ''
                        logger.info("🚀 DEBUG: Table extractor forcing CPU mode via environment variables")
                    else:
                        logger.info("🚀 DEBUG: Table extractor using GPU acceleration (if available)")
                    
                    logger.info("🚀 DEBUG: Creating DocumentConverter for table extraction with simple initialization...")
                    self.converter = DocumentConverter()
                    device_type = "GPU" if use_gpu else "CPU"
                    logger.info(f"🚀 DEBUG: Table DocumentConverter created with default settings ({device_type})")
                    
                except ImportError:
                    # Fallback to simple initialization if any imports fail
                    logger.info("🚀 DEBUG: Advanced configuration not available, using simple initialization...")
                    if not use_gpu:
                        import os
                        os.environ['CUDA_VISIBLE_DEVICES'] = ''
                        logger.info("🚀 DEBUG: Table extractor forcing CPU mode via environment variables")
                    
                    self.converter = DocumentConverter()
                    device_type = "GPU" if use_gpu else "CPU"
                    logger.info(f"🚀 DEBUG: Table DocumentConverter created with default settings ({device_type})")
                
                converter_time = time.time() - converter_start
                logger.info(f"🚀 DEBUG: DocumentConverter for tables created in {converter_time:.2f} seconds")
                
                total_time = time.time() - start_time
                device_type = "GPU" if use_gpu else "CPU"
                logger.info(f"Docling DocumentConverter for tables initialized with {device_type} in {total_time:.2f} seconds")
                
            except Exception as e:
                logger.warning(f"Failed to initialize Docling table extractor: {e}")
                import traceback
                logger.warning(f"🚀 DEBUG: Full table extractor traceback: {traceback.format_exc()}")
                self.converter = None
    
    def extract_tables_from_document(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract tables using Docling's comprehensive document analysis
        """
        print(f"\n🚀 DEBUG: Starting Docling extraction for {file_path}")
        
        if not DOCLING_AVAILABLE or not self.converter:
            print("🚀 DEBUG: Docling not available, falling back to legacy methods")
            return []
        
        # First attempt with current configuration
        try:
            return self._extract_tables_internal(file_path)
        except Exception as e:
            # Check if it's a CUDA memory error
            from gpu_utils import is_cuda_memory_error, clear_gpu_memory
            
            if is_cuda_memory_error(e):
                logger.warning(f"🔄 CUDA memory error in table extraction, switching to CPU-only mode: {e}")
                
                # Clear GPU memory cache
                clear_gpu_memory()
                
                # Force CPU mode and reinitialize converter
                self._switch_to_cpu_mode()
                
                # Retry with CPU
                try:
                    return self._extract_tables_internal(file_path, retry=True)
                except Exception as retry_error:
                    logger.error(f"🚀 DEBUG: Table extraction failed even with CPU mode: {retry_error}")
                    return []  # Return empty list instead of raising
            else:
                logger.error(f"🚀 DEBUG: Docling extraction failed: {e}")
                return []
    
    def _switch_to_cpu_mode(self):
        """Switch table extractor to CPU-only mode"""
        try:
            logger.info("🔄 Switching table extractor to CPU-only mode...")
            
            # Force CPU configuration
            os.environ['CUDA_VISIBLE_DEVICES'] = ''
            
            # Reinitialize converter with CPU settings
            try:
                # Simple CPU-only reinitialization for table extractor
                logger.info("🔄 Reinitializing Table DocumentConverter with CPU-only mode...")
                self.converter = DocumentConverter()
                logger.info("🔄 Table DocumentConverter reinitialized with CPU-only settings")
                
            except Exception as reinit_error:
                logger.error(f"Failed to reinitialize table converter in CPU mode: {reinit_error}")
                raise
        except Exception as e:
            logger.error(f"Failed to switch table extractor to CPU mode: {e}")
            raise
    
    def _extract_tables_internal(self, file_path: str, retry: bool = False) -> List[Dict[str, Any]]:
        """Internal table extraction logic"""
        tables = []
        
        try:
            # Convert document using Docling (v2.43.0 simplified API)
            print("🚀 DEBUG: Converting document with Docling...")
            conversion_results = list(self.converter.convert(file_path))
            
            if not conversion_results:
                print("🚀 DEBUG: No conversion results")
                return []
            
            # Get the document from conversion results (v2.43.0 returns tuples)
            document = None
            for result_tuple in conversion_results:
                if isinstance(result_tuple, tuple) and len(result_tuple) == 2:
                    key, value = result_tuple
                    if key == 'document':
                        document = value
                        break
            
            if not document:
                print("🚀 DEBUG: Document not found in conversion results")
                return []
            
            print(f"🚀 DEBUG: Document converted successfully")
            print(f"  Document type: {type(document)}")
            print(f"  Tables: {len(document.tables) if hasattr(document, 'tables') else 'N/A'}")
            print(f"  Main text segments: {len(document.main_text) if hasattr(document, 'main_text') else 'N/A'}")
            
            # Extract tables directly from document.tables (Docling v2.43.0)
            if hasattr(document, 'tables') and document.tables:
                print(f"🚀 DEBUG: Found {len(document.tables)} tables in document")
                
                for table_idx, table in enumerate(document.tables):
                    print(f"🚀 DEBUG: Processing table {table_idx}: {type(table).__name__}")
                    
                    # Extract table data using Docling v2.43.0 API
                    table_data = self._extract_docling_table_data_v2(table, table_idx + 1)
                    
                    if table_data:
                        print(f"🚀 DEBUG: Successfully extracted table {table_idx + 1}")
                        tables.append(table_data)
                    else:
                        print(f"🚀 DEBUG: Failed to extract data from table {table_idx + 1}")
            
            # Also check main_text for table-like structures
            if hasattr(document, 'main_text') and document.main_text:
                for text_idx, text_obj in enumerate(document.main_text):
                    if hasattr(text_obj, 'text') and text_obj.text:
                        if self._detect_table_in_text(text_obj.text):
                            print(f"🚀 DEBUG: Found table-like text in segment {text_idx}")
                            
                            # Try to parse table from text
                            table_data = self._parse_table_from_text(text_obj.text, len(tables) + 1)
                            if table_data:
                                print(f"🚀 DEBUG: Successfully parsed table from text element {text_idx + 1}")
                                tables.append(table_data)
            
            print(f"🚀 DEBUG: Docling extraction completed. Found {len(tables)} tables")
            return tables
            
        except Exception as e:
            print(f"🚀 DEBUG: Docling extraction failed: {e}")
            logger.error(f"Docling table extraction failed: {e}")
            return []
    
    def _extract_docling_table_data(self, table_element, table_number: int) -> Optional[Dict[str, Any]]:
        """Extract table data from Docling table element"""
        print(f"\n🚀 DEBUG: _extract_docling_table_data for table {table_number}")
        
        try:
            # Get table structure from Docling element
            if hasattr(table_element, 'cells') and table_element.cells:
                # Docling provides structured cell data
                table_data = self._build_table_from_cells(table_element.cells)
            elif hasattr(table_element, 'text') and table_element.text:
                # Fallback to text parsing
                table_data = self._parse_table_text(table_element.text)
            else:
                print("🚀 DEBUG: No extractable table data found in element")
                return None
            
            if not table_data or len(table_data) < 2:
                print("🚀 DEBUG: Insufficient table data extracted")
                return None
            
            print(f"🚀 DEBUG: Raw table data extracted: {len(table_data)} rows")
            print(f"  First row: {table_data[0] if table_data else 'N/A'}")
            
            # Validate table structure
            if not self._validate_docling_table(table_data):
                print("🚀 DEBUG: Table validation failed - not a real table")
                return None
            
            # Process table data into DataFrame
            processed_table = self._process_docling_table(table_data, table_number)
            
            if processed_table:
                print(f"🚀 DEBUG: Successfully processed Docling table {table_number}")
                return processed_table
            else:
                print(f"🚀 DEBUG: Failed to process Docling table {table_number}")
                return None
                
        except Exception as e:
            print(f"🚀 DEBUG: Error extracting Docling table data: {e}")
            logger.warning(f"Failed to extract Docling table data: {e}")
            return None
    
    def _build_table_from_cells(self, cells) -> List[List[str]]:
        """Build table structure from Docling cell objects"""
        print("🚀 DEBUG: Building table from Docling cells")
        
        try:
            # Group cells by row and column
            cell_grid = {}
            max_row, max_col = 0, 0
            
            for cell in cells:
                if hasattr(cell, 'bbox') and hasattr(cell, 'text'):
                    # Use bounding box to determine position
                    row = getattr(cell, 'row', 0)
                    col = getattr(cell, 'col', 0)
                    
                    # Fallback: estimate position from bbox if row/col not available
                    if not hasattr(cell, 'row') and hasattr(cell, 'bbox'):
                        # Simple heuristic: sort by y-coordinate for rows, x-coordinate for columns
                        row = int(cell.bbox.top / 20)  # Rough row estimation
                        col = int(cell.bbox.left / 50)  # Rough column estimation
                    
                    cell_grid[(row, col)] = str(cell.text).strip()
                    max_row = max(max_row, row)
                    max_col = max(max_col, col)
            
            # Build table data from grid
            table_data = []
            for row in range(max_row + 1):
                row_data = []
                for col in range(max_col + 1):
                    cell_content = cell_grid.get((row, col), "")
                    row_data.append(cell_content)
                table_data.append(row_data)
            
            print(f"🚀 DEBUG: Built table with {len(table_data)} rows and {max_col + 1} columns")
            return table_data
            
        except Exception as e:
            print(f"🚀 DEBUG: Error building table from cells: {e}")
            return []
    
    def _parse_table_text(self, text: str) -> List[List[str]]:
        """Parse table from text representation"""
        print("🚀 DEBUG: Parsing table from text")
        
        lines = text.strip().split('\n')
        table_data = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Try different delimiters
            if '|' in line:
                row = [cell.strip() for cell in line.split('|') if cell.strip()]
            elif '\t' in line:
                row = [cell.strip() for cell in line.split('\t')]
            elif ',' in line and line.count(',') >= 2:
                row = [cell.strip() for cell in line.split(',')]
            else:
                # Split by multiple spaces
                row = [cell.strip() for cell in re.split(r'\s{2,}', line) if cell.strip()]
            
            if row and len(row) >= 2:
                table_data.append(row)
        
        print(f"🚀 DEBUG: Parsed {len(table_data)} rows from text")
        return table_data
    
    def _validate_docling_table(self, table_data: List[List]) -> bool:
        """Validate that extracted data represents a real table"""
        print(f"🚀 DEBUG: Validating Docling table with {len(table_data)} rows")
        
        if not table_data or len(table_data) < 2:
            print("🚀 DEBUG: Too few rows for a table")
            return False
        
        # Check column consistency
        row_lengths = [len(row) for row in table_data if row]
        if not row_lengths or max(row_lengths) < 2:
            print("🚀 DEBUG: Need at least 2 columns")
            return False
        
        # Check for reasonable column consistency
        avg_cols = sum(row_lengths) / len(row_lengths)
        col_variation = max(row_lengths) - min(row_lengths)
        
        if col_variation > avg_cols * 0.5:  # Too much variation
            print(f"🚀 DEBUG: Too much column variation: {col_variation} vs avg {avg_cols}")
            return False
        
        # Check for table-like content patterns
        structured_rows = 0
        total_cells = 0
        long_cells = 0
        
        for row in table_data[:5]:  # Check first 5 rows
            if len(row) >= 2:
                structured_rows += 1
                
                for cell in row:
                    total_cells += 1
                    if len(str(cell).strip()) > 100:  # Very long cell
                        long_cells += 1
        
        # Too many long cells suggests paragraph text
        if total_cells > 0 and (long_cells / total_cells) > 0.3:
            print(f"🚀 DEBUG: Too many long cells ({long_cells}/{total_cells}) - likely text")
            return False
        
        if structured_rows < 2:
            print(f"🚀 DEBUG: Not enough structured rows ({structured_rows})")
            return False
        
        print("🚀 DEBUG: Docling table validation passed")
        return True
    
    def _process_docling_table(self, table_data: List[List], table_number: int) -> Optional[Dict[str, Any]]:
        """Process validated table data into final format"""
        print(f"🚀 DEBUG: Processing Docling table {table_number}")
        
        try:
            # Clean table data
            cleaned_data = []
            for row in table_data:
                cleaned_row = [str(cell).strip() for cell in row]
                if any(cell for cell in cleaned_row):  # Keep rows with at least one non-empty cell
                    cleaned_data.append(cleaned_row)
            
            if len(cleaned_data) < 2:
                print("🚀 DEBUG: Not enough data after cleaning")
                return None
            
            # Extract headers and data
            headers = cleaned_data[0]
            data_rows = cleaned_data[1:]
            
            # Ensure consistent row lengths
            max_cols = len(headers)
            for row in data_rows:
                while len(row) < max_cols:
                    row.append("")
                if len(row) > max_cols:
                    row = row[:max_cols]
            
            # Create DataFrame
            df = pd.DataFrame(data_rows, columns=headers)
            
            # Clean DataFrame
            df = self._clean_docling_dataframe(df)
            
            if df.empty:
                print("🚀 DEBUG: DataFrame empty after cleaning")
                return None
            
            # Final validation
            if not self._validate_dataframe_structure(df):
                print("🚀 DEBUG: DataFrame structure validation failed")
                return None
            
            result = {
                'data': df,
                'page_number': None,  # Docling provides page info separately
                'slide_number': None,
                'table_number': table_number,
                'bbox': None,
                'extraction_method': 'docling_advanced',
                'confidence': 0.95  # Higher confidence for Docling extraction
            }
            
            print(f"🚀 DEBUG: Successfully processed Docling table {table_number}")
            return result
            
        except Exception as e:
            print(f"🚀 DEBUG: Error processing Docling table: {e}")
            return None
    
    def _clean_docling_dataframe(self, df: pd.DataFrame) -> pd.DataFrame:
        """Clean DataFrame extracted from Docling"""
        if df.empty:
            return df
        
        print(f"🚀 DEBUG: Cleaning Docling DataFrame with shape {df.shape}")
        
        # Clean column names
        new_columns = []
        for i, col in enumerate(df.columns):
            col_str = str(col).strip()
            if not col_str or col_str.lower() in ['none', 'nan', '']:
                new_columns.append(f'Column_{i+1}')
            else:
                # Clean but preserve meaningful names
                cleaned = re.sub(r'\s+', ' ', col_str).strip()
                new_columns.append(cleaned)
        
        df.columns = new_columns
        
        # Clean cell values
        for col in df.columns:
            df[col] = df[col].astype(str).str.strip()
            df[col] = df[col].replace(['None', 'nan', '', 'null'], pd.NA)
        
        # Remove empty rows and columns
        df = df.dropna(how='all').reset_index(drop=True)
        df = df.dropna(axis=1, how='all')
        
        print(f"🚀 DEBUG: Cleaned DataFrame shape: {df.shape}")
        return df
    
    def _validate_dataframe_structure(self, df: pd.DataFrame) -> bool:
        """Final validation of DataFrame structure"""
        if df.empty or len(df) < 1:
            return False
        
        # Check for meaningful data
        non_empty_cols = sum(1 for col in df.columns if not df[col].isna().all())
        if non_empty_cols < 2:
            return False
        
        # Check cell length patterns
        total_cells = 0
        long_cells = 0
        
        for col in df.columns:
            for val in df[col].dropna():
                total_cells += 1
                if len(str(val)) > 80:
                    long_cells += 1
        
        # Reject if too many cells are very long
        if total_cells > 0 and (long_cells / total_cells) > 0.4:
            return False
        
        return True
    
    def _extract_docling_table_data_v2(self, table, table_number: int) -> Optional[Dict[str, Any]]:
        """Extract table data from Docling v2.43.0 Table object"""
        print(f"\n🚀 DEBUG: _extract_docling_table_data_v2 for table {table_number}")
        
        try:
            # Use the export_to_dataframe method from Docling v2.43.0
            if hasattr(table, 'export_to_dataframe'):
                print("🚀 DEBUG: Using export_to_dataframe method")
                df = table.export_to_dataframe()
                
                if df is not None and not df.empty:
                    print(f"🚀 DEBUG: Successfully got DataFrame with shape {df.shape}")
                    
                    # Clean the DataFrame
                    df = self._clean_docling_dataframe(df)
                    
                    if df.empty or not self._validate_dataframe_structure(df):
                        print("🚀 DEBUG: DataFrame validation failed")
                        return None
                    
                    result = {
                        'data': df,
                        'page_number': None,  # Docling v2.43.0 provides page info separately
                        'slide_number': None,
                        'table_number': table_number,
                        'bbox': None,
                        'extraction_method': 'docling_v2.43.0',
                        'confidence': 0.95,  # High confidence for Docling extraction
                        'docling_version': '2.43.0'
                    }
                    
                    print(f"🚀 DEBUG: Successfully extracted table {table_number} with Docling v2.43.0")
                    return result
                else:
                    print("🚀 DEBUG: export_to_dataframe returned empty DataFrame")
            
            # Fallback: use text representation if available
            elif hasattr(table, 'text') and table.text:
                print("🚀 DEBUG: Fallback to text parsing")
                table_data = self._parse_table_text(table.text)
                
                if table_data and self._validate_docling_table(table_data):
                    processed_table = self._process_docling_table(table_data, table_number)
                    if processed_table:
                        processed_table['extraction_method'] = 'docling_v2.43.0_text'
                        processed_table['docling_version'] = '2.43.0'
                        return processed_table
            
            print("🚀 DEBUG: No extractable table data found")
            return None
            
        except Exception as e:
            print(f"🚀 DEBUG: Error extracting Docling v2.43.0 table data: {e}")
            logger.warning(f"Failed to extract Docling v2.43.0 table data: {e}")
            return None
    
    def _detect_table_in_text(self, text: str) -> bool:
        """Detect if text contains table-like structures"""
        if not text or len(text.strip()) < 50:
            return False
        
        lines = text.strip().split('\n')
        if len(lines) < 3:
            return False
        
        # Check for table indicators
        pipe_lines = sum(1 for line in lines if '|' in line)
        tab_lines = sum(1 for line in lines if '\t' in line)
        comma_lines = sum(1 for line in lines if line.count(',') >= 2)
        
        # If significant portion has table delimiters
        total_lines = len(lines)
        if (pipe_lines / total_lines) > 0.3 or (tab_lines / total_lines) > 0.3 or (comma_lines / total_lines) > 0.3:
            return True
        
        return False
    
    def _parse_table_from_text(self, text: str, element_number: int) -> Optional[Dict[str, Any]]:
        """Parse table from text element"""
        table_data = self._parse_table_text(text)
        
        if not table_data or not self._validate_docling_table(table_data):
            return None
        
        return self._process_docling_table(table_data, element_number)
class TableExtractor:
    """Enhanced table extractor using Docling as primary method with fallbacks"""
    
    def __init__(self):
        # Initialize Docling extractor
        self.docling_extractor = DoclingTableExtractor() if DOCLING_AVAILABLE else None
        
        # Fallback format support
        self.supported_formats = {
            '.pdf': self.extract_pdf_tables,
            '.pptx': self.extract_pptx_tables,
            '.ppt': self.extract_pptx_tables,
            '.docx': self.extract_docx_tables,
            '.doc': self.extract_docx_tables,
            '.csv': self.extract_csv_tables,
            '.xlsx': self.extract_excel_tables,
            '.xls': self.extract_excel_tables
        }
    
    def extract_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract tables from a document file using Docling as primary method
        
        Returns:
            List of table dictionaries with structure:
            {
                'data': pandas.DataFrame,
                'page_number': int or None,
                'slide_number': int or None,
                'table_number': int,
                'bbox': tuple or None (x1, y1, x2, y2),
                'extraction_method': str,
                'confidence': float
            }
        """
        file_ext = Path(file_path).suffix.lower()
        
        print(f"\n🎯 DEBUG: Starting table extraction for {file_path}")
        print(f"  File extension: {file_ext}")
        print(f"  Docling available: {DOCLING_AVAILABLE}")
        
        # Try Docling first for PDF files (primary method)
        if file_ext == '.pdf' and self.docling_extractor:
            print("🎯 DEBUG: Attempting Docling extraction")
            tables = self.docling_extractor.extract_tables_from_document(file_path)
            
            if tables:
                print(f"🎯 DEBUG: Docling extraction successful - found {len(tables)} tables")
                return tables
            else:
                print("🎯 DEBUG: Docling extraction found no tables, trying fallback methods")
        
        # Fallback to format-specific extractors
        if file_ext not in self.supported_formats:
            logger.warning(f"File format {file_ext} not supported for table extraction")
            return []
        
        try:
            print(f"🎯 DEBUG: Using fallback extraction method for {file_ext}")
            tables = self.supported_formats[file_ext](file_path)
            print(f"🎯 DEBUG: Fallback extraction found {len(tables)} tables")
            return tables
        except Exception as e:
            logger.error(f"Error extracting tables from {file_path}: {e}")
            return []
    
    def extract_pdf_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract tables from PDF files - enhanced with Docling integration"""
        tables = []
        
        # Method 1: Docling (primary method - most accurate)
        if self.docling_extractor:
            print("🎯 DEBUG: PDF extraction using Docling (primary)")
            tables.extend(self.docling_extractor.extract_tables_from_document(file_path))
            
            if tables:
                print(f"🎯 DEBUG: Docling found {len(tables)} tables, returning")
                return tables
        
        # Method 2: pdfplumber fallback (for simple tables)
        if PDFPLUMBER_AVAILABLE:
            print("🎯 DEBUG: PDF extraction using pdfplumber (fallback)")
            fallback_tables = self._extract_pdf_tables_pdfplumber_fallback(file_path)
            tables.extend(fallback_tables)
            
            if fallback_tables:
                print(f"🎯 DEBUG: pdfplumber fallback found {len(fallback_tables)} tables")
        
        return tables
    
    def _extract_pdf_tables_pdfplumber_fallback(self, file_path: str) -> List[Dict[str, Any]]:
        """Simplified pdfplumber extraction as fallback when Docling fails"""
        print(f"\n📋 DEBUG: pdfplumber fallback extraction for {file_path}")
        tables = []
        
        try:
            with pdfplumber.open(file_path) as pdf:
                print(f"📋 DEBUG: PDF has {len(pdf.pages)} pages")
                
                for page_num, page in enumerate(pdf.pages, 1):
                    print(f"📋 DEBUG: Processing page {page_num}")
                    
                    try:
                        # Use simple extraction with basic settings
                        page_tables = page.extract_tables()
                        print(f"📋 DEBUG: Found {len(page_tables)} potential tables")
                        
                        # Quick filter for obvious false positives
                        if len(page_tables) > 5:  # Too many suggests false detection
                            print("📋 DEBUG: Too many potential tables - likely false detection")
                            continue
                        
                        for table_num, table_data in enumerate(page_tables, 1):
                            if table_data and len(table_data) > 1 and len(table_data) < 30:  # Reasonable size
                                # Basic validation
                                if self._basic_table_validation(table_data):
                                    processed_table = self._process_fallback_table(
                                        table_data, page_num, table_num
                                    )
                                    
                                    if processed_table:
                                        tables.append(processed_table)
                                        print(f"📋 DEBUG: Successfully processed fallback table {table_num}")
                                        
                    except Exception as e:
                        print(f"📋 DEBUG: Page {page_num} extraction failed: {e}")
                        continue
                        
        except Exception as e:
            print(f"📋 DEBUG: pdfplumber fallback extraction failed: {e}")
            logger.error(f"pdfplumber fallback extraction failed: {e}")
        
        return tables
    
    def _basic_table_validation(self, table_data: List[List]) -> bool:
        """Basic validation for fallback extraction"""
        if not table_data or len(table_data) < 2:
            return False
        
        # Check column consistency and reasonable content
        row_lengths = [len(row) for row in table_data if row]
        if not row_lengths or max(row_lengths) < 2:
            return False
        
        # Check for extremely long cells (paragraph text)
        for row in table_data[:3]:
            for cell in row:
                if cell and len(str(cell)) > 150:  # Very long cell
                    return False
        
        return True
    
    def _process_fallback_table(self, table_data: List[List], page_num: int, table_num: int) -> Optional[Dict[str, Any]]:
        """Process table data for fallback extraction"""
        try:
            # Clean data
            cleaned_data = []
            for row in table_data:
                cleaned_row = [str(cell).strip() if cell else "" for cell in row]
                if any(cell for cell in cleaned_row):
                    cleaned_data.append(cleaned_row)
            
            if len(cleaned_data) < 2:
                return None
            
            # Create DataFrame
            headers = cleaned_data[0]
            data_rows = cleaned_data[1:]
            
            # Ensure consistent lengths
            max_cols = len(headers)
            for row in data_rows:
                while len(row) < max_cols:
                    row.append("")
                if len(row) > max_cols:
                    row = row[:max_cols]
            
            df = pd.DataFrame(data_rows, columns=headers)
            
            # Basic cleaning
            for col in df.columns:
                df[col] = df[col].astype(str).str.strip()
                df[col] = df[col].replace(['None', 'nan', ''], pd.NA)
            
            df = df.dropna(how='all').reset_index(drop=True)
            
            if df.empty:
                return None
            
            return {
                'data': df,
                'page_number': page_num,
                'slide_number': None,
                'table_number': table_num,
                'bbox': None,
                'extraction_method': 'pdfplumber_fallback',
                'confidence': 0.6  # Lower confidence for fallback
            }
            
        except Exception as e:
            print(f"📋 DEBUG: Error processing fallback table: {e}")
            return None
    
    def extract_pptx_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract tables from PowerPoint files"""
        tables = []
        
        if not PPTX_AVAILABLE:
            return tables
        
        print(f"📊 DEBUG: Extracting tables from PowerPoint: {file_path}")
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_tables = []
                
                for shape in slide.shapes:
                    if hasattr(shape, 'table') and shape.has_table:
                        table_data = []
                        
                        for row in shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text.strip())
                            table_data.append(row_data)
                        
                        if table_data and len(table_data) > 1:
                            try:
                                df = pd.DataFrame(table_data[1:], columns=table_data[0])
                                df = self._clean_simple_dataframe(df)
                                
                                if not df.empty:
                                    tables.append({
                                        'data': df,
                                        'page_number': None,
                                        'slide_number': slide_num,
                                        'table_number': len(slide_tables) + 1,
                                        'bbox': None,
                                        'extraction_method': 'python_pptx',
                                        'confidence': 0.9
                                    })
                                    slide_tables.append(df)
                            except Exception as e:
                                logger.warning(f"Failed to process table in slide {slide_num}: {e}")
        
        except Exception as e:
            logger.error(f"PowerPoint table extraction failed: {e}")
        
        print(f"📊 DEBUG: PowerPoint extraction found {len(tables)} tables")
        return tables
    
    def extract_docx_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract tables from Word documents"""
        tables = []
        
        if not DOCX_AVAILABLE:
            return tables
        
        print(f"📄 DEBUG: Extracting tables from Word document: {file_path}")
        
        try:
            doc = DocxDocument(file_path)
            
            for table_num, table in enumerate(doc.tables, 1):
                table_data = []
                
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                
                if table_data and len(table_data) > 1:
                    try:
                        df = pd.DataFrame(table_data[1:], columns=table_data[0])
                        df = self._clean_simple_dataframe(df)
                        
                        if not df.empty:
                            tables.append({
                                'data': df,
                                'page_number': None,
                                'slide_number': None,
                                'table_number': table_num,
                                'bbox': None,
                                'extraction_method': 'python_docx',
                                'confidence': 0.9
                            })
                    except Exception as e:
                        logger.warning(f"Failed to process table {table_num}: {e}")
        
        except Exception as e:
            logger.error(f"Word document table extraction failed: {e}")
        
        print(f"📄 DEBUG: Word document extraction found {len(tables)} tables")
        return tables
    
    def extract_csv_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract tables from CSV files"""
        tables = []
        
        try:
            df = pd.read_csv(file_path)
            
            if not df.empty:
                tables.append({
                    'data': df,
                    'page_number': None,
                    'slide_number': None,
                    'table_number': 1,
                    'bbox': None,
                    'extraction_method': 'pandas_csv',
                    'confidence': 1.0
                })
        
        except Exception as e:
            logger.error(f"CSV table extraction failed: {e}")
        
        return tables
    
    def extract_excel_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract tables from Excel files"""
        tables = []
        
        try:
            excel_file = pd.ExcelFile(file_path)
            
            for sheet_num, sheet_name in enumerate(excel_file.sheet_names, 1):
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                if not df.empty:
                    tables.append({
                        'data': df,
                        'page_number': None,
                        'slide_number': sheet_num,
                        'table_number': 1,
                        'bbox': None,
                        'extraction_method': 'pandas_excel',
                        'confidence': 1.0,
                        'sheet_name': sheet_name
                    })
        
        except Exception as e:
            logger.error(f"Excel table extraction failed: {e}")
        
        return tables
    
    def _clean_simple_dataframe(self, df: pd.DataFrame) -> pd.DataFrame:
        """Simple DataFrame cleaning for non-PDF formats"""
        if df.empty:
            return df
        
        # Clean column names
        new_columns = []
        for i, col in enumerate(df.columns):
            col_str = str(col).strip()
            if not col_str or col_str.lower() in ['none', 'nan', '']:
                new_columns.append(f'Column_{i+1}')
            else:
                new_columns.append(col_str)
        
        df.columns = new_columns
        
        # Clean cell values
        for col in df.columns:
            df[col] = df[col].astype(str).str.strip()
            df[col] = df[col].replace(['None', 'nan', ''], pd.NA)
        
        # Remove empty rows and columns
        df = df.dropna(how='all').reset_index(drop=True)
        df = df.dropna(axis=1, how='all')
        
        return df


def detect_table_in_text(text: str) -> bool:
    """
    Detect if text contains table-like structures
    Enhanced for Docling integration
    """
    if not text or len(text.strip()) < 30:
        return False
        
    lines = text.strip().split('\n')
    
    if len(lines) < 2:
        return False
    
    # Check for common table indicators with improved patterns
    table_indicators = [
        # Multiple consecutive lines with delimiters
        lambda: sum(1 for line in lines if '|' in line or '\t' in line) >= 2,
        
        # Lines with consistent column structure
        lambda: _has_consistent_columns(lines),
        
        # Table headers pattern (improved)
        lambda: any(re.search(r'^[\w\s]+\s*[\|\t]\s*[\w\s]+', line) for line in lines[:3]),
        
        # CSV-like structure with better detection
        lambda: sum(1 for line in lines if line.count(',') >= 2 and len(line.split(',')) >= 3) >= 2,
        
        # Tab-separated values with minimum columns
        lambda: sum(1 for line in lines if line.count('\t') >= 1 and len(line.split('\t')) >= 2) >= 2,
        
        # Structured data patterns (numbers, percentages, etc.)
        lambda: _has_structured_data_patterns(lines),
    ]
    
    return any(indicator() for indicator in table_indicators)


def _has_consistent_columns(lines: List[str]) -> bool:
    """Check if lines have consistent column structure (enhanced)"""
    if len(lines) < 3:
        return False
    
    # Check for pipe-separated columns
    pipe_counts = [line.count('|') for line in lines[:6] if line.strip()]
    if len(pipe_counts) >= 3:
        unique_counts = set(pipe_counts)
        if len(unique_counts) <= 2 and max(pipe_counts) >= 1:
            return True
    
    # Check for tab-separated columns
    tab_counts = [line.count('\t') for line in lines[:6] if line.strip()]
    if len(tab_counts) >= 3:
        unique_counts = set(tab_counts)
        if len(unique_counts) <= 2 and max(tab_counts) >= 1:
            return True
    
    # Check for space-separated columns (multiple spaces as delimiter)
    space_separated = []
    for line in lines[:6]:
        if line.strip():
            # Split by multiple spaces
            parts = [part.strip() for part in re.split(r'\s{2,}', line.strip()) if part.strip()]
            if len(parts) >= 2:
                space_separated.append(len(parts))
    
    if len(space_separated) >= 3:
        unique_counts = set(space_separated)
        if len(unique_counts) <= 2 and max(space_separated) >= 2:
            return True
    
    return False


def _has_structured_data_patterns(lines: List[str]) -> bool:
    """Check for structured data patterns that indicate table content"""
    if len(lines) < 3:
        return False
    
    pattern_lines = 0
    total_lines = min(len(lines), 8)  # Check first 8 lines
    
    for line in lines[:total_lines]:
        line = line.strip()
        if not line:
            continue
        
        # Look for patterns that suggest tabular data
        has_numbers = bool(re.search(r'\d+\.?\d*%?', line))  # Numbers or percentages
        has_delimited_structure = '|' in line or '\t' in line or line.count(',') >= 2
        has_short_entries = len(line.split()) >= 2 and all(len(word) <= 30 for word in line.split())
        
        if (has_numbers and has_delimited_structure) or (has_delimited_structure and has_short_entries):
            pattern_lines += 1
    
    # If at least 40% of lines show structured patterns
    return pattern_lines >= (total_lines * 0.4)


def create_adaptive_table_chunks(table_data: pd.DataFrame, table_info: Dict[str, Any], 
                                max_chunk_size: int = 7800, preserve_context: bool = True) -> List[Dict[str, Any]]:
    """
    Create adaptive table chunks based on table size and content
    """
    print(f"\n🔍 DEBUG: create_adaptive_table_chunks called")
    print(f"  table_data shape: {table_data.shape if not table_data.empty else 'EMPTY'}")
    print(f"  table_info: {table_info}")
    print(f"  max_chunk_size: {max_chunk_size}")
    print(f"  preserve_context: {preserve_context}")
    
    if table_data.empty:
        print("🔍 DEBUG: Table data is empty, returning empty list")
        return []
    
    # Estimate size per row
    sample_row = table_data.iloc[0] if len(table_data) > 0 else pd.Series()
    estimated_row_size = sum(len(str(val)) for val in sample_row.values) + 20  # 20 for formatting
    estimated_row_size = max(estimated_row_size, 30)  # Minimum estimate
    
    # Calculate optimal chunk size - keep small to medium tables as single chunks
    if len(table_data) <= 10:  # Increased threshold for single chunk
        max_rows_per_chunk = len(table_data)
        print("🔍 DEBUG: Small/medium table detected, using single chunk")
    else:
        max_rows_per_chunk = max(5, min(20, max_chunk_size // estimated_row_size))
        print(f"🔍 DEBUG: Large table detected, using {max_rows_per_chunk} rows per chunk")
        print(f"  Estimated size per row: {estimated_row_size}")
    
    print("🔍 DEBUG: Calling create_semantic_table_chunks")
    return create_semantic_table_chunks(table_data, table_info, max_rows_per_chunk)


def create_semantic_table_chunks(table_data: pd.DataFrame, table_info: Dict[str, Any], 
                                max_rows_per_chunk: int = 10) -> List[Dict[str, Any]]:
    """
    Create semantically grouped table chunks
    """
    print(f"\n🔍 DEBUG: create_semantic_table_chunks called")
    print(f"  table_data shape: {table_data.shape if not table_data.empty else 'EMPTY'}")
    print(f"  table_info: {table_info}")
    print(f"  max_rows_per_chunk: {max_rows_per_chunk}")
    
    if table_data.empty:
        print("🔍 DEBUG: Table data is empty, returning empty list")
        return []
    
    chunks = []
    
    # Clean the DataFrame
    print("🔍 DEBUG: Cleaning DataFrame")
    cleaned_df = table_data.copy()
    cleaned_df = cleaned_df.fillna('')
    
    print(f"  After fillna: {cleaned_df.shape}")
    print(f"  Columns: {list(cleaned_df.columns)}")
    print(f"  Data sample:")
    
    # Print actual data rows to see what we're working with
    for i, (idx, row) in enumerate(cleaned_df.iterrows()):
        if i < 5:  # Show first 5 rows
            non_empty_values = [str(val) for val in row.values if val and str(val).strip()]
            print(f"    Row {i}: {non_empty_values}")
    
    print(f"  Total non-empty data rows: {len([idx for idx, row in cleaned_df.iterrows() if any(val and str(val).strip() for val in row.values)])}")
    
    # Convert all columns to string
    print("🔍 DEBUG: Converting columns to string")
    for col in cleaned_df.columns:
        cleaned_df[col] = cleaned_df[col].astype(str).str.strip()
        print(f"  Converted column '{col}'")
    
    print(f"🔍 DEBUG: Final cleaned DataFrame:")
    print(cleaned_df.head(3))
    
    # Get table title/identifier
    table_identifier = f"Table {table_info.get('table_number', 1)}"
    if table_info.get('page_number'):
        table_identifier += f" (Page {table_info['page_number']})"
    elif table_info.get('slide_number'):
        table_identifier += f" (Slide {table_info['slide_number']})"
    
    print(f"🔍 DEBUG: Table identifier: {table_identifier}")
    
    # Split into chunks
    total_rows = len(cleaned_df)
    print(f"🔍 DEBUG: Total rows: {total_rows}, max_rows_per_chunk: {max_rows_per_chunk}")
    
    # For small tables, use single chunk
    if max_rows_per_chunk >= total_rows:
        print("🔍 DEBUG: Creating single chunk for small table")
        chunk_content = create_single_table_chunk(cleaned_df, table_identifier)
        print(f"  Single chunk content length: {len(chunk_content)}")
        print(f"  Content preview: {chunk_content[:200]}...")
        
        chunk_metadata = create_table_chunk_metadata(table_info, 0)
        print(f"  Chunk metadata: {chunk_metadata}")
        
        chunks.append({
            'content': chunk_content,
            'metadata': chunk_metadata
        })
        
        print(f"🔍 DEBUG: Created 1 single chunk, returning")
        return chunks
    
    # Multiple chunks needed
    print("🔍 DEBUG: Creating multiple chunks")
    for chunk_start in range(0, total_rows, max_rows_per_chunk):
        chunk_end = min(chunk_start + max_rows_per_chunk, total_rows)
        chunk_df = cleaned_df.iloc[chunk_start:chunk_end]
        
        print(f"  Creating chunk for rows {chunk_start}-{chunk_end-1}")
        print(f"  Chunk df shape: {chunk_df.shape}")
        
        chunk_content = create_multi_table_chunk(
            chunk_df, table_identifier,
            chunk_start, chunk_end, total_rows
        )
        
        print(f"  Multi-chunk content length: {len(chunk_content)}")
        print(f"  Content preview: {chunk_content[:200]}...")
        
        # Create metadata
        chunk_metadata = create_table_chunk_metadata(table_info, chunk_index=chunk_start // max_rows_per_chunk)
        chunk_metadata.update({
            'chunk_type': 'table_group',
            'row_start': chunk_start,
            'row_end': chunk_end - 1,
            'row_count': len(chunk_df),
            'chunking_strategy': 'semantic',
            'formatting_style': 'structured'
        })
        
        print(f"  Multi-chunk metadata: {chunk_metadata}")
        
        chunks.append({
            'content': chunk_content,
            'metadata': chunk_metadata
        })
    
    print(f"🔍 DEBUG: Created {len(chunks)} chunks total")
    return chunks


def create_single_table_chunk(df: pd.DataFrame, table_identifier: str) -> str:
    """Create optimized single chunk for small tables"""
    print(f"\n🔍 DEBUG: create_single_table_chunk called")
    print(f"  df shape: {df.shape}")
    print(f"  table_identifier: {table_identifier}")
    
    lines = []
    
    # Add title and context
    lines.append(f"# {table_identifier}")
    lines.append("")
    lines.append("This is a data table extracted from the document.")
    lines.append("")
    
    # Use markdown table format for better readability
    if len(df.columns) <= 10:  # Most tables should fit this
        # Create header row with proper column names
        header_cells = []
        for col in df.columns:
            # Clean column names for better display
            clean_col = str(col).replace('\n', ' ').replace('\r', ' ').strip()
            header_cells.append(clean_col)
        
        header_row = "| " + " | ".join(header_cells) + " |"
        lines.append(header_row)
        
        # Create separator row
        separator_row = "|" + "|".join("-" * (len(cell) + 2) for cell in header_cells) + "|"
        lines.append(separator_row)
        
        # Add data rows with proper cell formatting
        for _, row in df.iterrows():
            data_cells = []
            for val in row.values:
                # Clean cell values and handle multi-line content
                clean_val = str(val).replace('\n', ' ').replace('\r', ' ').strip()
                if not clean_val or clean_val.lower() in ['nan', 'none', '']:
                    clean_val = "-"
                data_cells.append(clean_val)
            
            data_row = "| " + " | ".join(data_cells) + " |"
            lines.append(data_row)
    else:
        # For very wide tables, use key-value format
        lines.append("**Table Data:**")
        lines.append("")
        for row_idx, (_, row) in enumerate(df.iterrows(), 1):
            lines.append(f"**Row {row_idx}:**")
            for col, val in row.items():
                clean_val = str(val).strip()
                if clean_val and clean_val.lower() not in ['nan', 'none', '']:
                    lines.append(f"- {col}: {clean_val}")
            lines.append("")
    
    result = '\n'.join(lines)
    print(f"🔍 DEBUG: Single chunk result length: {len(result)}")
    print(f"  Preview: {result[:300]}...")
    
    return result


def create_multi_table_chunk(chunk_df: pd.DataFrame, table_identifier: str,
                           chunk_start: int, chunk_end: int, total_rows: int) -> str:
    """Create chunk for multi-chunk tables with context preservation"""
    print(f"\n🔍 DEBUG: create_multi_table_chunk called")
    print(f"  chunk_df shape: {chunk_df.shape}")
    print(f"  chunk_start: {chunk_start}, chunk_end: {chunk_end}, total_rows: {total_rows}")
    
    lines = []
    
    # Chunk title with context
    chunk_title = f"═══ {table_identifier} - Rows {chunk_start + 1} to {chunk_end} ═══"
    lines.append(chunk_title)
    lines.append("")
    
    # Add table context
    lines.append(f"📋 Table Context: {len(chunk_df.columns)} columns, {total_rows} total rows")
    lines.append(f"📊 This Chunk: {len(chunk_df)} rows ({chunk_start + 1}-{chunk_end})")
    lines.append("")
    
    # Use markdown table format
    if len(chunk_df.columns) <= 8:
        # Create header row
        header_row = "| " + " | ".join(str(col) for col in chunk_df.columns) + " |"
        lines.append(header_row)
        
        # Create separator row
        separator_row = "|" + "|".join("--------" for _ in chunk_df.columns) + "|"
        lines.append(separator_row)
        
        # Add data rows
        for _, row in chunk_df.iterrows():
            data_row = "| " + " | ".join(str(val) for val in row.values) + " |"
            lines.append(data_row)
    else:
        # For wide tables, use key-value format
        for row_idx, (_, row) in enumerate(chunk_df.iterrows()):
            lines.append(f"Row {chunk_start + row_idx + 1}:")
            for col, val in row.items():
                if str(val).strip():
                    lines.append(f"  {col}: {val}")
            lines.append("")
    
    result = '\n'.join(lines)
    print(f"🔍 DEBUG: Multi-chunk result length: {len(result)}")
    
    return result


def create_table_chunk_metadata(table_info: Dict[str, Any], chunk_index: int = 0) -> Dict[str, Any]:
    """Create metadata for table chunks"""
    metadata = {
        'chunk_type': 'table',
        'table_number': table_info.get('table_number', 1),
        'extraction_method': table_info.get('extraction_method', 'unknown'),
        'confidence': table_info.get('confidence', 0.8),
        'chunk_index': chunk_index
    }
    
    if table_info.get('page_number'):
        metadata['page_number'] = table_info['page_number']
    
    if table_info.get('slide_number'):
        metadata['slide_number'] = table_info['slide_number']
    
    return metadata


def identify_key_columns(df: pd.DataFrame) -> List[str]:
    """Identify key columns in a table"""
    key_columns = []
    
    for col in df.columns:
        # Check if column looks like an identifier or primary key
        col_lower = str(col).lower()
        if any(keyword in col_lower for keyword in ['id', 'name', 'title', 'category', 'type']):
            key_columns.append(col)
        
        # Check if column has unique or mostly unique values
        try:
            unique_ratio = len(df[col].unique()) / len(df)
            if unique_ratio > 0.8:  # More than 80% unique values
                key_columns.append(col)
        except:
            pass
    
    return key_columns[:3]  # Limit to top 3 key columns


# Global table extractor instance
_table_extractor_instance = None


def get_table_extractor() -> TableExtractor:
    """Get singleton table extractor instance with Docling support"""
    global _table_extractor_instance
    if _table_extractor_instance is None:
        _table_extractor_instance = TableExtractor()
        if DOCLING_AVAILABLE:
            logger.info("Table extractor initialized with Docling support")
        else:
            logger.warning("Table extractor initialized without Docling - using fallback methods")
    return _table_extractor_instance


def create_contextual_table_chunks(table_data: pd.DataFrame, table_info: Dict[str, Any], 
                                  max_chunk_size: int = 7800, preserve_context: bool = True,
                                  surrounding_text: str = None) -> List[Dict[str, Any]]:
    """
    Create contextual table chunks (alias for create_adaptive_table_chunks for backward compatibility)
    
    Args:
        table_data: DataFrame containing the table data
        table_info: Metadata about the table
        max_chunk_size: Maximum size for each chunk
        preserve_context: Whether to preserve context across chunks
        surrounding_text: Text context surrounding the table (unused but kept for compatibility)
    """
    print(f"\n🔍 DEBUG: create_contextual_table_chunks called with surrounding_text: {bool(surrounding_text)}")
    return create_adaptive_table_chunks(table_data, table_info, max_chunk_size, preserve_context)
