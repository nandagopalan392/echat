import os
import tempfile
import logging
from pathlib import Path
from typing import List, Dict, Optional
import base64
from io import BytesIO
from PIL import Image, ImageDraw, ImageFont
import html

logger = logging.getLogger(__name__)

class PPTXHTMLConverter:
    """Convert PPTX slides to HTML format for better web rendering"""
    
    def __init__(self, temp_dir: Optional[str] = None):
        self.temp_dir = temp_dir or tempfile.gettempdir()
    
    def _inches_to_px(self, value_in_inches: float) -> str:
        """Convert inches to pixels (1 inch ≈ 96px)"""
        return f"{value_in_inches * 96:.1f}px"
    
    def _get_slide_dimensions_px(self, prs) -> tuple:
        """Get slide dimensions in pixels"""
        width_inches = prs.slide_width.inches if hasattr(prs.slide_width, 'inches') else 10
        height_inches = prs.slide_height.inches if hasattr(prs.slide_height, 'inches') else 7.5
        return (width_inches * 96, height_inches * 96)
        
    def convert_pptx_to_html_slides(self, pptx_path: str) -> Dict:
        """
        Convert PPTX slides to HTML format
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            Dictionary with slide information in HTML format
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            prs = Presentation(pptx_path)
            slides = []
            
            # Get slide dimensions in pixels
            slide_width_px, slide_height_px = self._get_slide_dimensions_px(prs)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_html = self._convert_slide_to_html(slide, slide_num, slide_width_px, slide_height_px)
                
                slides.append({
                    'slide_number': slide_num,
                    'html_content': slide_html,
                    'format': 'html',
                    'width': slide_width_px / 96,  # Keep in inches for compatibility
                    'height': slide_height_px / 96,
                    'width_px': slide_width_px,
                    'height_px': slide_height_px,
                    'is_html': True
                })
            
            return {
                'slides': slides,
                'total_slides': len(prs.slides),
                'conversion_method': 'html',
                'has_html': True
            }
            
        except Exception as e:
            logger.error(f"PPTX to HTML conversion error: {e}")
            return self._create_fallback_html(pptx_path)
    
    def _convert_slide_to_html(self, slide, slide_num: int, slide_width_px: float, slide_height_px: float) -> str:
        """Convert a single slide to HTML"""
        html_content = f"""
        <div class="pptx-slide" style="
            position: relative;
            width: {slide_width_px:.1f}px;
            height: {slide_height_px:.1f}px;
            background: white;
            border: 1px solid #eee;
            overflow: hidden;
            margin: 0 auto;
        ">
        """
        
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or hasattr(shape, 'text'):
                    html_content += self._convert_text_shape_to_html(shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    html_content += self._convert_table_to_html(shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    html_content += self._convert_image_to_html(shape, slide_num)
                elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    html_content += self._convert_shape_to_html(shape)
                    
        except Exception as e:
            logger.warning(f"Error processing slide {slide_num}: {e}")
            html_content += f'<div style="padding: 20px; color: red;">Error processing slide content: {str(e)}</div>'
        
        html_content += """
        </div>
        """
        
        return html_content
    
    def _convert_text_shape_to_html(self, shape) -> str:
        """Convert text shape to HTML"""
        if not hasattr(shape, 'text') or not shape.text.strip():
            return ""
        
        text = html.escape(shape.text.strip())
        
        # Try to get text formatting
        try:
            # Get font size and style information
            font_size = "16px"
            font_weight = "normal"
            text_align = "left"
            
            if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    if hasattr(run.font, 'size') and run.font.size:
                        font_size = f"{run.font.size.pt}px"
                    if hasattr(run.font, 'bold') and run.font.bold:
                        font_weight = "bold"
                
                if hasattr(para, 'alignment'):
                    alignment_map = {
                        1: "left",    # PP_ALIGN.LEFT
                        2: "center",  # PP_ALIGN.CENTER
                        3: "right",   # PP_ALIGN.RIGHT
                    }
                    text_align = alignment_map.get(para.alignment, "left")
            
            # Get position information and convert to pixels
            left_px = self._inches_to_px(shape.left.inches if hasattr(shape.left, 'inches') else 0)
            top_px = self._inches_to_px(shape.top.inches if hasattr(shape.top, 'inches') else 0)
            width_px = self._inches_to_px(shape.width.inches if hasattr(shape.width, 'inches') else 2)
            height_px = self._inches_to_px(shape.height.inches if hasattr(shape.height, 'inches') else 1)
            
            style = f"""
                position: absolute;
                left: {left_px};
                top: {top_px};
                width: {width_px};
                height: {height_px};
                font-size: {font_size};
                font-weight: {font_weight};
                text-align: {text_align};
                padding: 4px;
                box-sizing: border-box;
                overflow: hidden;
                word-wrap: break-word;
            """
            
            return f'<div class="text-shape" style="{style}">{text}</div>'
            
        except Exception as e:
            logger.warning(f"Error formatting text shape: {e}")
            return f'<div class="text-shape simple">{text}</div>'
    
    def _convert_table_to_html(self, shape) -> str:
        """Convert table shape to HTML"""
        if not hasattr(shape, 'table'):
            return ""
        
        try:
            table = shape.table
            html_table = '<table class="slide-table" style="border-collapse: collapse; margin: 10px;">'
            
            for row in table.rows:
                html_table += '<tr>'
                for cell in row.cells:
                    cell_text = html.escape(cell.text.strip()) if cell.text.strip() else "&nbsp;"
                    html_table += f'<td style="border: 1px solid #ccc; padding: 8px;">{cell_text}</td>'
                html_table += '</tr>'
            
            html_table += '</table>'
            return html_table
            
        except Exception as e:
            logger.warning(f"Error converting table: {e}")
            return '<div class="table-error">Table content unavailable</div>'
    
    def _convert_image_to_html(self, shape, slide_num: int) -> str:
        """Convert image shape to HTML"""
        try:
            if hasattr(shape, 'image'):
                # Extract image data and convert to base64
                image = shape.image
                image_bytes = image.blob
                
                # Determine image format
                image_format = "png"
                if image_bytes.startswith(b'\xff\xd8'):
                    image_format = "jpeg"
                elif image_bytes.startswith(b'\x89PNG'):
                    image_format = "png"
                elif image_bytes.startswith(b'GIF'):
                    image_format = "gif"
                
                # Convert to base64
                image_b64 = base64.b64encode(image_bytes).decode('utf-8')
                data_url = f"data:image/{image_format};base64,{image_b64}"
                
                # Get position information and convert to pixels
                left_px = self._inches_to_px(shape.left.inches if hasattr(shape.left, 'inches') else 0)
                top_px = self._inches_to_px(shape.top.inches if hasattr(shape.top, 'inches') else 0)
                width_px = self._inches_to_px(shape.width.inches if hasattr(shape.width, 'inches') else 2)
                height_px = self._inches_to_px(shape.height.inches if hasattr(shape.height, 'inches') else 2)
                
                style = f"""
                    position: absolute;
                    left: {left_px};
                    top: {top_px};
                    width: {width_px};
                    height: {height_px};
                    object-fit: contain;
                """
                
                return f'<img class="slide-image" src="{data_url}" style="{style}" alt="Slide {slide_num} Image" />'
            
        except Exception as e:
            logger.warning(f"Error converting image: {e}")
        
        return '<div class="image-placeholder">Image unavailable</div>'
    
    def _convert_shape_to_html(self, shape) -> str:
        """Convert other shapes to HTML"""
        try:
            if hasattr(shape, 'text') and shape.text.strip():
                return self._convert_text_shape_to_html(shape)
            else:
                # For other shapes, create a placeholder
                shape_type = getattr(shape, 'shape_type', 'unknown')
                return f'<div class="shape-placeholder">Shape: {shape_type}</div>'
        except Exception as e:
            logger.warning(f"Error converting shape: {e}")
            return '<div class="shape-error">Shape unavailable</div>'
    
    def _create_fallback_html(self, pptx_path: str) -> Dict:
        """Create fallback HTML when conversion fails"""
        try:
            from pptx import Presentation
            
            prs = Presentation(pptx_path)
            slides = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # Extract basic text content
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text_content.append(html.escape(shape.text.strip()))
                
                # Create simple HTML with proper container
                slide_html = f"""
                <div class="pptx-slide fallback" style="
                    position: relative;
                    width: 960px;
                    height: 540px;
                    background: white;
                    border: 1px solid #eee;
                    padding: 20px;
                    box-sizing: border-box;
                    margin: 0 auto;
                ">
                    <div class="slide-header" style="margin-bottom: 15px;">
                        <h3 style="margin: 0; color: #333;">Slide {slide_num}</h3>
                    </div>
                    <div class="slide-content" style="font-size: 14px; line-height: 1.5;">
                        {"<br><br>".join(text_content) if text_content else "<p>No text content available</p>"}
                    </div>
                </div>
                """
                
                slides.append({
                    'slide_number': slide_num,
                    'html_content': slide_html,
                    'format': 'html',
                    'width': 10,
                    'height': 7.5,
                    'is_html': True,
                    'is_fallback': True
                })
            
            return {
                'slides': slides,
                'total_slides': len(prs.slides),
                'conversion_method': 'html_fallback',
                'has_html': True
            }
            
        except Exception as e:
            logger.error(f"Fallback HTML conversion error: {e}")
            return {
                'slides': [],
                'total_slides': 0,
                'conversion_method': 'error',
                'has_html': False,
                'error': str(e)
            }

def convert_pptx_to_html_slides(pptx_file_path: str) -> Dict:
    """
    Main function to convert PPTX to HTML slides
    
    Args:
        pptx_file_path: Path to the PPTX file
        
    Returns:
        Dictionary with HTML slide information
    """
    converter = PPTXHTMLConverter()
    return converter.convert_pptx_to_html_slides(pptx_file_path)
